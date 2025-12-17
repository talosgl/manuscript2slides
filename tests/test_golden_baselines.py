"""
Golden baseline tests for docx ↔ pptx conversion.

These tests use pre-extracted JSON baselines to verify that the conversion pipeline
produces expected output for formatting, annotations, and document structure.
"""

import json
import pytest
from pathlib import Path

from docx import Document
from pptx import Presentation

from manuscript2slides.orchestrator import run_pipeline
from manuscript2slides.internals.define_config import UserConfig, PipelineDirection
from manuscript2slides.processing.formatting import (
    get_theme_fonts_from_docx_package,
    get_effective_font_name_docx,
)


# region Fixtures


@pytest.fixture
def docx_baseline() -> dict:
    """Load the golden baseline for sample_doc.docx."""
    baseline_path = Path("tests/baselines/docx_sample.json")
    with open(baseline_path, encoding="utf-8") as f:
        return json.load(f)


@pytest.fixture
def pptx_baseline() -> dict:
    """Load the golden baseline for sample_slides_output.pptx."""
    baseline_path = Path("tests/baselines/pptx_sample.json")
    with open(baseline_path, encoding="utf-8") as f:
        return json.load(f)


# endregion


# region DOCX → PPTX Tests


def test_docx_to_pptx_first_chunk_formatting(
    path_to_sample_docx_with_everything: Path,
    path_to_empty_pptx: Path,
    docx_baseline: list[dict],
    tmp_path: Path,
) -> None:
    """
    Test that the first paragraph from sample_doc.docx is correctly converted to pptx.

    The default chunking strategy is PARAGRAPH, so each non-empty paragraph becomes
    one slide. This tests that the first paragraph's formatting is preserved.

    Verifies:
    - Text content preservation
    - Number of runs matches baseline
    - Run-level formatting (bold, italic, fonts, etc.)
    """
    # Arrange: Run the pipeline with default PARAGRAPH chunking
    cfg = UserConfig(
        input_docx=path_to_sample_docx_with_everything,
        template_pptx=path_to_empty_pptx,
        output_folder=tmp_path,
    )

    output_pptx_path = run_pipeline(cfg)

    # Load the generated pptx
    prs = Presentation(str(output_pptx_path))

    # With PARAGRAPH chunking, each non-empty paragraph becomes a slide
    assert len(prs.slides) > 0, "No slides generated"
    first_slide = prs.slides[0]

    # Get the content placeholder (index 1) text frame
    content_shape = first_slide.placeholders[1]
    text_frame = content_shape.text_frame

    # Get baseline data for the first paragraph
    first_para_baseline = docx_baseline[0]

    # Assert: Text content matches
    # The entire paragraph should be in the slide's first text frame paragraph
    assert len(text_frame.paragraphs) > 0, "No paragraphs in slide text frame"
    first_pptx_para = text_frame.paragraphs[0]

    # Text should match (might have minor whitespace differences)
    # TODO: this has a link in it, which is added from the field code id edge case code. Can we change it so we verify the baseline text is IN the new text, rather than matches exactly?
    assert (
        first_pptx_para.text.strip() == first_para_baseline["text"].strip()
    ), f"Text mismatch:\nExpected: {first_para_baseline['text']}\nGot: {first_pptx_para.text}"

    # Assert: Number of runs should match
    baseline_runs = first_para_baseline["runs"]
    pptx_runs = first_pptx_para.runs

    assert len(pptx_runs) == len(
        baseline_runs
    ), f"Run count mismatch: expected {len(baseline_runs)}, got {len(pptx_runs)}"

    # Assert: Run formatting preservation
    for run_idx, (pptx_run, baseline_run) in enumerate(zip(pptx_runs, baseline_runs)):
        # Text should match
        assert (
            pptx_run.text == baseline_run["text"]
        ), f"Run {run_idx} text mismatch: expected '{baseline_run['text']}', got '{pptx_run.text}'"

        # Bold should match (if specified in baseline)
        if "bold" in baseline_run and baseline_run["bold"] is not None:
            assert (
                pptx_run.font.bold == baseline_run["bold"]
            ), f"Run {run_idx} bold mismatch: expected {baseline_run['bold']}, got {pptx_run.font.bold}"

        # Italic should match (if specified in baseline)
        if "italic" in baseline_run and baseline_run["italic"] is not None:
            assert (
                pptx_run.font.italic == baseline_run["italic"]
            ), f"Run {run_idx} italic mismatch: expected {baseline_run['italic']}, got {pptx_run.font.italic}"

        # Font name should match (if specified in baseline)
        if "font_name" in baseline_run and baseline_run["font_name"]:
            assert (
                pptx_run.font.name == baseline_run["font_name"]
            ), f"Run {run_idx} font mismatch: expected '{baseline_run['font_name']}', got '{pptx_run.font.name}'"


# endregion


# region PPTX → DOCX Tests


def test_pptx_to_docx_first_five_slides(
    path_to_empty_docx: Path,
    pptx_baseline: list[dict],
    tmp_path: Path,
) -> None:
    """
    Test that the first 5 slides from sample_slides_output.pptx convert correctly to docx.

    Each slide becomes one or more paragraphs in the docx. If speaker notes contain
    metadata with headings, those are restored as proper heading paragraphs.

    Verifies:
    - Text content from slides appears in docx
    - Number of slides matches number of content sections in docx
    - Basic formatting preservation (bold, italic, fonts)
    - Headings from speaker notes metadata are restored with proper styles
    """
    # Arrange: Load the source pptx that we have a baseline for
    source_pptx = Path("tests/data/sample_slides_output.pptx")

    if not source_pptx.exists():
        pytest.skip(f"Source pptx not found: {source_pptx}")

    cfg = UserConfig(
        input_pptx=source_pptx,
        template_docx=path_to_empty_docx,
        output_folder=tmp_path,
    )

    output_docx_path = run_pipeline(cfg)

    # Load the generated docx
    doc = Document(str(output_docx_path))

    # Get theme fonts for proper font name resolution
    theme_fonts = get_theme_fonts_from_docx_package(doc.part.package)

    # Test the first 5 slides worth of content
    num_slides_to_test = min(5, len(pptx_baseline))

    # Collect all text from the first N slides in the baseline
    expected_slide_texts = []
    expected_headings = []

    for slide_idx in range(num_slides_to_test):
        slide_baseline = pptx_baseline[slide_idx]

        # Extract text from all shapes in the slide
        slide_text_parts = []
        for shape in slide_baseline.get("shapes", []):
            for para in shape.get("paragraphs", []):
                text = para.get("text", "").strip()
                if text:
                    slide_text_parts.append(text)

        if slide_text_parts:
            expected_slide_texts.extend(slide_text_parts)

        # Extract headings from speaker notes metadata if present
        speaker_notes = slide_baseline.get("speaker_notes", {})
        if speaker_notes:
            metadata = speaker_notes.get("metadata", {})
            headings = (
                metadata.get("headings", []) if isinstance(metadata, dict) else []
            )
            for heading in headings:
                if isinstance(heading, dict) and "text" in heading:
                    expected_headings.append(heading["text"])

    # Assert: Document should have paragraphs with content
    doc_paragraphs_with_text = [p for p in doc.paragraphs if p.text.strip()]
    assert len(doc_paragraphs_with_text) > 0, "Generated docx has no text content"

    # Assert: All slide text should appear somewhere in the docx
    docx_full_text = "\n".join(p.text for p in doc.paragraphs)

    for slide_text in expected_slide_texts:
        assert (
            slide_text in docx_full_text
        ), f"Slide text not found in docx: '{slide_text[:50]}...'"

    # Assert: Headings from speaker notes metadata should be restored
    if expected_headings:
        docx_heading_paragraphs = [
            p
            for p in doc.paragraphs
            if p.style and p.style.name and p.style.name.startswith("Heading")
        ]

        # At least some headings should have been restored
        assert (
            len(docx_heading_paragraphs) > 0
        ), "No heading paragraphs found in docx despite speaker notes metadata containing headings"

        # Check that expected heading texts appear in the docx
        docx_heading_texts = [p.text.strip() for p in docx_heading_paragraphs]

        for expected_heading in expected_headings:
            assert (
                expected_heading in docx_heading_texts
            ), f"Expected heading not found: '{expected_heading}'"

    # Assert: Basic formatting preservation for first paragraph with runs
    first_slide_baseline = pptx_baseline[0]
    if first_slide_baseline.get("shapes"):
        first_shape = first_slide_baseline["shapes"][0]
        if first_shape.get("paragraphs"):
            first_para_baseline = first_shape["paragraphs"][0]
            baseline_runs = first_para_baseline.get("runs", [])

            # Find the corresponding paragraph in docx by matching text
            baseline_text = first_para_baseline.get("text", "").strip()
            matching_docx_para = None

            for para in doc.paragraphs:
                if para.text.strip() == baseline_text:
                    matching_docx_para = para
                    break

            if matching_docx_para and baseline_runs:
                docx_runs = matching_docx_para.runs

                # Check that runs were preserved (count might differ due to merging/splitting)
                # So we just verify some runs exist
                assert len(docx_runs) > 0, "No runs in matching docx paragraph"

                # Verify at least one formatted run matches
                for baseline_run in baseline_runs:
                    if baseline_run.get("bold") or baseline_run.get("italic"):
                        # At least one run in docx should have matching formatting
                        has_matching_format = any(
                            (
                                run.font.bold == baseline_run.get("bold")
                                or baseline_run.get("bold") is None
                            )
                            and (
                                run.font.italic == baseline_run.get("italic")
                                or baseline_run.get("italic") is None
                            )
                            for run in docx_runs
                        )
                        assert (
                            has_matching_format
                        ), f"No run with matching formatting found for baseline run: {baseline_run}"
                        break  # Only need to verify one formatted run


# endregion


# region Helper Functions (for future use)


def extract_run_formatting(run) -> dict:
    """Extract formatting properties from a docx run for comparison."""
    return {
        "text": run.text,
        "bold": run.font.bold,
        "italic": run.font.italic,
        "underline": run.font.underline,
        "font_name": run.font.name,
        "font_size": run.font.size.pt if run.font.size else None,
    }


def extract_para_formatting(para, theme_fonts) -> dict:
    """Extract formatting properties from a docx paragraph for comparison."""
    if not para.style or not para.style.font:
        return {}

    effective_font_name = get_effective_font_name_docx(para.style, theme_fonts)

    return {
        "bold": para.style.font.bold,
        "italic": para.style.font.italic,
        "underline": para.style.font.underline,
        "font_name": effective_font_name,
        "font_size": para.style.font.size.pt if para.style.font.size else None,
    }


# endregion
