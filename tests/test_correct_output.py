"""Tests for correctness of pipeline output based on manual test cases performed during development."""

# pyright: reportAttributeAccessIssue=false
import pytest
from pathlib import Path

import pptx
from pptx.text.text import TextFrame
from pptx.text.text import _Hyperlink as Hyperlink_pptx
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE as MSO_TEXT_UNDERLINE_TYPE_pptx
from pptx.text.text import _Paragraph as Paragraph_pptx
from pptx.text.text import _Run as Run_pptx
from pptx.dml.color import RGBColor as RGBColor_pptx


import docx
import docx.document
from docx.text.hyperlink import Hyperlink as Hyperlink_docx
from docx.text.run import Run as Run_docx
from docx.shared import RGBColor as RGBColor_docx
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_COLOR_INDEX, WD_UNDERLINE


from manuscript2slides.internals.define_config import UserConfig
from tests import helpers


# region docx2pptx tests

# When converting from sample_doc.docx -> standard pptx output, with these options:
#     - preserve experimental formatting
#     - keep all annotations
#     - preserve metadata in speaker notes
# ...the results should match the assertions in the tests below.

def test_where_are_data_slide(output_pptx: Path) -> None:
    """Verify formatting preservation (bold, italic, color, typeface) and metadata/comments in speaker notes."""

    # Action: note that the "Action" here is running the docx2pptx pipeline, which we do
    # in a fixture one time instead of repeatedly for many tests.
    prs = pptx.Presentation(output_pptx)

    # Arrange: Find the slide and paragraph containing the formatted text we are going to assert against.
    slide, para = helpers.find_first_slide_containing(prs, "Where are Data?")

    # Assert
    # Case: The paragraph should be bolded at the paragraph-level
    assert para.font.bold is True

    # Case: "are" within "Where are Data?" should be in red
    # We have to do exact match, or the URL run will be returned
    run_result = helpers.find_first_pptx_run_in_para_with_exact_match(para, "are")

    # Verify color object exists and has RGB value before checking specific color
    assert run_result.font.color is not None
    assert (
        hasattr(run_result.font.color, "rgb") and run_result.font.color.rgb is not None
    )
    assert run_result.font.color.rgb == RGBColor_pptx(0xFF, 0x00, 0x00)  # red
    assert run_result.font.italic is True

    # Case: Paragraph text should contain link text from a field code hyperlink
    assert "http" in para.text

    notes_text_frame: TextFrame = helpers.get_slide_notes_text(slide)

    # Case: Speaker notes should contain copied comments (UserConfig.display_comments is true)
    assert "COMMENTS FROM SOURCE DOCUMENT" in notes_text_frame.text
    assert '"text": "What happens if' in notes_text_frame.text
    assert "What happens if there's a threaded comment?" in notes_text_frame.text

    # Case: Speaker notes should contain formatting data for the heading (preserve_docx_metadata_in_speaker_notes is true)
    assert '"headings": [' in notes_text_frame.text
    assert '"name": "Heading 1"' in notes_text_frame.text


def test_author_slide(output_pptx: Path) -> None:
    """Verify heading-level formatting (color, italics) is copied to paragraph level and metadata preserved in speaker notes."""

    prs = pptx.Presentation(output_pptx)
    slide, para = helpers.find_first_slide_containing(prs, "by J. King-Yost")

    # Case: The text should be gray and in italics
    # NOTE: Because this is a heading, the formatting is applied at the paragraph-level, not the run-level.
    # We will explicitly check the paragraph's font here; it is possible the test may need to be adapted
    # to be more flexible and to check both (either/or) in the future.

    # Verify color object exists and has RGB value before checking specific color
    assert para.font.color is not None
    assert hasattr(para.font.color, "rgb") and para.font.color.rgb is not None
    assert para.font.color.rgb == RGBColor_pptx(89, 89, 89)  # gray
    assert para.font.italic is True

    notes_text_frame: TextFrame = helpers.get_slide_notes_text(slide)

    assert "JSON METADATA FROM SOURCE DOCUMENT" in notes_text_frame.text
    assert '"headings": [' in notes_text_frame.text
    assert '"name": "Heading 3"' in notes_text_frame.text


def test_in_a_cold_concrete_underground_tunnel_slide(output_pptx: Path) -> None:
    """Verify standard formatting (italic, bold, color, underline, hyperlinks), experimental formatting (highlights), and speaker notes (footnotes, comments, metadata)."""
    prs = pptx.Presentation(output_pptx)
    slide, para = helpers.find_first_slide_containing(
        prs, "In a cold concrete underground tunnel"
    )

    # Case: Hyperlinks are preserved
    # "In a cold concrete underground tunnel" should be a functional hyperlink to https://dataepics.webflow.io/stories/where-are-data
    hyperlink_run = helpers.find_first_pptx_run_in_para_containing(
        para, "In a cold concrete underground"
    )

    # Verify hyperlink object exists and is correct type before checking URL
    assert (
        hyperlink_run.hyperlink is not None
        and isinstance(hyperlink_run.hyperlink, Hyperlink_pptx)
        and hyperlink_run.hyperlink.address is not None
    )
    assert (
        hyperlink_run.hyperlink.address
        == "https://dataepics.webflow.io/stories/where-are-data"
    )

    # Case: Hyperlinks AND standard formatting is preserved in the same run
    # the "tunnel" part of the clause should be in italics and maintain the same hyperlink as the above.
    hyperlink_run_formatted = helpers.find_first_pptx_run_in_para_with_exact_match(
        para, "tunnel"
    )  # Tunnel shows up in multiple spots in this paragraph so we're relying on this being the first one

    # Verify hyperlink object exists and is correct type before checking URL
    assert (
        hyperlink_run_formatted.hyperlink is not None
        and isinstance(hyperlink_run_formatted.hyperlink, Hyperlink_pptx)
        and hyperlink_run_formatted.hyperlink.address is not None
    )
    assert (
        hyperlink_run_formatted.hyperlink.address
        == "https://dataepics.webflow.io/stories/where-are-data"
    )
    assert hyperlink_run_formatted.font.italic is True

    # Case: Standard formatting is preserved - "splayed" should be bolded
    bolded_run = helpers.find_first_pptx_run_in_para_with_exact_match(para, "splayed")
    assert bolded_run.font.bold is True

    # Case: "three dozen directions" should be italic
    italic_run = helpers.find_first_pptx_run_in_para_with_exact_match(
        para, "three dozen directions"
    )
    assert italic_run.font.italic is True

    # Case: "Dust covers the cables" should be in red text
    red_run = helpers.find_first_pptx_run_in_para_containing(
        para, "Dust covers the cables"
    )
    # Verify color object exists and has RGB value before checking specific color
    assert red_run.font.color is not None
    assert hasattr(red_run.font.color, "rgb") and red_run.font.color.rgb is not None
    assert red_run.font.color.rgb == RGBColor_pptx(0xFF, 0x00, 0x00)  # red

    # Verify this run has a custom font set on it
    assert red_run.font.name == "Georgia"

    # Case: "She could wipe them clean without too much risk" should be underlined
    ul_run = helpers.find_first_pptx_run_in_para_containing(
        para, "She could wipe them clean without too much risk"
    )
    assert ul_run.font.underline is True

    # Case: Experimental formatting is preserved ("buzzing" should be highlighted in yellow)
    highlight_run = helpers.find_first_pptx_run_in_para_with_exact_match(
        para, "buzzing"
    )
    assert helpers.run_has_highlight(highlight_run)
    hl_color = helpers.get_run_highlight_color(highlight_run)
    assert hl_color == "FFFF00"

    # Case: "read those stories" should be underlined... double-underlined!
    double_ul_run = helpers.find_first_pptx_run_in_para_containing(
        para, "read those stories"
    )
    assert double_ul_run.font.underline == MSO_TEXT_UNDERLINE_TYPE_pptx.DOUBLE_LINE

    # Speaker notes checks:
    notes_text_frame = helpers.get_slide_notes_text(slide)

    # Case: Comment is copied in when display_comments is enabled
    # Verify both comment section header and content are present
    assert (
        "COMMENTS FROM SOURCE DOCUMENT" in notes_text_frame.text
        and "Sample comment" in notes_text_frame.text
    )

    # Case: Footnote is copied in when display_footnotes is enabled
    # Verify footnote section, content, and hyperlink are all present
    assert (
        "FOOTNOTES FROM SOURCE DOCUMENT" in notes_text_frame.text
        and "1. James Griffiths." in notes_text_frame.text
        and "Hyperlinks:" in notes_text_frame.text
        and r"https://www.cnn.com/2019/07/25/asia/internet-undersea-cables-intl-hnk/index.html"
        in notes_text_frame.text
    )

    # Case: Footnote JSON is preserved when "preserve_metadata..." is enabled, with its hyperlink
    assert '"hyperlinks": [' in notes_text_frame.text

    # Case: experimental formatting data is copied in when "preserve_metadata..." is enabled
    assert "JSON METADATA FROM SOURCE DOCUMENT" in notes_text_frame.text
    # Verify the experimental_formatting JSON structure and all expected fields for "buzzing" highlight
    assert (
        '"experimental_formatting": [' in notes_text_frame.text
        and '"ref_text": "buzzing"' in notes_text_frame.text
        and '"highlight_color_enum": "YELLOW"' in notes_text_frame.text
        and '"formatting_type": "highlight"' in notes_text_frame.text
    )
    # Verify the footnote data was also preserved in the JSON metadata
    assert (
        '"footnotes": [' in notes_text_frame.text
        and '"hyperlinks": [' in notes_text_frame.text
        and '"note_type": "footnote"' in notes_text_frame.text
    )


def test_endnotes_slide(output_pptx: Path) -> None:
    """Verify endnotes and their metadata (including hyperlinks) are copied to speaker notes."""
    prs = pptx.Presentation(output_pptx)
    slide, para = helpers.find_first_slide_containing(prs, "Vedantam, Shankar")

    notes_text_frame = helpers.get_slide_notes_text(slide)

    # Case: Endnotes get copied into the speaker notes when display_endnotes is True
    assert "ENDNOTES FROM SOURCE DOCUMENT" in notes_text_frame.text
    assert "Another sample endnote with a url" in notes_text_frame.text
    assert "Hyperlinks:" in notes_text_frame.text
    assert (
        r"https://dataepics.webflow.io/stories/where-are-data" in notes_text_frame.text
    )

    # Case: Endnotes metadata get copied into the speaker notes when preserve metadata is True
    assert (
        "JSON METADATA FROM SOURCE DOCUMENT" in notes_text_frame.text
        and '"endnotes": [' in notes_text_frame.text
        and '"hyperlinks": [' in notes_text_frame.text
        and '"reference_text": "Vedantam, Shankar, and Iain McGilchrist.'
        in notes_text_frame.text
        and '"note_type": "endnote"' in notes_text_frame.text
    )


def test_vanilla_docx_theme_font_does_not_carry_into_pptx_output(
    output_pptx: Path,
) -> None:
    """Verify docx theme font is not copied to pptx (pptx theme fonts should be respected)."""
    prs = pptx.Presentation(output_pptx)

    _slide, para = helpers.find_first_slide_containing(
        prs, "In a cold concrete underground tunnel"
    )

    assert para.runs[0].font.name is None


def test_speaker_notes_empty_if_json_and_annotations_off(
    output_pptx_default_options: Path,
) -> None:
    """Verify speaker notes are empty when all annotation and metadata display options are disabled."""

    prs = pptx.Presentation(output_pptx_default_options)

    slide, _para = helpers.find_first_slide_containing(
        prs, "In a cold concrete underground tunnel"
    )
    notes_text_frame = helpers.get_slide_notes_text(slide)

    assert notes_text_frame.text.strip() == ""


# endregion


# region pptx2docx
def test_header_formatting_restoration_pptx2docx(output_docx: Path) -> None:
    """Verify heading styles (Heading 1-3) and formatting (color, italics, alignment) are restored from pptx to docx."""

    doc = docx.Document(str(output_docx))

    first_header_para = helpers.find_first_docx_para_containing(doc, "Where are Data?")

    # Case: Heading 1 data gets restored
    assert (
        first_header_para.style
        and first_header_para.style.name
        and first_header_para.style.name == "Heading 1"
    )

    # Case Heading 1 is centered
    assert first_header_para.paragraph_format.alignment == WD_ALIGN_PARAGRAPH.CENTER

    # Case: "are" in this header is in red text
    red_are_run = helpers.find_first_docx_run_in_para_exact_match(
        first_header_para, "are"
    )

    # Verify color object exists and has RGB value before checking specific color
    assert red_are_run.font.color is not None
    assert (
        hasattr(red_are_run.font.color, "rgb")
        and red_are_run.font.color.rgb is not None
    )
    assert red_are_run.font.color.rgb == RGBColor_docx(0xFF, 0x00, 0x00)  # red

    second_header_para = helpers.find_first_docx_para_containing(doc, "by J. King-Yost")

    # Case: Heading 2 styling is restored
    assert (
        second_header_para.style
        and second_header_para.style.name
        and second_header_para.style.name == "Heading 2"
    )

    # Case: Heading 2 is in italics and centered
    assert second_header_para.style.font.italic is True
    assert second_header_para.paragraph_format.alignment == WD_ALIGN_PARAGRAPH.CENTER

    # Case Heading 3 styling is restored, and in gray text
    third_header_para = helpers.find_first_docx_para_containing(
        doc, "Intro to Data: Places we are embodied"
    )

    assert (
        third_header_para.style
        and third_header_para.style.name
        and third_header_para.style.name == "Heading 3"
    )
    assert third_header_para.style.font.color is not None
    assert (
        hasattr(third_header_para.style.font.color, "rgb")
        and third_header_para.style.font.color.rgb is not None
    )
    assert third_header_para.style.font.color.rgb == RGBColor_docx(
        0x99, 0x99, 0x99
    )  # gray


def test_run_formatting_restoration_pptx2docx(output_docx: Path) -> None:
    """Verify run-level formatting (bold, italic, color, underline, highlights, hyperlinks, typeface) is restored from pptx to docx."""

    doc = docx.Document(str(output_docx))

    cold_underground_para = helpers.find_first_docx_para_containing(
        doc, "In a cold concrete underground tunnel"
    )

    # Case: Check that the expected texts are hyperlinks
    first_hyperlink = helpers.find_first_docx_item_in_para_containing(
        cold_underground_para, "In a cold concrete underground"
    )
    assert isinstance(first_hyperlink, Hyperlink_docx)
    second_hyperlink = helpers.find_first_docx_item_in_para_containing(
        cold_underground_para, "tunnel"
    )
    assert isinstance(second_hyperlink, Hyperlink_docx)

    # Case: Hyperlink text maintains italic formatting
    hyperlink_inner_run = helpers.find_inner_run_containing(second_hyperlink, "tunnel")
    assert hyperlink_inner_run.italic is True

    # Case: Bold formatting is preserved
    bold_run = helpers.find_first_docx_item_in_para_containing(
        cold_underground_para, "splayed"
    )
    assert bold_run.bold is True

    # Case: Italics formatting is preserved
    italic_run = helpers.find_first_docx_item_in_para_containing(
        cold_underground_para, "three dozen directions"
    )
    assert italic_run.italic is True

    # Case: Yellow highlight is preserved
    hl_run = helpers.find_first_docx_item_in_para_containing(
        cold_underground_para, "buzzing"
    )
    assert isinstance(hl_run, Run_docx)
    assert hl_run.font.highlight_color == WD_COLOR_INDEX.YELLOW

    # Case: Red color and Georgia typeface are preserved
    red_typeface_run = helpers.find_first_docx_item_in_para_containing(
        cold_underground_para, "Dust covers the cables"
    )
    assert isinstance(red_typeface_run, Run_docx)
    assert red_typeface_run.font.color is not None
    assert (
        hasattr(red_typeface_run.font.color, "rgb")
        and red_typeface_run.font.color.rgb is not None
    )
    assert red_typeface_run.font.color.rgb == RGBColor_docx(0xFF, 0x00, 0x00)  # red
    assert red_typeface_run.font.name == "Georgia"

    ul_run = helpers.find_first_docx_item_in_para_containing(
        cold_underground_para,
        "She could wipe them clean without too much risk",
    )
    assert ul_run.underline is True

    # Case: Double underline is preserved
    double_ul_run = helpers.find_first_docx_item_in_para_containing(
        cold_underground_para, "read those stories"
    )
    assert isinstance(double_ul_run, Run_docx)
    assert double_ul_run.font.underline == WD_UNDERLINE.DOUBLE


def test_annotations_are_restored_pptx2docx(output_docx: Path) -> None:
    """Verify comments, footnotes, and endnotes are restored from pptx speaker notes to docx comments."""
    doc = docx.Document(str(output_docx))

    # Case: At least one comment contains "Sample Comment"
    assert any(
        "Sample comment" in comment.text for comment in doc.comments
    ), "Expected to find 'Sample comment' in at least one comment"

    # Case: At least one comment contains "Footnote:"
    assert any(
        "Footnote:" in comment.text for comment in doc.comments
    ), "Expected to find 'Footnote:' in at least one comment"

    # Case: At least one comment contains endnote text
    assert any(
        "Endnote:" in comment.text
        and "Another sample endnote with a url" in comment.text
        for comment in doc.comments
    ), "Expected to find endnote content in at least one comment"


# endregion
