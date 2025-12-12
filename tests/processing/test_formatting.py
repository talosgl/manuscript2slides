"""Tests for standard formatting"""

# tests/test_formatting.py

# pyright: reportPrivateUsage=false
# pyright: reportAttributeAccessIssue=false
# pyright: reportIndexIssue=false
# pyright: reportOptionalMemberAccess=false
import logging
from pathlib import Path

import pytest
from docx import Document, document
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_COLOR_INDEX
from docx.shared import RGBColor as RGBColor_docx
from docx.text.font import Font as Font_docx
from docx.text.paragraph import Paragraph as Paragraph_docx
from docx.text.run import Run as Run_docx
from pptx import Presentation, presentation
from pptx.dml.color import RGBColor as RGBColor_pptx
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.text.text import Font as Font_pptx
from pptx.text.text import TextFrame
from pptx.text.text import _Paragraph as Paragraph_pptx
from pptx.text.text import _Run as Run_pptx

from manuscript2slides.internals.define_config import UserConfig
from manuscript2slides.processing import formatting

# region fixtures


@pytest.fixture
def blank_docx(path_to_empty_docx: Path) -> document.Document:
    """An empty docx object."""
    return Document(str(path_to_empty_docx))


# sanity checks
def test_blank_docx_is_fresh_per_test(blank_docx: document.Document) -> None:
    """Verify that blank_docx fixture creates a new instance for each test"""
    # Add a paragraph to "mutate" it2
    blank_docx.add_paragraph("Test mutation")
    # docx requires at least 1 paragraph to save, so the starting count was 1
    assert len(blank_docx.paragraphs) == 2


def test_blank_docx_is_fresh_per_test_2(blank_docx: document.Document) -> None:
    """This test should get a fresh blank_docx without the paragraph from the previous test"""
    # If the fixture is truly fresh, there should be 1 paragraphs
    # docx requires at least 1 paragraph to save, so the starting count was 1
    assert len(blank_docx.paragraphs) == 1


@pytest.fixture
def pptx_w_twenty_empty_slides(
    path_to_pptx_w_twenty_empty_slides: Path,
) -> presentation.Presentation:
    """An empty pptx object that contains 20 empty slides so we don't need to add slides or slide layout stuff to be adding paragraphs and runs.

    NOTE: Returns a fresh Presentation object for each test to avoid mutation/pollution across tests.
    This is because by default, pytest fixtures with no explicit scope parameter use scope="function",
    which means they create a new instance for each test.
    """
    return Presentation(path_to_pptx_w_twenty_empty_slides)


@pytest.fixture
def docx_formatting_runs_dict(
    path_to_sample_docx_with_formatting: Path,
) -> dict[str, Run_docx]:
    """Pre-validated runs from test_formatting.docx for easy access."""

    docx_with_formatting = Document(str(path_to_sample_docx_with_formatting))

    test_err = "Test cannot proceed because test document does not match expectations; "

    # Bold, italics, underline
    bolded_run = docx_with_formatting.paragraphs[2].runs[0]
    assert bolded_run.font.bold is True, test_err + "expected bold run at para 2, run 0"

    italic_run = docx_with_formatting.paragraphs[3].runs[0]
    assert italic_run.font.italic is True, (
        test_err + "expected italic in run at para 3, run 0"
    )

    bold_and_italic_run = docx_with_formatting.paragraphs[4].runs[0]
    assert (
        bold_and_italic_run.font.bold is True
        and bold_and_italic_run.font.italic is True
    ), (test_err + "expected italic + bold at para 4, run 0")

    underline_run = docx_with_formatting.paragraphs[5].runs[0]
    assert underline_run.font.underline is True, (
        test_err + "expected underline at para 5, run 0"
    )

    bold_and_underline_run = docx_with_formatting.paragraphs[6].runs[0]
    assert (
        bold_and_underline_run.font.bold is True
        and bold_and_underline_run.font.underline is True
    ), (test_err + "expected bold + underline at para 6, run 0")

    underline_and_italic_run = docx_with_formatting.paragraphs[7].runs[0]
    assert (
        underline_and_italic_run.font.underline is True
        and underline_and_italic_run.font.italic is True
    ), (test_err + "expected underline + italic at para 7, run 0")

    all_three_run = docx_with_formatting.paragraphs[8].runs[0]
    assert (
        all_three_run.font.bold is True
        and all_three_run.font.italic is True
        and all_three_run.font.underline is True
    ), (test_err + "expected bold + italic + underline at para 8, run 0")

    # font.name or Typeface
    typeface_run = docx_with_formatting.paragraphs[9].runs[0]
    assert typeface_run.font.name is not None and typeface_run.font.name == "Georgia", (
        test_err + "expected font.name to be set to 'Georgia' at para 9, run 0"
    )

    # Size
    large_size_run = docx_with_formatting.paragraphs[11].runs[0]
    assert large_size_run.font.size.pt == 48

    small_size_run = docx_with_formatting.paragraphs[12].runs[0]
    assert small_size_run.font.size.pt == 8

    # Color
    color_para = docx_with_formatting.paragraphs[14]
    red_run = color_para.runs[1]
    assert red_run.font.color.rgb is not None, (
        test_err + "expected font.color.rgb to be set at para 14, run 1"
    )
    assert red_run.font.color.rgb == RGBColor_docx(0xFF, 0x00, 0x00), (
        test_err
        + "expected font.color.rgb to be red (RGBColor_docx(0xFF, 0x00, 0x00)) at para 14, run 1"
    )

    blue_run = color_para.runs[3]
    assert blue_run.font.color.rgb is not None, (
        test_err + "expected font.color.rgb at para 14, run 3"
    )
    assert blue_run.font.color.rgb == RGBColor_docx(0x44, 0x72, 0xC4), (
        test_err
        + "expected font.color.rgb to be blue RGBColor_docx(0x44, 0x72, 0xC4)) at para 14, run 1"
    )

    # Experimental formatting
    hl_run = docx_with_formatting.paragraphs[34].runs[0]
    assert hl_run.font.highlight_color is not None, (
        test_err + "expected highlight_color at para 34, run 0"
    )

    complex_hl_para = docx_with_formatting.paragraphs[35]
    assert "highlighted, but also" in complex_hl_para.text, (
        test_err + "expected text 'highlighted, but also' in para 35"
    )
    hl_run2 = complex_hl_para.runs[0]
    assert hl_run2.font.highlight_color is not None, (
        test_err + "expected highlight_color at para 35, run 0"
    )

    hl_and_underlined = complex_hl_para.runs[1]
    assert (
        hl_and_underlined.font.highlight_color is not None
        and hl_and_underlined.font.underline is True
    ), (test_err + "expected highlight + underline at para 35, run 1")

    hl_only = complex_hl_para.runs[2]
    assert hl_only.font.highlight_color is not None, (
        test_err + "expected highlight_color at para 35, run 2"
    )

    hl_and_italic = complex_hl_para.runs[3]
    assert (
        hl_and_italic.font.highlight_color is not None
        and hl_and_italic.font.italic is True
    ), (test_err + "expected highlight + italic at para 35, run 3")

    hl_and_bold = complex_hl_para.runs[5]
    assert (
        hl_and_bold.font.highlight_color is not None and hl_and_bold.font.bold is True
    ), (test_err + "expected highlight + bold at para 35, run 5")

    super_sub_para = docx_with_formatting.paragraphs[36]
    assert "superscript" in super_sub_para.text, (
        test_err + "expected text 'superscript' in para 36"
    )
    subscript_run = super_sub_para.runs[1]
    assert subscript_run.font.subscript is True, (
        test_err + "expected subscript at para 36, run 1"
    )

    superscript_run = super_sub_para.runs[3]
    assert superscript_run.font.superscript is True, (
        test_err + "expected superscript at para 36, run 3"
    )

    all_caps_run = docx_with_formatting.paragraphs[37].runs[0]
    assert all_caps_run.font.all_caps is True, (
        test_err + "expected all_caps at para 37, run 0"
    )

    small_caps_run = docx_with_formatting.paragraphs[38].runs[0]
    assert small_caps_run.font.small_caps is True, (
        test_err + "expected small_caps at para 38, run 0"
    )

    strike_para = docx_with_formatting.paragraphs[39]
    single_str_run = strike_para.runs[1]
    assert single_str_run.font.strike is True, (
        test_err + "expected strike at para 39, run 1"
    )

    dbl_str_run = strike_para.runs[3]
    assert dbl_str_run.font.double_strike is True, (
        test_err + "expected double_strike at para 39, run 3"
    )

    return {
        "bold": bolded_run,
        "italic": italic_run,
        "bold_and_italic": bold_and_italic_run,
        "underline": underline_run,
        "bold_and_underline": bold_and_underline_run,
        "underline_and_italic": underline_and_italic_run,
        "all_three": all_three_run,
        "typeface": typeface_run,
        "red_run": red_run,
        "blue_run": blue_run,
        "large_size": large_size_run,
        "small_size": small_size_run,
        "highlight": hl_run,
        "highlight_2": hl_run2,
        "highlight_and_underline": hl_and_underlined,
        "highlight_only": hl_only,
        "highlight_and_italic": hl_and_italic,
        "highlight_and_bold": hl_and_bold,
        "subscript": subscript_run,
        "superscript": superscript_run,
        "all_caps": all_caps_run,
        "small_caps": small_caps_run,
        "single_strike": single_str_run,
        "double_strike": dbl_str_run,
    }


@pytest.fixture
def pptx_formatting_runs_dict(
    path_to_sample_pptx_with_formatting: Path,
) -> dict[str, Run_pptx]:
    """Pre-validated runs from test_formatting_expected_output.pptx for easy access."""
    pptx_with_formatting = Presentation(path_to_sample_pptx_with_formatting)
    test_err = "Test cannot proceed because test document does not match expectations; "

    bold_slide = pptx_with_formatting.slides[1]
    bold_paragraphs = bold_slide.shapes.placeholders[1].text_frame.paragraphs
    bold_run = bold_paragraphs[0].runs[0]
    assert bold_run.font.bold == True, (
        test_err + "expected bold run at slide 1, placeholder 1, para 0, run 0"
    )

    italic_run = (
        pptx_with_formatting.slides[2]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert italic_run.font.italic is True, (
        test_err + "expected italic run at slide 2, placeholder 1, para 0, run 0"
    )

    bold_and_italic_run = (
        pptx_with_formatting.slides[3]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert (
        bold_and_italic_run.font.bold is True
        and bold_and_italic_run.font.italic is True
    ), (test_err + "expected bold + italic at slide 3, placeholder 1, para 0, run 0")

    underline_run = (
        pptx_with_formatting.slides[4]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert underline_run.font.underline is True, (
        test_err + "expected underline at slide 4, placeholder 1, para 0, run 0"
    )

    bold_and_underline_run = (
        pptx_with_formatting.slides[5]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert (
        bold_and_underline_run.font.bold is True
        and bold_and_underline_run.font.underline is True
    ), (test_err + "expected bold + underline at slide 5, placeholder 1, para 0, run 0")

    underline_and_italic_run = (
        pptx_with_formatting.slides[6]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert (
        underline_and_italic_run.font.underline is True
        and underline_and_italic_run.font.italic is True
    ), (
        test_err
        + "expected underline + italic at slide 6, placeholder 1, para 0, run 0"
    )

    all_three_run = (
        pptx_with_formatting.slides[7]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert (
        all_three_run.font.bold is True
        and all_three_run.font.italic is True
        and all_three_run.font.underline is True
    ), (
        test_err
        + "expected bold + italic + underline at slide 7, placeholder 1, para 0, run 0"
    )

    typeface_run = (
        pptx_with_formatting.slides[8]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert typeface_run.font.name is not None and typeface_run.font.name == "Georgia", (
        test_err
        + "expected font.name to be 'Georgia' at slide 8, placeholder 1, para 0, run 0"
    )

    # skip 9

    large_size_run = (
        pptx_with_formatting.slides[10]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert large_size_run.font.size.pt == 48, (
        test_err + "expected font size 48pt at slide 10, placeholder 1, para 0, run 0"
    )

    small_size_run = (
        pptx_with_formatting.slides[11]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert small_size_run.font.size.pt == 8, (
        test_err + "expected font size 8pt at slide 11, placeholder 1, para 0, run 0"
    )

    # skip 12

    # Color - assuming similar to docx where there are multiple runs
    color_para = (
        pptx_with_formatting.slides[13].shapes.placeholders[1].text_frame.paragraphs[0]
    )
    red_run = color_para.runs[1]  # guessing run 1 for red
    assert red_run.font.color.rgb is not None and red_run.font.color.rgb == (
        255,
        0,
        0,
    ), (
        test_err + "expected red color at slide 13, placeholder 1, para 0, run 1"
    )
    blue_run = color_para.runs[3]  # guessing run 3 for blue
    assert blue_run.font.color.rgb is not None and blue_run.font.color.rgb == (
        68,
        114,
        196,
    ), (
        test_err + "expected blue color at slide 13, placeholder 1, para 0, run 3"
    )

    # skip 14

    # Experimental formatting - using XML assertions
    hl_run = (
        pptx_with_formatting.slides[15]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    # Check for highlight in XML since pptx doesn't expose as property
    hl_xml = hl_run._r.xml
    assert "a:highlight" in hl_xml, (
        test_err
        + "expected highlight element in XML at slide 15, placeholder 1, para 0, run 0"
    )

    # Complex highlight paragraph
    hl_para = (
        pptx_with_formatting.slides[16].shapes.placeholders[1].text_frame.paragraphs[0]
    )
    # Get the highlighted runs from this paragraph
    hl_run2 = hl_para.runs[0]
    assert "a:highlight" in hl_run2._r.xml, (
        test_err + "expected highlight in run 0 at slide 16"
    )

    hl_and_underlined = hl_para.runs[1]
    assert (
        "a:highlight" in hl_and_underlined._r.xml
        and hl_and_underlined.font.underline is True
    ), (test_err + "expected highlight + underline at slide 16, run 1")

    super_sub_para = (
        pptx_with_formatting.slides[17].shapes.placeholders[1].text_frame.paragraphs[0]
    )
    subscript_run = super_sub_para.runs[1]
    baseline_sub = subscript_run.font._element.get("baseline")
    assert baseline_sub is not None and int(baseline_sub) < 0, (
        test_err + "expected negative baseline for subscript at slide 17, para 0, run 1"
    )

    superscript_run = super_sub_para.runs[3]
    baseline_super = superscript_run.font._element.get("baseline")
    assert baseline_super is not None and int(baseline_super) > 0, (
        test_err
        + "expected positive baseline for superscript at slide 17, para 0, run 3"
    )

    # All caps
    all_caps_run = (
        pptx_with_formatting.slides[18]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert all_caps_run.font._element.get("cap") == "all", (
        test_err
        + "expected cap='all' attribute at slide 18, placeholder 1, para 0, run 0"
    )

    # Small caps
    small_caps_run = (
        pptx_with_formatting.slides[19]
        .shapes.placeholders[1]
        .text_frame.paragraphs[0]
        .runs[0]
    )
    assert small_caps_run.font._element.get("cap") == "small", (
        test_err
        + "expected cap='small' attribute at slide 19, placeholder 1, para 0, run 0"
    )

    # Strike paragraph
    strike_para = (
        pptx_with_formatting.slides[20].shapes.placeholders[1].text_frame.paragraphs[0]
    )
    single_str_run = strike_para.runs[1]
    assert single_str_run.font._element.get("strike") == "sngStrike", (
        test_err + "expected strike='sngStrike' attribute at slide 20, para 0, run 1"
    )

    dbl_str_run = strike_para.runs[3]
    assert dbl_str_run.font._element.get("strike") == "dblStrike", (
        test_err + "expected strike='dblStrike' attribute at slide 20, para 0, run 3"
    )

    return {
        "bold": bold_run,
        "italic": italic_run,
        "bold_and_italic": bold_and_italic_run,
        "underline": underline_run,
        "bold_and_underline": bold_and_underline_run,
        "underline_and_italic": underline_and_italic_run,
        "all_three": all_three_run,
        "typeface": typeface_run,
        "large_size": large_size_run,
        "small_size": small_size_run,
        "red_run": red_run,
        "blue_run": blue_run,
        "highlight": hl_run,
        "highlight_2": hl_run2,
        "highlight_and_underline": hl_and_underlined,
        "subscript": subscript_run,
        "superscript": superscript_run,
        "all_caps": all_caps_run,
        "small_caps": small_caps_run,
        "single_strike": single_str_run,
        "double_strike": dbl_str_run,
    }


# endregion


# region test helpers


def _create_target_pptx_run(pptx_obj: presentation.Presentation) -> Run_pptx:
    """Helper to create a target pptx run for testing."""
    slide = pptx_obj.slides[0]
    text_frame = slide.shapes.placeholders[1].text_frame
    target_para = text_frame.paragraphs[0]
    target_run = target_para.add_run()
    target_run.text = "Test pptx Run"
    return target_run


def _create_target_docx_run(docx_obj: document.Document) -> Run_docx:
    """Helper to create a target docx run for testing."""
    paragraph = docx_obj.add_paragraph()
    target_run = paragraph.add_run("Test docx Run. ")
    return target_run


# endregion


# region docx2pptx RUN tests


# region _copy_basic_font_formatting tests


@pytest.mark.parametrize(
    "run_key,expected_attrs",
    [
        ("typeface", {"name": "Georgia"}),
        ("bold", {"bold": True}),
        ("italic", {"italic": True}),
        ("underline", {"underline": True}),
        ("bold_and_italic", {"bold": True, "italic": True}),
        ("bold_and_underline", {"bold": True, "underline": True}),
        ("underline_and_italic", {"underline": True, "italic": True}),
        ("all_three", {"bold": True, "italic": True, "underline": True}),
    ],
)
def test_copy_basic_font_formatting_docx2pptx(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
    run_key: str,
    expected_attrs: dict,
) -> None:
    """Test that _copy_basic_font_formatting correctly copies formatting from docx to pptx."""
    source_run = docx_formatting_runs_dict[run_key]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)

    formatting._copy_basic_font_formatting(source_run.font, target_run.font)

    for attr, expected_value in expected_attrs.items():
        actual_value = getattr(target_run.font, attr)
        assert actual_value is not None and actual_value == expected_value


# endregion

# region _copy_font_size_formatting tests


@pytest.mark.parametrize(
    "run_key,expected_size_pt",
    [
        ("large_size", 48),
        ("small_size", 8),
    ],
)
def test_copy_font_size_formatting_docx2pptx(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
    run_key: str,
    expected_size_pt: int,
) -> None:
    """Test that _copy_font_size_formatting correctly copies font size from docx to pptx."""
    source_run = docx_formatting_runs_dict[run_key]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)

    formatting._copy_font_size_formatting(source_run.font, target_run.font)

    assert target_run.font.size.pt == expected_size_pt


# endregion

# region _copy_font_color_formatting tests


def test_copy_font_color_formatting_preserves_color(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that _copy_font_color_formatting correctly copies colors."""
    source_run = docx_formatting_runs_dict["red_run"]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)

    formatting._copy_font_color_formatting(source_run.font, target_run.font)

    assert target_run.font.color.rgb is not None and target_run.font.color.rgb == (
        255,
        0,
        0,
    )


# endregion

# region _copy_experimental_formatting_docx2pptx tests


def test_copy_experimental_formatting_docx2pptx_preserves_highlight(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that experimental formatting correctly copies highlight from docx to pptx."""
    source_run = docx_formatting_runs_dict["highlight"]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    metadata = []

    formatting._copy_experimental_formatting_docx2pptx(source_run, target_run, metadata)

    # Check XML contains highlight
    assert "a:highlight" in target_run._r.xml
    # Check metadata was recorded
    assert len(metadata) == 1
    assert metadata[0]["formatting_type"] == "highlight"
    assert metadata[0]["ref_text"] == source_run.text
    assert metadata[0]["highlight_color_enum"] == source_run.font.highlight_color.name


def test_copy_experimental_formatting_docx2pptx_preserves_single_strike(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that experimental formatting correctly copies single strike from docx to pptx."""
    source_run = docx_formatting_runs_dict["single_strike"]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    metadata = []

    formatting._copy_experimental_formatting_docx2pptx(source_run, target_run, metadata)

    # Check XML attribute
    assert target_run.font._element.get("strike") == "sngStrike"
    # Check metadata was recorded
    assert len(metadata) == 1
    assert metadata[0]["formatting_type"] == "strike"
    assert metadata[0]["ref_text"] == source_run.text


def test_copy_experimental_formatting_docx2pptx_preserves_double_strike(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that experimental formatting correctly copies double strike from docx to pptx."""
    source_run = docx_formatting_runs_dict["double_strike"]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    metadata = []

    formatting._copy_experimental_formatting_docx2pptx(source_run, target_run, metadata)

    # Check XML attribute
    assert target_run.font._element.get("strike") == "dblStrike"
    # Check metadata was recorded
    assert len(metadata) == 1
    assert metadata[0]["formatting_type"] == "double_strike"
    assert metadata[0]["ref_text"] == source_run.text


def test_copy_experimental_formatting_docx2pptx_preserves_subscript(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that experimental formatting correctly copies subscript from docx to pptx."""
    source_run = docx_formatting_runs_dict["subscript"]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    metadata = []

    formatting._copy_experimental_formatting_docx2pptx(source_run, target_run, metadata)

    # Check XML attribute (negative baseline)
    baseline = target_run.font._element.get("baseline")
    assert baseline is not None and int(baseline) < 0
    # Check metadata was recorded
    assert len(metadata) == 1
    assert metadata[0]["formatting_type"] == "subscript"
    assert metadata[0]["ref_text"] == source_run.text


def test_copy_experimental_formatting_docx2pptx_preserves_superscript(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that experimental formatting correctly copies superscript from docx to pptx."""
    source_run = docx_formatting_runs_dict["superscript"]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    metadata = []

    formatting._copy_experimental_formatting_docx2pptx(source_run, target_run, metadata)

    # Check XML attribute (positive baseline)
    baseline = target_run.font._element.get("baseline")
    assert baseline is not None and int(baseline) > 0
    # Check metadata was recorded
    assert len(metadata) == 1
    assert metadata[0]["formatting_type"] == "superscript"
    assert metadata[0]["ref_text"] == source_run.text


def test_copy_experimental_formatting_docx2pptx_preserves_all_caps(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that experimental formatting correctly copies all caps from docx to pptx."""
    source_run = docx_formatting_runs_dict["all_caps"]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    metadata = []

    formatting._copy_experimental_formatting_docx2pptx(source_run, target_run, metadata)

    # Check XML attribute
    assert target_run.font._element.get("cap") == "all"
    # Check metadata was recorded
    assert len(metadata) == 1
    assert metadata[0]["formatting_type"] == "all_caps"
    assert metadata[0]["ref_text"] == source_run.text


def test_copy_experimental_formatting_docx2pptx_preserves_small_caps(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that experimental formatting correctly copies small caps from docx to pptx."""
    source_run = docx_formatting_runs_dict["small_caps"]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    metadata = []

    formatting._copy_experimental_formatting_docx2pptx(source_run, target_run, metadata)

    # Check XML attribute
    assert target_run.font._element.get("cap") == "small"
    # Check metadata was recorded
    assert len(metadata) == 1
    assert metadata[0]["formatting_type"] == "small_caps"
    assert metadata[0]["ref_text"] == source_run.text


# endregion

# region copy_run_formatting_docx2pptx tests basics


def test_copy_run_formatting_docx2pptx_copies_all_basic_formatting(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that copy_run_formatting_docx2pptx copies typeface, bold, italic, and underline formatting attributes."""

    # Use a run with multiple formatting attributes
    source_run = docx_formatting_runs_dict["all_three"]  # bold + italic + underline
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)

    formatting.copy_run_formatting_docx2pptx(
        source_run, target_run, [], UserConfig(experimental_formatting_on=False)
    )

    assert target_run.font.bold is True
    assert target_run.font.italic is True
    assert target_run.font.underline is True


def test_copy_run_formatting_docx2pptx_copies_typeface(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that orchestrator copy_run_formatting_docx2pptx properly copies typeface."""

    source_run = docx_formatting_runs_dict["typeface"]
    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)

    formatting.copy_run_formatting_docx2pptx(
        source_run, target_run, [], UserConfig(experimental_formatting_on=False)
    )

    assert target_run.font.name is not None and target_run.font.name == "Georgia"


def test_copy_run_formatting_docx2pptx_copies_size_formatting(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that orchestrator copy_run_formatting_docx2pptx properly copies size."""

    small_size_source_run = docx_formatting_runs_dict["small_size"]
    small_size_target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    formatting.copy_run_formatting_docx2pptx(
        small_size_source_run,
        small_size_target_run,
        [],
        UserConfig(experimental_formatting_on=False),
    )

    large_size_source_run = docx_formatting_runs_dict["large_size"]
    large_size_target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    formatting.copy_run_formatting_docx2pptx(
        large_size_source_run,
        large_size_target_run,
        [],
        UserConfig(experimental_formatting_on=False),
    )

    assert small_size_target_run.font.size.pt == 8
    assert large_size_target_run.font.size.pt == 48


def test_copy_run_formatting_docx2pptx_copies_color_formatting(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that orchestrator copy_run_formatting_docx2pptx properly copies color."""

    red_source_run = docx_formatting_runs_dict["red_run"]
    red_target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    formatting.copy_run_formatting_docx2pptx(
        red_source_run,
        red_target_run,
        [],
        UserConfig(experimental_formatting_on=False),
    )

    blue_source_run = docx_formatting_runs_dict["blue_run"]
    blue_target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    formatting.copy_run_formatting_docx2pptx(
        blue_source_run,
        blue_target_run,
        [],
        UserConfig(experimental_formatting_on=False),
    )

    assert (
        red_target_run.font.color.rgb is not None
        and red_target_run.font.color.rgb == (255, 0, 0)
    )
    assert (
        blue_target_run.font.color.rgb is not None
        and blue_target_run.font.color.rgb == (68, 114, 196)
    )


# endregion

# region copy_run_formatting_docx2pptx experimental tests


def test_copy_run_formatting_docx2pptx_copies_experimental_when_enabled(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that orchestrator copy_run_formatting_docx2pptx copies experimental formatting when enabled."""
    # Test with highlight
    hl_source = docx_formatting_runs_dict["highlight"]
    hl_target = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    hl_metadata = []
    formatting.copy_run_formatting_docx2pptx(
        hl_source, hl_target, hl_metadata, UserConfig(experimental_formatting_on=True)
    )
    assert "a:highlight" in hl_target._r.xml
    assert len(hl_metadata) == 1

    # Test with subscript
    sub_source = docx_formatting_runs_dict["subscript"]
    sub_target = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    sub_metadata = []
    formatting.copy_run_formatting_docx2pptx(
        sub_source,
        sub_target,
        sub_metadata,
        UserConfig(experimental_formatting_on=True),
    )
    baseline = sub_target.font._element.get("baseline")
    assert baseline is not None and int(baseline) < 0
    assert len(sub_metadata) == 1

    # Test with all caps
    caps_source = docx_formatting_runs_dict["all_caps"]
    caps_target = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    caps_metadata = []
    formatting.copy_run_formatting_docx2pptx(
        caps_source,
        caps_target,
        caps_metadata,
        UserConfig(experimental_formatting_on=True),
    )
    assert caps_target.font._element.get("cap") == "all"
    assert len(caps_metadata) == 1


def test_copy_run_formatting_docx2pptx_skips_experimental_when_disabled(
    docx_formatting_runs_dict: dict[str, Run_docx],
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that orchestrator copy_run_formatting_docx2pptx skips experimental formatting when disabled."""
    # Test with highlight
    hl_source = docx_formatting_runs_dict["highlight"]
    hl_target = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    hl_metadata = []
    formatting.copy_run_formatting_docx2pptx(
        hl_source, hl_target, hl_metadata, UserConfig(experimental_formatting_on=False)
    )
    assert "a:highlight" not in hl_target._r.xml
    assert len(hl_metadata) == 0

    # Test with subscript
    sub_source = docx_formatting_runs_dict["subscript"]
    sub_target = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    sub_metadata = []
    formatting.copy_run_formatting_docx2pptx(
        sub_source,
        sub_target,
        sub_metadata,
        UserConfig(experimental_formatting_on=False),
    )
    baseline = sub_target.font._element.get("baseline")
    assert baseline is None
    assert len(sub_metadata) == 0


# endregion


# endregion docx2pptx RUN tests

# region docx2pptx paragraph formatting

# region copy_paragraph_formatting_docx2pptx


def test_copy_paragraph_formatting_docx2pptx_happy_path(
    path_to_sample_docx_with_formatting: Path,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that paragraph formatting is copied as expected in the docx2pptx direction;
    right now this only includes color, size, and alignment."""
    # Arrange
    docx_with_formatting = Document(str(path_to_sample_docx_with_formatting))

    # This paragraph has these 3 properties, just alias so it it is easier to reference.
    # (Also we can split to reference other paragraphs in future but keep the var names.)
    para_with_color = para_with_size = para_with_center_alignment = (
        docx_with_formatting.paragraphs[1]
    )

    slide = pptx_w_twenty_empty_slides.slides[0]
    text_frame: TextFrame = slide.shapes.placeholders[1].text_frame

    target_para_for_color = text_frame.add_paragraph()
    target_para_for_size = text_frame.add_paragraph()
    target_para_for_alignment = text_frame.add_paragraph()

    # Action
    formatting.copy_paragraph_formatting_docx2pptx(
        source_para=para_with_color, target_para=target_para_for_color
    )
    formatting.copy_paragraph_formatting_docx2pptx(
        source_para=para_with_size, target_para=target_para_for_size
    )
    formatting.copy_paragraph_formatting_docx2pptx(
        source_para=para_with_center_alignment, target_para=target_para_for_alignment
    )

    # Assert
    assert (
        target_para_for_color.font.color.rgb is not None
        and target_para_for_color.font.color.rgb
        == (47, 84, 150)  # blue from the test_formatting.docx theme
    )
    assert (
        target_para_for_size.font.size.pt is not None
        and target_para_for_size.font.size.pt == 16  # expected size from heading 1
    )
    assert (
        target_para_for_alignment.alignment is not None
        and target_para_for_alignment.alignment == PP_ALIGN.CENTER
    )


def test_copy_paragraph_formatting_docx2pptx_advanced(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """
    Test that typeface and bold/italic/underline from paragraph styles work for
    docx2pptx when explicitly set on style.font.

    NOTE: I can't for the life of me figure out how to actually make a document HAVE
    these properties at the paragraph level (they're typically run-level). However,
    the code DOES attempt to copy these from source_para.style.font, so we test that
    it works when explicitly set.
    """
    # Arrange - Create fresh paragraphs and explicitly set style font properties
    para_with_bold = blank_docx.add_paragraph("Bold paragraph")
    para_with_bold.style.font.bold = True

    para_with_italic = blank_docx.add_paragraph("Italic paragraph")
    para_with_italic.style.font.italic = True

    para_with_underline = blank_docx.add_paragraph("Underline paragraph")
    para_with_underline.style.font.underline = True

    para_with_typeface = blank_docx.add_paragraph("Georgia paragraph")
    para_with_typeface.style.font.name = "Georgia"

    slide = pptx_w_twenty_empty_slides.slides[0]
    text_frame: TextFrame = slide.shapes.placeholders[1].text_frame

    target_para_for_bold = text_frame.add_paragraph()
    target_para_for_italic = text_frame.add_paragraph()
    target_para_for_underline = text_frame.add_paragraph()
    target_para_for_typeface = text_frame.add_paragraph()

    # Act
    formatting.copy_paragraph_formatting_docx2pptx(
        source_para=para_with_bold, target_para=target_para_for_bold
    )
    formatting.copy_paragraph_formatting_docx2pptx(
        source_para=para_with_italic, target_para=target_para_for_italic
    )
    formatting.copy_paragraph_formatting_docx2pptx(
        source_para=para_with_underline, target_para=target_para_for_underline
    )
    formatting.copy_paragraph_formatting_docx2pptx(
        source_para=para_with_typeface, target_para=target_para_for_typeface
    )

    # Assert
    assert target_para_for_bold.font.bold is True
    assert target_para_for_italic.font.italic is True
    assert target_para_for_underline.font.underline is True
    assert target_para_for_typeface.font.name == "Georgia"


# endregion

# region _copy_paragraph_font_name_docx2pptx


def test_copy_paragraph_font_name_docx2pptx_resolves_theme_fonts(
    path_to_sample_docx_with_formatting: Path,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that _copy_paragraph_font_name_docx2pptx resolves theme fonts correctly."""
    # Arrange
    docx_with_formatting = Document(str(path_to_sample_docx_with_formatting))
    # Paragraph 1 is Heading 1 which uses the Major theme font (Calibri Light)
    para_with_heading_style_from_theme = docx_with_formatting.paragraphs[1]

    slide = pptx_w_twenty_empty_slides.slides[0]
    text_frame: TextFrame = slide.shapes.placeholders[1].text_frame
    target_para = text_frame.add_paragraph()

    # Act
    formatting._copy_paragraph_font_name_docx2pptx(
        para_with_heading_style_from_theme, target_para
    )

    # Assert - Should resolve majorHAnsi to "Calibri Light" from the theme
    assert target_para.font.name == "Calibri Light"


# endregion

# region _copy_paragraph_alignment_docx2pptx


def test_copy_paragraph_alignment_docx2pptx_from_style(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that paragraph alignment is copied from style when no direct formatting is applied."""
    # Arrange - Create a paragraph with center alignment from style
    source_para = blank_docx.add_paragraph("Test text")
    # Set alignment on the style's paragraph format
    source_para.style.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

    slide = pptx_w_twenty_empty_slides.slides[0]
    text_frame: TextFrame = slide.shapes.placeholders[1].text_frame
    target_para = text_frame.add_paragraph()

    # Act
    formatting._copy_paragraph_alignment_docx2pptx(source_para, target_para)

    # Assert - Should map WD_ALIGN_PARAGRAPH.CENTER to PP_ALIGN.CENTER (enum value 2)
    assert target_para.alignment == PP_ALIGN.CENTER


def test_copy_paragraph_alignment_docx2pptx_direct_formatting_overrides_style(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that direct paragraph alignment overrides style alignment (higher priority)."""
    # Arrange - Create a paragraph with CENTER alignment in style but RIGHT in direct formatting
    source_para = blank_docx.add_paragraph("Test text")
    # Set alignment on the style (lower priority)
    source_para.style.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER
    # Set direct formatting (higher priority - should win)
    source_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    slide = pptx_w_twenty_empty_slides.slides[0]
    text_frame: TextFrame = slide.shapes.placeholders[1].text_frame
    target_para = text_frame.add_paragraph()

    # Act
    formatting._copy_paragraph_alignment_docx2pptx(source_para, target_para)

    # Assert - Should use direct formatting (RIGHT), not style (CENTER)
    assert target_para.alignment == PP_ALIGN.RIGHT


def test_copy_paragraph_alignment_docx2pptx_no_alignment_set(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that paragraph alignment copy handles None/unset alignment gracefully."""
    # Arrange - Create a paragraph with no explicit alignment
    source_para = blank_docx.add_paragraph("Test text")

    slide = pptx_w_twenty_empty_slides.slides[0]
    text_frame: TextFrame = slide.shapes.placeholders[1].text_frame
    target_para = text_frame.add_paragraph()

    # Act
    formatting._copy_paragraph_alignment_docx2pptx(source_para, target_para)

    # Assert - Alignment should remain None or unchanged
    # (The function doesn't set alignment if source has none)
    assert target_para.alignment is None or target_para.alignment == PP_ALIGN.LEFT


# endregion


# region get_style_font_name_with_fallback_docx
def test_get_effective_font_name_docx_explicit_font(
    blank_docx: document.Document,
) -> None:
    """Test that the function returns explicit font names from a style."""
    # Arrange - Create a paragraph with Normal style and set explicit font
    para = blank_docx.add_paragraph("Test")
    para.style.font.name = "Georgia"

    # Act
    assert para.style is not None
    result = formatting.get_effective_font_name_docx(para.style)

    # Assert
    assert result == "Georgia"


def test_get_effective_font_name_docx_traverses_style_hierarchy(
    blank_docx: document.Document,
) -> None:
    """Test that the function traverses the style hierarchy to find font names."""
    # Arrange - Create a base style with explicit font, then a derived style without one
    # Create base style with explicit font
    base_style = blank_docx.styles.add_style("BaseStyle", WD_STYLE_TYPE.PARAGRAPH)
    base_style.font.name = "Georgia"

    # Create a custom style based on BaseStyle without setting its own font
    custom_style = blank_docx.styles.add_style("CustomStyle", WD_STYLE_TYPE.PARAGRAPH)
    custom_style.base_style = base_style
    # Don't set custom_style.font.name - it should inherit from base_style

    para = blank_docx.add_paragraph("Test", style="CustomStyle")

    # Act
    assert para.style is not None
    result = formatting.get_effective_font_name_docx(para.style)

    # Assert - Should find "Georgia" from the base style
    assert result == "Georgia"


# endregion

# endregion


# region pptx2docx RUN tests


# region _copy_basic_font_formatting tests (pptx2docx)


@pytest.mark.parametrize(
    "run_key,expected_attrs",
    [
        ("typeface", {"name": "Georgia"}),
        ("bold", {"bold": True}),
        ("italic", {"italic": True}),
        ("underline", {"underline": True}),
        ("bold_and_italic", {"bold": True, "italic": True}),
        ("bold_and_underline", {"bold": True, "underline": True}),
        ("underline_and_italic", {"underline": True, "italic": True}),
        ("all_three", {"bold": True, "italic": True, "underline": True}),
    ],
)
def test_copy_basic_font_formatting_pptx2docx(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
    run_key: str,
    expected_attrs: dict,
) -> None:
    """Test that _copy_basic_font_formatting correctly copies formatting from pptx to docx."""
    source_run = pptx_formatting_runs_dict[run_key]
    target_run = _create_target_docx_run(blank_docx)

    formatting._copy_basic_font_formatting(source_run.font, target_run.font)

    for attr, expected_value in expected_attrs.items():
        actual_value = getattr(target_run.font, attr)
        assert actual_value is not None and actual_value == expected_value


# endregion

# region _copy_font_size_formatting tests (pptx2docx)


@pytest.mark.parametrize(
    "run_key,expected_size_pt",
    [
        ("large_size", 48),
        ("small_size", 8),
    ],
)
def test_copy_font_size_formatting_pptx2docx(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
    run_key: str,
    expected_size_pt: int,
) -> None:
    """Test that _copy_font_size_formatting correctly copies font size from pptx to docx."""
    source_run = pptx_formatting_runs_dict[run_key]
    target_run = _create_target_docx_run(blank_docx)

    formatting._copy_font_size_formatting(source_run.font, target_run.font)

    assert target_run.font.size.pt == expected_size_pt


# endregion

# region _copy_font_color_formatting tests (pptx2docx)


def test_copy_font_color_formatting_pptx2docx_preserves_color(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that _copy_font_color_formatting correctly copies colors from pptx to docx."""
    source_run = pptx_formatting_runs_dict["red_run"]
    target_run = _create_target_docx_run(blank_docx)

    formatting._copy_font_color_formatting(source_run.font, target_run.font)

    assert target_run.font.color.rgb is not None and target_run.font.color.rgb == (
        255,
        0,
        0,
    )


# endregion

# region copy_run_formatting_pptx2docx tests basics


def test_copy_run_formatting_pptx2docx_copies_all_basic_formatting(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that copy_run_formatting_pptx2docx copies typeface, bold, italic, and underline formatting attributes."""

    # Use a run with multiple formatting attributes
    source_run = pptx_formatting_runs_dict["all_three"]  # bold + italic + underline
    target_run = _create_target_docx_run(blank_docx)

    formatting.copy_run_formatting_pptx2docx(
        source_run, target_run, UserConfig(experimental_formatting_on=False)
    )

    assert target_run.font.bold is True
    assert target_run.font.italic is True
    assert target_run.font.underline is True


def test_copy_run_formatting_pptx2docx_copies_typeface(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that orchestrator copy_run_formatting_pptx2docx properly copies typeface."""

    source_run = pptx_formatting_runs_dict["typeface"]
    target_run = _create_target_docx_run(blank_docx)

    formatting.copy_run_formatting_pptx2docx(
        source_run, target_run, UserConfig(experimental_formatting_on=False)
    )

    assert target_run.font.name is not None and target_run.font.name == "Georgia"


def test_copy_run_formatting_pptx2docx_copies_size_formatting(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that orchestrator copy_run_formatting_pptx2docx properly copies size."""

    small_size_source_run = pptx_formatting_runs_dict["small_size"]
    small_size_target_run = _create_target_docx_run(blank_docx)
    formatting.copy_run_formatting_pptx2docx(
        small_size_source_run,
        small_size_target_run,
        UserConfig(experimental_formatting_on=False),
    )

    large_size_source_run = pptx_formatting_runs_dict["large_size"]
    large_size_target_run = _create_target_docx_run(blank_docx)
    formatting.copy_run_formatting_pptx2docx(
        large_size_source_run,
        large_size_target_run,
        UserConfig(experimental_formatting_on=False),
    )

    assert small_size_target_run.font.size.pt == 8
    assert large_size_target_run.font.size.pt == 48


def test_copy_run_formatting_pptx2docx_copies_color_formatting(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that orchestrator copy_run_formatting_pptx2docx properly copies color."""

    red_source_run = pptx_formatting_runs_dict["red_run"]
    red_target_run = _create_target_docx_run(blank_docx)
    formatting.copy_run_formatting_pptx2docx(
        red_source_run,
        red_target_run,
        UserConfig(experimental_formatting_on=False),
    )

    blue_source_run = pptx_formatting_runs_dict["blue_run"]
    blue_target_run = _create_target_docx_run(blank_docx)
    formatting.copy_run_formatting_pptx2docx(
        blue_source_run,
        blue_target_run,
        UserConfig(experimental_formatting_on=False),
    )

    assert (
        red_target_run.font.color.rgb is not None
        and red_target_run.font.color.rgb == (255, 0, 0)
    )
    assert (
        blue_target_run.font.color.rgb is not None
        and blue_target_run.font.color.rgb == (68, 114, 196)
    )


# endregion

# region _copy_experimental_formatting_pptx2docx tests


def test_copy_experimental_formatting_pptx2docx_preserves_highlight(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that experimental formatting correctly copies highlight from pptx to docx."""
    source_run = pptx_formatting_runs_dict["highlight"]
    target_run = _create_target_docx_run(blank_docx)

    formatting._copy_experimental_formatting_pptx2docx(source_run, target_run)

    assert target_run.font.highlight_color is not None


def test_copy_experimental_formatting_pptx2docx_preserves_single_strike(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that experimental formatting correctly copies single strike from pptx to docx."""
    source_run = pptx_formatting_runs_dict["single_strike"]
    target_run = _create_target_docx_run(blank_docx)

    formatting._copy_experimental_formatting_pptx2docx(source_run, target_run)

    assert target_run.font.strike is True


def test_copy_experimental_formatting_pptx2docx_preserves_double_strike(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that experimental formatting correctly copies double strike from pptx to docx."""
    source_run = pptx_formatting_runs_dict["double_strike"]
    target_run = _create_target_docx_run(blank_docx)

    formatting._copy_experimental_formatting_pptx2docx(source_run, target_run)

    assert target_run.font.double_strike is True


def test_copy_experimental_formatting_pptx2docx_preserves_subscript(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that experimental formatting correctly copies subscript from pptx to docx."""
    source_run = pptx_formatting_runs_dict["subscript"]
    target_run = _create_target_docx_run(blank_docx)

    formatting._copy_experimental_formatting_pptx2docx(source_run, target_run)

    assert target_run.font.subscript is True


def test_copy_experimental_formatting_pptx2docx_preserves_superscript(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that experimental formatting correctly copies superscript from pptx to docx."""
    source_run = pptx_formatting_runs_dict["superscript"]
    target_run = _create_target_docx_run(blank_docx)

    formatting._copy_experimental_formatting_pptx2docx(source_run, target_run)

    assert target_run.font.superscript is True


def test_copy_experimental_formatting_pptx2docx_preserves_all_caps(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that experimental formatting correctly copies all caps from pptx to docx."""
    source_run = pptx_formatting_runs_dict["all_caps"]
    target_run = _create_target_docx_run(blank_docx)

    formatting._copy_experimental_formatting_pptx2docx(source_run, target_run)

    assert target_run.font.all_caps is True


def test_copy_experimental_formatting_pptx2docx_preserves_small_caps(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that experimental formatting correctly copies small caps from pptx to docx."""
    source_run = pptx_formatting_runs_dict["small_caps"]
    target_run = _create_target_docx_run(blank_docx)

    formatting._copy_experimental_formatting_pptx2docx(source_run, target_run)

    assert target_run.font.small_caps is True


# endregion

# region copy_run_formatting_pptx2docx experimental tests


def test_copy_run_formatting_pptx2docx_copies_experimental_when_enabled(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that orchestrator copy_run_formatting_pptx2docx copies experimental formatting when enabled."""
    # Test with highlight
    hl_source = pptx_formatting_runs_dict["highlight"]
    hl_target = _create_target_docx_run(blank_docx)
    formatting.copy_run_formatting_pptx2docx(
        hl_source, hl_target, UserConfig(experimental_formatting_on=True)
    )
    assert hl_target.font.highlight_color is not None

    # Test with subscript
    sub_source = pptx_formatting_runs_dict["subscript"]
    sub_target = _create_target_docx_run(blank_docx)
    formatting.copy_run_formatting_pptx2docx(
        sub_source, sub_target, UserConfig(experimental_formatting_on=True)
    )
    assert sub_target.font.subscript is True

    # Test with all caps
    caps_source = pptx_formatting_runs_dict["all_caps"]
    caps_target = _create_target_docx_run(blank_docx)
    formatting.copy_run_formatting_pptx2docx(
        caps_source, caps_target, UserConfig(experimental_formatting_on=True)
    )
    assert caps_target.font.all_caps is True


def test_copy_run_formatting_pptx2docx_skips_experimental_when_disabled(
    pptx_formatting_runs_dict: dict[str, Run_pptx],
    blank_docx: document.Document,
) -> None:
    """Test that orchestrator copy_run_formatting_pptx2docx skips experimental formatting when disabled."""
    # Test with highlight
    hl_source = pptx_formatting_runs_dict["highlight"]
    hl_target = _create_target_docx_run(blank_docx)
    formatting.copy_run_formatting_pptx2docx(
        hl_source, hl_target, UserConfig(experimental_formatting_on=False)
    )
    assert hl_target.font.highlight_color is None

    # Test with subscript
    sub_source = pptx_formatting_runs_dict["subscript"]
    sub_target = _create_target_docx_run(blank_docx)
    formatting.copy_run_formatting_pptx2docx(
        sub_source, sub_target, UserConfig(experimental_formatting_on=False)
    )
    assert sub_target.font.subscript is not True


# endregion


# endregion

# region pptx2docx PARA tests


# region copy_paragraph_formatting_pptx2docx tests
def test_copy_paragraph_formatting_pptx2docx_copies_font_name(
    path_to_pptx_w_twenty_empty_slides: Path, path_to_empty_docx: Path
) -> None:
    """Test that copy_paragraph_formatting_pptx2docx copies font name from pptx to docx.

    This test validates that get_effective_font_name_pptx works correctly by testing
    through the orchestrator function copy_paragraph_formatting_pptx2docx.
    """
    pptx_with_formatting: presentation.Presentation = Presentation(
        path_to_pptx_w_twenty_empty_slides
    )
    new_docx = Document(str(path_to_empty_docx))
    source_slide = pptx_with_formatting.slides[0]
    source_paragraphs = source_slide.shapes.placeholders[1].text_frame.paragraphs
    source_pptx_para: Paragraph_pptx = source_paragraphs[0]
    assert (
        source_pptx_para.font.name is None
    ), f"Test cannot proceed; test data source_pptx_para.font.name should be None but is {source_pptx_para.font.name}"

    target_docx_para = new_docx.add_paragraph()

    formatting.copy_paragraph_formatting_pptx2docx(
        source_para=source_pptx_para, target_para=target_docx_para
    )

    assert target_docx_para.style.font.name == "Times New Roman"


def test_copy_paragraph_formatting_pptx2docx_copies_alignment(
    path_to_sample_pptx_with_formatting: Path, path_to_empty_docx: Path
) -> None:
    """Test that copy_paragraph_formatting_pptx2docx copies alignment from pptx to docx."""

    pptx_with_formatting: presentation.Presentation = Presentation(
        path_to_sample_pptx_with_formatting
    )

    new_docx = Document(str(path_to_empty_docx))
    source_slide = pptx_with_formatting.slides[0]
    source_paragraphs = source_slide.shapes.placeholders[1].text_frame.paragraphs
    source_pptx_para: Paragraph_pptx = source_paragraphs[0]
    assert (
        source_pptx_para.alignment == PP_ALIGN.CENTER
    ), f"Test cannot proceed because test data is not in expected form. Source_pptx_para.alignment should be PP_ALIGN.CENTER but is {source_pptx_para.alignment}."
    target_docx_para = new_docx.add_paragraph()

    formatting.copy_paragraph_formatting_pptx2docx(
        source_para=source_pptx_para, target_para=target_docx_para
    )

    assert target_docx_para.alignment == WD_ALIGN_PARAGRAPH.CENTER


# endregion


# region test alignment map
def test_alignment_map_wd2pp_completeness() -> None:
    """Verify all common WD_ALIGN values map to PP_ALIGN values"""
    # Test the common alignment cases
    assert formatting.ALIGNMENT_MAP_WD2PP[WD_ALIGN_PARAGRAPH.LEFT] == PP_ALIGN.LEFT
    assert formatting.ALIGNMENT_MAP_WD2PP[WD_ALIGN_PARAGRAPH.CENTER] == PP_ALIGN.CENTER
    assert formatting.ALIGNMENT_MAP_WD2PP[WD_ALIGN_PARAGRAPH.RIGHT] == PP_ALIGN.RIGHT
    assert (
        formatting.ALIGNMENT_MAP_WD2PP[WD_ALIGN_PARAGRAPH.JUSTIFY] == PP_ALIGN.JUSTIFY
    )

    # Edge case: Multiple JUSTIFY types map to same value
    assert (
        formatting.ALIGNMENT_MAP_WD2PP[WD_ALIGN_PARAGRAPH.JUSTIFY_HI]
        == PP_ALIGN.JUSTIFY
    )
    assert (
        formatting.ALIGNMENT_MAP_WD2PP[WD_ALIGN_PARAGRAPH.JUSTIFY_MED]
        == PP_ALIGN.JUSTIFY
    )


def test_alignment_map_pp2wd_is_inverse() -> None:
    """Verify reverse map is constructed correctly"""
    # Note: This won't be perfect inverse due to many-to-one mappings
    # (multiple JUSTIFY variants map to one PP_ALIGN.JUSTIFY)
    for pp_val, wd_val in formatting.ALIGNMENT_MAP_PP2WD.items():
        # The round-trip should work for PP->WD->PP
        assert formatting.ALIGNMENT_MAP_WD2PP[wd_val] == pp_val


# endregion

# endregion


# region apply_experimental_formatting_from_metadata tests


def test_apply_experimental_formatting_from_metadata_highlight(
    blank_docx: document.Document,
) -> None:
    """Test that highlight formatting is applied from metadata correctly"""
    target_run = _create_target_docx_run(blank_docx)
    format_info = {
        "formatting_type": "highlight",
        "ref_text": "test text",
        "highlight_color_enum": "YELLOW",
    }

    formatting.apply_experimental_formatting_from_metadata(target_run, format_info)

    assert target_run.font.highlight_color == WD_COLOR_INDEX.YELLOW


@pytest.mark.parametrize(
    "formatting_type,expected_attr",
    [
        ("strike", "strike"),
        ("double_strike", "double_strike"),
        ("subscript", "subscript"),
        ("superscript", "superscript"),
        ("all_caps", "all_caps"),
        ("small_caps", "small_caps"),
    ],
)
def test_apply_experimental_formatting_from_metadata_bool_formatting(
    blank_docx: document.Document,
    formatting_type: str,
    expected_attr: str,
) -> None:
    """Test that boolean formatting types (strike, caps, sub/superscript) are applied from metadata correctly"""
    target_run = _create_target_docx_run(blank_docx)

    # Verify initial state is False (or None/not set)
    initial_value = getattr(target_run.font, expected_attr)
    assert (
        initial_value is not True
    ), f"Expected {expected_attr} to start as False/None, but was {initial_value}"

    format_info = {
        "formatting_type": formatting_type,
        "ref_text": "test text",
    }

    formatting.apply_experimental_formatting_from_metadata(target_run, format_info)

    actual_value = getattr(target_run.font, expected_attr)
    assert actual_value is True


def test_apply_experimental_formatting_from_metadata_unknown_formatting_type(
    blank_docx: document.Document,
) -> None:
    """Test that unknown formatting_type in the JSON/dict is handled gracefully"""
    target_run = _create_target_docx_run(blank_docx)
    format_info = {
        "formatting_type": "unknown_type",
        "ref_text": "test",
    }

    # Should not crash - just do nothing
    formatting.apply_experimental_formatting_from_metadata(target_run, format_info)

    # Run should be unchanged (no attributes set to True)
    assert target_run.font.strike is not True
    assert target_run.font.all_caps is not True


def test_apply_experimental_formatting_from_metadata_invalid_highlight_enum(
    blank_docx: document.Document, caplog: pytest.LogCaptureFixture
) -> None:
    """Test that invalid highlight_color_enum logs warning"""
    target_run = _create_target_docx_run(blank_docx)
    format_info = {
        "formatting_type": "highlight",
        "ref_text": "test",
        "highlight_color_enum": "INVALID_COLOR",
    }
    # Act
    formatting.apply_experimental_formatting_from_metadata(target_run, format_info)

    # Assert - should log a warning
    with caplog.at_level(logging.DEBUG):
        assert "Could not restore highlight color. Invalid enum" in caplog.text
    # Or just verify highlight_color is None
    assert target_run.font.highlight_color is None


# endregion


# region edge case tests


def test_copy_run_formatting_with_empty_text_docx2pptx(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that copying formatting works when run has empty text."""
    # Create a run with empty text but formatting
    source_para = blank_docx.add_paragraph()
    source_run = source_para.add_run("")
    source_run.font.bold = True
    source_run.font.italic = True

    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)

    formatting.copy_run_formatting_docx2pptx(
        source_run, target_run, [], UserConfig(experimental_formatting_on=False)
    )

    # Should copy formatting even with empty text
    assert target_run.font.bold is True
    assert target_run.font.italic is True
    assert target_run.text == ""


def test_copy_run_formatting_with_whitespace_only_text_docx2pptx(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that copying formatting works when run has only whitespace."""
    # Create a run with whitespace
    source_para = blank_docx.add_paragraph()
    source_run = source_para.add_run("   ")
    source_run.font.bold = True
    source_run.font.highlight_color = 7  # Yellow

    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)
    metadata = []

    formatting.copy_run_formatting_docx2pptx(
        source_run, target_run, metadata, UserConfig(experimental_formatting_on=True)
    )

    # Should copy basic formatting
    assert target_run.font.bold is True
    # Should NOT apply experimental formatting (whitespace-only text is skipped)
    assert len(metadata) == 0


def test_copy_run_formatting_with_none_font_properties_docx2pptx(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that copying formatting handles None font properties gracefully."""
    # Create a run with no explicit formatting (all None)
    source_para = blank_docx.add_paragraph()
    source_run = source_para.add_run("Plain text")

    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)

    # This should not crash
    formatting.copy_run_formatting_docx2pptx(
        source_run, target_run, [], UserConfig(experimental_formatting_on=False)
    )

    # Text should be copied
    assert target_run.text == "Plain text"
    # Font properties remain as they were (or None/inherited)


def test_copy_run_formatting_with_long_text_docx2pptx(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that copying formatting works with very long text."""
    long_text = "A" * 10000  # 10k characters

    source_para = blank_docx.add_paragraph()
    source_run = source_para.add_run(long_text)
    source_run.font.bold = True

    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)

    formatting.copy_run_formatting_docx2pptx(
        source_run, target_run, [], UserConfig(experimental_formatting_on=False)
    )

    assert target_run.font.bold is True
    assert target_run.text == long_text
    assert len(target_run.text) == 10000


def test_copy_run_formatting_with_unicode_text_docx2pptx(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that copying formatting works with Unicode and special characters."""
    unicode_text = "Hello 世界 🌍 café naïve résumé"

    source_para = blank_docx.add_paragraph()
    source_run = source_para.add_run(unicode_text)
    source_run.font.italic = True

    target_run = _create_target_pptx_run(pptx_w_twenty_empty_slides)

    formatting.copy_run_formatting_docx2pptx(
        source_run, target_run, [], UserConfig(experimental_formatting_on=False)
    )

    assert target_run.font.italic is True
    assert target_run.text == unicode_text


def test_copy_run_formatting_with_empty_text_pptx2docx(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that copying formatting works when pptx run has empty text."""
    # Create a pptx run with empty text but formatting
    slide = pptx_w_twenty_empty_slides.slides[0]
    text_frame = slide.shapes.placeholders[1].text_frame
    source_para = text_frame.paragraphs[0]
    source_run = source_para.add_run()
    source_run.text = ""
    source_run.font.bold = True
    source_run.font.italic = True

    target_run = _create_target_docx_run(blank_docx)

    formatting.copy_run_formatting_pptx2docx(
        source_run, target_run, UserConfig(experimental_formatting_on=False)
    )

    # Should copy formatting even with empty text
    assert target_run.font.bold is True
    assert target_run.font.italic is True
    assert target_run.text == ""


def test_copy_run_formatting_with_unicode_text_pptx2docx(
    blank_docx: document.Document,
    pptx_w_twenty_empty_slides: presentation.Presentation,
) -> None:
    """Test that copying formatting works with Unicode from pptx to docx."""
    unicode_text = "Привет 你好 مرحبا 🎉"

    slide = pptx_w_twenty_empty_slides.slides[0]
    text_frame = slide.shapes.placeholders[1].text_frame
    source_para = text_frame.paragraphs[0]
    source_run = source_para.add_run()
    source_run.text = unicode_text
    source_run.font.bold = True

    target_run = _create_target_docx_run(blank_docx)

    formatting.copy_run_formatting_pptx2docx(
        source_run, target_run, UserConfig(experimental_formatting_on=False)
    )

    assert target_run.font.bold is True
    assert target_run.text == unicode_text


# endregion


# region test color map
def test_color_map_hex_has_expected_colors() -> None:
    """Verify color map contains expected highlight colors"""
    # Test a few key colors
    assert formatting.COLOR_MAP_HEX[WD_COLOR_INDEX.YELLOW] == "FFFF00"
    assert formatting.COLOR_MAP_HEX[WD_COLOR_INDEX.PINK] == "FF00FF"
    assert formatting.COLOR_MAP_HEX[WD_COLOR_INDEX.BLACK] == "000000"


def test_color_map_from_hex_is_inverse() -> None:
    """Verify reverse color map is constructed correctly"""
    for wd_color, hex_val in formatting.COLOR_MAP_HEX.items():
        assert formatting.COLOR_MAP_FROM_HEX[hex_val] == wd_color


def test_color_map_completeness() -> None:
    """Verify all WD_COLOR_INDEX highlight colors are mapped (except AUTO)"""
    # List all possible highlight colors from WD_COLOR_INDEX enum
    # We exclude AUTO (0) since it represents "automatic/default" rather than a specific color
    expected_colors = [
        WD_COLOR_INDEX.BLACK,
        WD_COLOR_INDEX.BLUE,
        WD_COLOR_INDEX.TURQUOISE,
        WD_COLOR_INDEX.BRIGHT_GREEN,
        WD_COLOR_INDEX.PINK,
        WD_COLOR_INDEX.RED,
        WD_COLOR_INDEX.YELLOW,
        WD_COLOR_INDEX.WHITE,
        WD_COLOR_INDEX.DARK_BLUE,
        WD_COLOR_INDEX.TEAL,
        WD_COLOR_INDEX.GREEN,
        WD_COLOR_INDEX.VIOLET,
        WD_COLOR_INDEX.DARK_RED,
        WD_COLOR_INDEX.DARK_YELLOW,
        WD_COLOR_INDEX.GRAY_50,
        WD_COLOR_INDEX.GRAY_25,
    ]
    for color in expected_colors:
        assert color in formatting.COLOR_MAP_HEX, f"Missing color mapping for {color}"


# endregion
