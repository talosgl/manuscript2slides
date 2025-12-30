"""Shared test helper functions."""

# mypy: disable-error-code="import-untyped"
# pyright: reportArgumentType=false

import xml.etree.ElementTree as ET

from docx import document
from docx.text.hyperlink import Hyperlink as Hyperlink_docx
from docx.text.paragraph import Paragraph as Paragraph_docx
from docx.text.run import Run as Run_docx
from pptx import presentation
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.text.text import _Paragraph as Paragraph_pptx
from pptx.text.text import _Run as Run_pptx

from manuscript2slides.processing.populate_docx import get_slide_paragraphs


def find_first_docx_para_containing(
    doc: document.Document, search_text: str
) -> Paragraph_docx:
    """Find the first paragraph instance with the search text or raises an AssertionError if it can't be found."""
    if not doc.paragraphs:
        raise AssertionError(
            f"Test cannot proceed because there aren't any paragraphs in this document."
        )

    for para in doc.paragraphs:
        if search_text in para.text:
            return para

    raise AssertionError(
        f"Test cannot proceed because the required text '{search_text}' could not be found in document."
    )


def find_first_docx_run_in_para_exact_match(
    parent: Paragraph_docx | Hyperlink_docx, search_text: str
) -> Run_docx:
    """Find the first run in a docx paragraph or hyperlink with exact matching text, or raise AssertionError if it can't be found."""
    if not parent.runs:
        raise AssertionError(
            f"Test cannot proceed because this paragraph doesn't have any runs."
        )

    for run in parent.runs:
        if run.text == search_text:
            return run

    raise AssertionError(
        f"Test cannot proceed because the required text '{search_text}' could not be found in paragraph."
    )


def find_inner_run_containing(
    hyperlink_parent: Hyperlink_docx, search_text: str
) -> Run_docx:
    """Return the run inside this hyperlink containing the search text, or raise AssertionError if it can't be found."""
    if not hyperlink_parent.runs:
        raise AssertionError(
            f"Test cannot proceed because this hyperlink doesn't have any runs."
        )

    for run in hyperlink_parent.runs:
        if search_text in run.text:
            return run

    raise AssertionError(
        f"Test cannot proceed because the required text '{search_text}' could not be found in hyperlink parent."
    )


def find_first_docx_item_in_para_containing(
    para_parent: Paragraph_docx, search_text: str
) -> Run_docx | Hyperlink_docx:
    """Find the first run in a docx paragraph or hyperlink containing the search text,
    or raise AssertionError if it can't be found."""
    if not para_parent.runs and not para_parent.hyperlinks:
        raise AssertionError(
            f"Test cannot proceed because this paragraph doesn't have any inner contents (hyperlinks nor runs)."
        )

    for item in para_parent.iter_inner_content():
        if search_text in item.text:
            return item

    raise AssertionError(
        f"Test cannot proceed because the required text '{search_text}' could not be found in paragraph."
    )


def find_first_slide_containing(
    prs: presentation.Presentation, search_text: str
) -> tuple[Slide, Paragraph_pptx]:
    """Find the first slide that contains the given text anywhere in its shapes (but not searching its speaker notes),
    or raise an assertion error if it is not found.
    """

    if not prs.slides:
        raise AssertionError(
            f"Test cannot proceed because there's no slides in this presentation."
        )

    slides = list(prs.slides)

    for slide in slides:
        paragraphs: list[Paragraph_pptx] = get_slide_paragraphs(slide)

        for para in paragraphs:
            if search_text in para.text:
                return (slide, para)

    raise AssertionError(
        f"Test cannot proceed because the required text '{search_text}' could not be found."
    )


def find_first_pptx_run_in_para_containing(
    para: Paragraph_pptx, search_text: str
) -> Run_pptx:
    """Given a paragraph from a slide, find the run containing the search text, or raise an
    AssertionError if the run cannot be found.
    """
    if not para.text or (len(para.runs) < 1):
        raise AssertionError(
            f"Test cannot proceed because this paragraph has no text or runs."
        )

    for run in para.runs:
        if search_text in run.text:
            return run

    raise AssertionError(
        f"Test cannot proceed because the required text '{search_text}' could not be found in paragraph."
    )


def find_first_pptx_run_in_para_with_exact_match(
    para: Paragraph_pptx, exact_text: str
) -> Run_pptx:
    """Find a run whose text exactly matches the given text (not substring match), or raises an
    AssertionError if the run cannot be found.
    """
    if not para.text or len(para.runs) < 1:
        raise AssertionError(
            f"Test cannot proceed because this paragraph has no text or runs."
        )

    for run in para.runs:
        if run.text == exact_text:
            return run

    raise AssertionError(
        f"Test cannot proceed because the required text '{exact_text}' could not be found in paragraph."
    )


def get_slide_notes_text(slide: Slide) -> TextFrame:
    """Helper to modularize getting a slide's notes frame for testing against."""
    assert (
        slide.has_notes_slide is True
        and slide.notes_slide is not None
        and slide.notes_slide.notes_text_frame is not None
    ), (
        f"Test cannot proceed because we expected the slide, {slide.slide_id} to have a speaker notes section and it does not."
    )

    notes_text_frame: TextFrame = slide.notes_slide.notes_text_frame

    return notes_text_frame


def run_has_highlight(run: Run_pptx) -> bool:
    """Check if a pptx run has highlight formatting by inspecting its XML.

    Returns:
        True if the run has highlight formatting, False otherwise.
    """
    if not hasattr(run, "_r"):
        return False

    try:
        xml = run._r.xml
        return "a:highlight" in xml
    except (AttributeError, TypeError):
        return False


def get_run_highlight_color(run: Run_pptx) -> str | None:
    """Extract the highlight color from a pptx run's XML and return it as Hex color string (e.g., "FFFF00" for yellow),
    or return None if no highlight or color not found.
    """
    if not hasattr(run, "_r"):
        return None

    try:
        xml = run._r.xml
        if "a:highlight" not in xml:
            return None

        root = ET.fromstring(xml)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        highlight = root.find(".//a:highlight/a:srgbClr", ns)

        if highlight is not None:
            return highlight.get("val")

        return None
    except (AttributeError, TypeError, ET.ParseError):
        return None
