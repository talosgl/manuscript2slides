"""Tests to ensure we can create blank slide deck and blank document from our standard templates."""

# pyright: reportArgumentType=false

import pytest
from manuscript2slides import templates
from pathlib import Path
import pptx
import docx
from manuscript2slides.internals.define_config import UserConfig


# region test_create_empty_slide_deck
def test_create_empty_slide_deck_returns_empty_deck(
    path_to_sample_pptx_with_formatting: Path, path_to_sample_docx_with_formatting: Path
) -> None:
    """Ensure that the base presentation object is empty of slides,
    even if a pptx is passed in that template which contains slides."""
    # Arrange: Construct a cfg using a deck as template that we know has slides in it
    test_cfg = UserConfig(
        input_docx=str(path_to_sample_docx_with_formatting),
        template_pptx=str(path_to_sample_pptx_with_formatting),
    )

    # Action: Pass that to our function
    prs = templates.create_empty_slide_deck(test_cfg)

    # Assert the slide count is now 0
    assert len(prs.slides) == 0, f"prs.slides has {len(prs.slides)} slides in it."


def test_create_empty_slide_deck_fails_gracefully_without_slide_layout():
    # TODO/Strategy: I think we need to make a mock object with improper slide layout
    # Then we need to use pytest.raises to catch the error
    # and we need to make sure caplog contains helpful text.

    # with pytest.raises(ValueError):
    pass


# endregion


# region test_empty_document


def test_create_empty_document_returns_empty_doc(
    path_to_sample_docx_with_formatting: Path, path_to_sample_pptx_with_formatting: Path
) -> None:
    """Ensure that even if we're passed in a docx full of stuff as a template, we start with no pages/paragraphs."""
    # Arrange: Construct a cfg using a doc that we know has content in it
    test_cfg = UserConfig(
        input_pptx=str(path_to_sample_pptx_with_formatting),
        template_docx=str(path_to_sample_docx_with_formatting),
    )

    # Action: Pass that to our function
    docu = templates.create_empty_document(test_cfg)

    assert (
        len(docu.paragraphs) == 0
    ), f"docu.paragraphs has {len(docu.paragraphs)} paragraphs in it."


def test_create_empty_document_fails_gracefully_on_missing_styles():
    pass


# endregion
