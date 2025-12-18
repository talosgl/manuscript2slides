"""
Extract baseline JSON data from a pptx file for integration testing.

Usage:
    python tests/tools/extract_pptx_data.py

Extracts slide structure, formatting, and experimental XML properties.
"""

from pathlib import Path
import json
import sys
from typing import cast
import xml.etree.ElementTree as ET

from pptx.presentation import Presentation
from pptx.slide import Slides
from manuscript2slides import io
from manuscript2slides.annotations.restore_from_slides import split_speaker_notes

from extraction_utils import safe_pprint, filter_none_keep_false, rgb_to_hex


# ============================================================================
# CONFIGURATION - Edit these to change input/output files
# ============================================================================
INPUT_PPTX = "tests/data/sample_slides_output.pptx"
OUTPUT_JSON = "tests/baselines/pptx_sample.json"
# Alternative configs:
# INPUT_PPTX = "tests/data/test_formatting_expected_output.pptx"
# OUTPUT_JSON = "tests/baselines/pptx_formatting.json"
# ============================================================================


def main() -> None:
    """Load a test pptx and explore its structure."""

    # Configuration
    test_pptx_path = Path(INPUT_PPTX)

    if not test_pptx_path.exists():
        print(f"Error: {test_pptx_path} not found!", file=sys.stderr)
        print(
            f"Please check the INPUT_PPTX configuration at the top of this script.",
            file=sys.stderr,
        )
        sys.exit(1)

    # Load the pptx
    print(f"Loading {test_pptx_path}...")
    prs = io.load_and_validate_pptx(test_pptx_path)
    slides = cast(Slides, prs.slides)

    print(f"\nFound {len(slides)} slides")
    print("\n=== Extracting fixture data ===")

    data = []
    for idx, slide in enumerate(slides):
        slide_data = {
            "slide_number": idx + 1,
            "shapes": [],
        }

        # Extract speaker notes if present
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            speaker_notes_text = notes_text_frame.text

            # Use the restore function to parse the notes structure
            parsed_notes = split_speaker_notes(speaker_notes_text)

            # Only add speaker_notes if there's actual content
            if parsed_notes.user_notes or parsed_notes.metadata:
                slide_data["speaker_notes"] = {
                    "user_notes": parsed_notes.user_notes,
                    "metadata": parsed_notes.metadata,
                }

                # Also capture the individual annotation lists for easier testing
                if parsed_notes.comments:
                    slide_data["speaker_notes"]["comments"] = parsed_notes.comments
                if parsed_notes.footnotes:
                    slide_data["speaker_notes"]["footnotes"] = parsed_notes.footnotes
                if parsed_notes.endnotes:
                    slide_data["speaker_notes"]["endnotes"] = parsed_notes.endnotes
                if parsed_notes.headings:
                    slide_data["speaker_notes"]["headings"] = parsed_notes.headings
                if parsed_notes.experimental_formatting:
                    slide_data["speaker_notes"][
                        "experimental_formatting"
                    ] = parsed_notes.experimental_formatting

        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_data = {"text": shape.text, "paragraphs": []}

                for para in shape.text_frame.paragraphs:
                    para_font = {
                        "bold": para.font.bold,
                        "italic": para.font.italic,
                        "underline": para.font.underline,
                        "font_name": para.font.name,
                        "font_size": para.font.size.pt if para.font.size else None,
                    }
                    # Filter out None values but keep False
                    para_font = filter_none_keep_false(para_font)

                    para_data = {
                        "text": para.text,
                        "level": para.level,
                        "paragraph_font": para_font,
                        "runs": [],
                    }

                    for run_idx, run in enumerate(para.runs):
                        run_data = {
                            "run_number": run_idx,
                            "text": run.text,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "underline": run.font.underline,
                            "font_name": run.font.name,
                            "font_size": run.font.size.pt if run.font.size else None,
                        }

                        # Extract color if it's RGB
                        if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                            run_data["color"] = rgb_to_hex(run.font.color.rgb)

                        # Extract experimental formatting from XML (not exposed by python-pptx)
                        if hasattr(run, "_r") and hasattr(run.font, "_element"):
                            xml = run._r.xml
                            element = run.font._element

                            # Highlight - extract both presence and color
                            if "a:highlight" in xml:
                                run_data["has_highlight"] = True
                                # Parse XML to extract highlight color
                                try:
                                    root = ET.fromstring(xml)
                                    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                                    highlight = root.find(".//a:highlight/a:srgbClr", ns)
                                    if highlight is not None:
                                        hex_color = highlight.get("val")
                                        if hex_color:
                                            run_data["highlight_color"] = hex_color
                                except ET.ParseError:
                                    pass  # Just keep has_highlight=True without color

                            # Strike (single or double)
                            strike_val = element.get("strike")
                            if strike_val:
                                run_data["strike"] = (
                                    strike_val  # "sngStrike" or "dblStrike"
                                )

                            # Caps (all or small)
                            cap_val = element.get("cap")
                            if cap_val:
                                run_data["cap"] = cap_val  # "all" or "small"

                            # Subscript/Superscript (baseline)
                            baseline_val = element.get("baseline")
                            if baseline_val:
                                run_data["baseline"] = int(
                                    baseline_val
                                )  # negative=sub, positive=super

                        # Filter out None values but keep False
                        run_data = filter_none_keep_false(run_data)
                        para_data["runs"].append(run_data)

                    shape_data["paragraphs"].append(para_data)

                slide_data["shapes"].append(shape_data)

        data.append(slide_data)

    # Show what we got
    safe_pprint(data, width=120, item_type="slides")

    # Save to JSON
    output_path = Path(OUTPUT_JSON)
    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
        print(f"\nSaved to {output_path}")
    except (OSError, PermissionError) as e:
        print(f"Error: Failed to write output file: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"Fatal error: {e}", file=sys.stderr)
        sys.exit(1)
