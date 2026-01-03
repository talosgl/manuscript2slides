"""Test I/O functions"""

# tests/test_io.py
import logging
import sys
from pathlib import Path

import pytest
from docx import Document, document
from pptx import Presentation, presentation

from manuscript2slides.io import (
    _validate_path,
    load_and_validate_docx,
    load_and_validate_pptx,
    validate_docx_path,
    validate_pptx_path,
)


# region test _validate_path
def test_validate_path_raises_when_path_is_dir(
    tmp_path: Path, caplog: pytest.LogCaptureFixture
) -> None:
    """Ensure we raise when a folder is passed in instead of a file."""

    with caplog.at_level(logging.ERROR):
        with pytest.raises(ValueError, match="file"):
            _validate_path(tmp_path)
    # Assert: test error is captured in the log at error level
    assert "Path is not a file" in caplog.text


@pytest.mark.skipif(
    sys.platform == "win32",
    reason="Windows-style path detection only applies on non-Windows systems",
)
def test_validate_path_detects_windows_paths_on_unix(
    caplog: pytest.LogCaptureFixture,
) -> None:
    """Ensure we provide helpful error when user provides Windows-style path on Unix/macOS."""
    windows_style_path = r"C:\Users\test\file.docx"

    with caplog.at_level(logging.ERROR):
        with pytest.raises(FileNotFoundError, match="File not found"):
            _validate_path(windows_style_path)

    # Assert: check that we logged the helpful Windows-style path message
    assert "Windows-style path" in caplog.text
    assert "backslashes" in caplog.text
    assert "forward slashes" in caplog.text


# endregion


# region test_validate_docx_path
def test_validate_docx_path_accepts_valid_file(
    path_to_sample_docx_with_formatting: Path,
) -> None:
    """Test that validate_docx_path accepts valid .docx files"""
    result = validate_docx_path(path_to_sample_docx_with_formatting)

    assert result == path_to_sample_docx_with_formatting
    assert result.exists()


def test_validate_docx_path_rejects_wrong_extension(
    tmp_path: Path, caplog: pytest.LogCaptureFixture
) -> None:
    """Test that non-.docx files are rejected"""
    wrong_file = tmp_path / "not_a_docx.txt"
    wrong_file.touch()

    with pytest.raises(ValueError, match="Expected a .docx"):
        validate_docx_path(wrong_file)

    assert "Wrong file extension" in caplog.text


def test_validate_docx_path_rejects_old_doc_format(
    tmp_path: Path, caplog: pytest.LogCaptureFixture
) -> None:
    """Test that old .doc format is rejected with helpful message"""
    old_doc = tmp_path / "old.doc"
    old_doc.touch()

    with pytest.raises(ValueError, match="only supports .docx files. Please convert"):
        validate_docx_path(old_doc)

    assert "Unsupported .doc file" in caplog.text


# endregion


# region test_validate_pptx_path
def test_validate_pptx_path_accepts_valid_file(
    path_to_sample_pptx_with_formatting: Path,
) -> None:
    """Test that a valid file path passed in passes validation and is passed onward to caller."""
    result = validate_pptx_path(path_to_sample_pptx_with_formatting)
    assert result == path_to_sample_pptx_with_formatting


def test_validate_pptx_path_rejects_wrong_extension(
    tmp_path: Path, caplog: pytest.LogCaptureFixture
) -> None:
    """Verify we reject bad extensions and provide helpful error and logging messages."""
    wrong_ext = tmp_path / "my_slides.txt"
    wrong_ext.touch()  # make this a readable path

    with pytest.raises(ValueError, match="Expected a .pptx"):
        validate_pptx_path(wrong_ext)

    assert "expected .pptx, got" in caplog.text


def test_validate_pptx_path_rejects_old_ppt_format(
    tmp_path: Path, caplog: pytest.LogCaptureFixture
) -> None:
    """Verify we catch old file type .ppt as a special case and provide helpful error and logging messaging."""
    wrong_ext = tmp_path / "my_slides.ppt"
    wrong_ext.touch()  # make this a readable path

    with pytest.raises(ValueError, match="only supports .pptx files. Please convert"):
        validate_pptx_path(wrong_ext)

    assert "Unsupported .ppt file" in caplog.text


# endregion


# region test_load_and_validate_docx


def test_load_docx_returns_non_empty_document_object(
    path_to_sample_docx_with_formatting: Path,
) -> None:
    """Test that loading a valid docx returns a Document object that is non-empty"""
    doc = load_and_validate_docx(path_to_sample_docx_with_formatting)

    assert isinstance(doc, document.Document)
    assert len(doc.paragraphs) > 0  # Should have some content


def test_load_docx_rejects_missing_file() -> None:
    """Test that a missing file passed in raises FileNotFoundError and it's bubbled up through load_and_validate_docx"""
    fake_path = Path("i_dont_exist.docx")

    with pytest.raises(FileNotFoundError, match="File not found"):
        load_and_validate_docx(fake_path)


def test_load_docx_raises_on_perms_error(
    monkeypatch: pytest.MonkeyPatch, caplog: pytest.LogCaptureFixture
) -> None:
    """Test that when PermissionError is raised from downstream, it's raised in the function,
    provides a helpful message, and logs to the logger."""

    # Arrange:
    # Make a nested/scoped function that will raise a permission error.
    def mock_document_raise(*args, **kwargs) -> None:  # noqa: ANN003, ANN002
        raise PermissionError("No joy buddy")

    # Disguise our function as the docx.Document() constructor call.
    # Now when docx.Document() is called by the function we're calling/testing,
    # actually our fake function will be called instead.
    monkeypatch.setattr("docx.Document", mock_document_raise)

    # Act/Assert:
    with pytest.raises(PermissionError, match="open in another program"):
        # It doesn't matter what path we send, it's discarded.
        load_and_validate_docx(Path("fake_path.docx"))  # Test that raise happens

    # Assert: Test error is captured in log
    assert "Permission denied" in caplog.text


def test_load_docx_rejects_empty_docx(tmp_path: Path) -> None:
    """Test that document with no content raises ValueError"""
    empty = tmp_path / "empty.docx"
    doc = Document()
    # Don't add any paragraphs
    doc.save(str(empty))

    with pytest.raises(ValueError, match="no paragraphs"):
        load_and_validate_docx(empty)

    doc.add_paragraph()
    doc.save(str(empty))

    with pytest.raises(ValueError, match="no text content"):
        load_and_validate_docx(empty)


# endregion


# region test_load_and_validate_pptx
def test_load_pptx_returns_non_empty_document_object(
    path_to_sample_pptx_with_formatting: Path,
) -> None:
    """Test that a valid pptx that is non-empty returns a pptx object"""
    prs = load_and_validate_pptx(path_to_sample_pptx_with_formatting)

    assert isinstance(prs, presentation.Presentation)
    assert prs.slides  # Verify presentation has at least one slide


def test_load_pptx_rejects_missing_file(caplog: pytest.LogCaptureFixture) -> None:
    """Test file not found raises FileNotFound error and adds helpful message to log."""
    fake_path = Path("i_dont_exist.pptx")

    with pytest.raises(FileNotFoundError, match="File not found"):
        load_and_validate_pptx(fake_path)

    assert "not found" in caplog.text


def test_load_pptx_raises_on_perms_error(
    monkeypatch: pytest.MonkeyPatch, caplog: pytest.LogCaptureFixture
) -> None:
    """Test that PermissionsError is properly raised, bubbled up, and logged for unreachable files during pptx flow."""

    def mock_presentation_raise(*args, **kwargs) -> None:  # noqa: ANN003, ANN002
        raise PermissionError("ROAD CLOSED")

    # Disguise our above mock function as the pptx.Presentation() constructor
    monkeypatch.setattr("pptx.Presentation", mock_presentation_raise)

    with pytest.raises(PermissionError, match="open in another program"):
        # It doesn't matter what path we send, it's discarded.
        load_and_validate_pptx(Path("fake_path.pptx"))

    # Assert: Test error is captured in log
    assert "Permission denied" in caplog.text


def test_load_pptx_rejects_pptx_with_no_slides_or_no_text(
    tmp_path: Path, path_to_empty_pptx: Path
) -> None:
    """Ensure we reject empty slide decks (no slides OR slides, but no text) in pptx2docx flow."""
    no_slides = tmp_path / "no_slides.pptx"
    # We can't directly create a presentation without a template, so we provide one that's empty.
    prs: presentation.Presentation = Presentation(path_to_empty_pptx)
    # Don't add any slides
    prs.save(str(no_slides))

    with pytest.raises(ValueError, match="no slides"):
        load_and_validate_pptx(no_slides)

    slide_layout = prs.slide_layouts[0]

    prs.slides.add_slide(slide_layout)  # type: ignore
    prs.save(str(no_slides))

    with pytest.raises(ValueError, match="text content"):
        load_and_validate_pptx(no_slides)


# endregion
