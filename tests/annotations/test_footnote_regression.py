"""Regression test for the '1. 1' footnote double numbering issue.

This test verifies that when footnotes are copied to slide notes, they don't
get double-numbered. For example, the first footnote should render as:
  "1. James Griffiths..."
not:
  "1. 1 James Griffiths..."

The issue was caused by the regex in docx_xml.py not properly stripping the
leading number from footnote text that came from Word without a period after
the number (e.g., "1 Text" instead of "1. Text").
"""
import pytest
from pathlib import Path
from manuscript2slides.orchestrator import run_pipeline
from manuscript2slides.internals.define_config import UserConfig
from pptx import Presentation


def test_footnote_no_double_numbering(
    path_to_sample_docx_with_everything: Path,
    path_to_empty_pptx: Path,
    session_temp_dir: Path,
) -> None:
    """Verify footnotes don't have double numbering like '1. 1 Text'.

    Regression test for issue where footnotes would render as "1. 1 James Griffiths"
    instead of "1. James Griffiths" in slide notes.
    """
    # Run the pipeline with footnotes enabled
    cfg = UserConfig(
        input_docx=path_to_sample_docx_with_everything,
        template_pptx=path_to_empty_pptx,
        output_folder=session_temp_dir,
    )
    cfg.display_footnotes = True

    output_path = run_pipeline(cfg)
    prs = Presentation(output_path)

    # Find a slide with the "James Griffiths" footnote
    found_footnote_slide = False
    for slide in prs.slides:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text
            if 'FOOTNOTES' in notes_text and 'James Griffiths' in notes_text:
                found_footnote_slide = True

                # Check that there's no double numbering pattern
                # Should be "1. James Griffiths" NOT "1. 1 James Griffiths"
                assert '1. 1 James' not in notes_text, \
                    "Found double numbering '1. 1 James' in footnote text"
                assert '1.1 James' not in notes_text, \
                    "Found double numbering '1.1 James' in footnote text"

                # Verify the footnote is actually there with correct numbering
                assert '1. James Griffiths' in notes_text, \
                    "Expected footnote '1. James Griffiths' not found in slide notes"
                break

    assert found_footnote_slide, \
        "Could not find a slide with the 'James Griffiths' footnote to test against"
