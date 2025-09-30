# pyright: basic
"""
Convert Microsoft Word documents to PowerPoint presentations.

This tool processes .docx files and converts them into .pptx slide decks by chunking
the document content based on various strategies (paragraphs, headings, or page breaks).
Text formatting like bold, italics, and colors can optionally be preserved.

The main workflow:
1. Load a .docx file using python-docx
2. Chunk the content based on the selected strategy
3. Create slides from chunks using a PowerPoint template
4. Save the resulting .pptx file

Supported chunking methods:
- paragraph: Each paragraph becomes a slide
- page: New slide for each page break
- heading_flat: New slide for each heading (any level)
- heading_nested: New slide based on heading hierarchy

Example:
    python docx2pptx-text.py

    (Configure INPUT_DOCX_FILE and other constants before running)
"""
# region imports
from __future__ import annotations

# Standard library
import io
import platform
import re
import sys
from datetime import datetime
from enum import Enum
from pathlib import Path
from dataclasses import dataclass, field
from typing import TypeVar
import json

# import warnings

# Third-party libraries
import docx
import pptx
from docx import document
from docx.comments import Comment as Comment_docx
from docx.opc.part import Part
from docx.opc import constants
from docx.oxml.ns import qn
from docx.oxml.parser import OxmlElement as OxmlElement_docx

# from docx.table import Table as Table_docx
# from docx.text.hyperlink import Hyperlink as Hyperlink_docx
from docx.text.paragraph import Paragraph as Paragraph_docx
from docx.text.run import Run as Run_docx
from pptx import presentation
from pptx.dml.color import RGBColor as RGBColor_pptx
from pptx.slide import Slide, SlideLayout
from pptx.text.text import TextFrame, _Paragraph as Paragraph_pptx, _Run as Run_pptx  # type: ignore
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.oxml.xmlchemy import OxmlElement as OxmlElement_pptx
from pptx.util import Pt
import xml.etree.ElementTree as ET

from pptx.slide import NotesSlide

# RUN_TYPE = TypeVar("RUN_TYPE", Run_docx, Run_pptx)
from docx.shared import RGBColor as RGBColor_docx
from docx.text.font import Font as Font_docx
from pptx.text.text import Font as Font_pptx
from typing import Union

# endregion

# region colormap
from docx.enum.text import WD_COLOR_INDEX

COLOR_MAP_HEX = {
    WD_COLOR_INDEX.YELLOW: "FFFF00",
    WD_COLOR_INDEX.PINK: "FF00FF",
    WD_COLOR_INDEX.BLACK: "000000",
    WD_COLOR_INDEX.WHITE: "FFFFFF",
    WD_COLOR_INDEX.BLUE: "0000FF",
    WD_COLOR_INDEX.BRIGHT_GREEN: "00FF00",
    WD_COLOR_INDEX.DARK_BLUE: "000080",
    WD_COLOR_INDEX.DARK_RED: "800000",
    WD_COLOR_INDEX.DARK_YELLOW: "808000",
    WD_COLOR_INDEX.GRAY_25: "C0C0C0",
    WD_COLOR_INDEX.GRAY_50: "808080",
    WD_COLOR_INDEX.GREEN: "008000",
    WD_COLOR_INDEX.RED: "FF0000",
    WD_COLOR_INDEX.TEAL: "008080",
    WD_COLOR_INDEX.TURQUOISE: "00FFFF",
    WD_COLOR_INDEX.VIOLET: "800080",
}

COLOR_MAP_FROM_HEX = {v: k for k, v in COLOR_MAP_HEX.items()}
BASELINE_SUBSCRIPT_SMALL_FONT = "-50000"
BASELINE_SUBSCRIPT_LARGE_FONT = "-25000"
BASELINE_SUPERSCRIPT_SMALL_FONT = "60000"  # For fonts < 24pt
BASELINE_SUPERSCRIPT_LARGE_FONT = "30000"  # For fonts >= 24pt
#

# region Overarching TODOs
"""
Must-Implement v0 Features:
- Change consts configuration to use a class or similar
- Rearchitect to be multi-file

v1 features I'd like:
- Add a feature to split the output pptx or docx into multiple files based on slide or page count. Add default counts and allow user overrides for the default.
- Investigate if we can insert pptx sections safely enough (to allow for docx headings -> pptx sections, or other section-chunking)
    - If not, investigate if/when we want to mimic the same type of behavior with "segue slides"
- Create Documents/docx2pptx/input/output/resources structure
    - Copies sample files from app resources to user folders
    - Cleanup mode for debug runs    
- Add an actual logger
- Investigate how impossible non-local file input/output (OneDrive/SharePoint) would be; add to known limitations if not supportable.

Public v1
- What does is "done enough for public github repo mean"? 
    - "When I'm comfortable having strangers use it without asking me questions."
    - Engineer-audience documentation
    - Log & error messages that tell users what went wrong and how to fix it
    - Code that doesn't crash on common edge cases.

Public v2: UI
    - Build a simple UI with good enough UX that any non-tech-savvy writer can use it without friction

Public v3: Package/Distribution
    - Figure out how to package it into an installer for each non-mobile platform (Win 11, MacOS, Linux)

Stretch Wishlist Features:
-   Investigate linking slides or sections-of-slides or file chunks back to their source "place" in the original docx (og file if possible, or a copy where we insert the anchor)
-   Add support for importing .md and .txt; split by whitespaces or newline characters.
-   Add support to break chunks (of any type) at a word count threshold.

"""

"""
Known Issues & Limitations:
    -   We only support text content. No images, tables, etc., are copied between the formats, and we do not have plans 
        to support these in future.
    
    - "Sections" in both docx and pptx are not supported. TODO: investigate

    -   We do not support .doc or .ppt, only .docx. If you have a .doc file, convert it to .docx using Word, Google Docs, 
        or LibreOffice before processing.

    -   We do not support .ppt, only .pptx.

    -   Auto-fit text resizing in slides doesn't work. PowerPoint only applies auto-fit sizing when opened in the UI. 
        You can get around this manually with these steps:
            1. Open up the output presentation in PowerPoint Desktop > View > Slide Master
            2. Select the text frame object, right-click > Format Shape
            3. Click the Size & Properties icon {TODO ADD IMAGES}
            4. Click Text Box to see the options
            5. Toggle "Do not Autofit" and then back to "Shrink Text on overflow"
            6. Close Master View
            7. Now all the slides should have their text properly resized.

    -   Field code hyperlinks not supported - Some hyperlinks (like the sample_doc.docx's first "Where are Data?" link) 
        are stored using Word's field code format and display as plain text instead of clickable links. The exact 
        conditions that cause this format are unclear, but it may occur with hyperlinks in headings or certain copy/paste 
        scenarios. We think most normal hyperlinks will work fine. We try to report when we detect these are present but cannot
        reliably copy them as text into the body.

    -   ANNOTATIONS LIMITATIONS
        -   We collapse all comments, footnotes, and endnotes into a slide's speaker notes. PowerPoint itself doesn't 
            support real footnotes or endnotes at all. It does have a comments functionality, but the library used here 
            (python-pptx) doesn't support adding comments to slides yet. 

        -   Note that inline reference numbers (1, 2, 3, etc.) from the docx body are not preserved in the slide text - 
            only the annotation content appears in speaker notes.

        -   You can choose to preserve some comment metadata (author, timestamps) in plain text, but not threading.
    
    -   REVERSE FLOW LIMITATIONS
        -   The reverse flow (pptx2docx-text) is significantly less robust. Your original input document to the docx2pptx-text flow, and 
            the output document from a follow-up pptx2docx-text flow will not look the same. Expect to lose images, tables, footnotes, 
            endnotes, and fancy formatting. We attempt to preserve headings (text-matching based). Comments should be restored, but their 
            anchor positioning may be altered slightly.

        -   There will always be a blank line at the start of the reverse-pipeline document. When creating a new document with python-docx 
            using Document(), it inherently includes a single empty paragraph at the start. This behavior is mandated by the Open XML 
            .docx standard, which requires at least one paragraph within the w:body element in the document's XML structure.

"""
# endregion
# TODO: rename to docx2pptx-text / pptx2docx-text

# region CONSTANTS / config.py
# Get the directory where this script lives (NOT INTENDED FOR USER EDITING)
SCRIPT_DIR = Path(__file__).parent


# === docx2pptx Consts for script user to alter per-run ===

# The pptx file to use as the template for the slide deck
TEMPLATE_PPTX = SCRIPT_DIR / "resources" / "blank_template.pptx"
# You can make your own template with the master slide and master notes page
# to determine how the output will look. You can customize things like font, paragraph style,
# slide size, slide layout...

# Desired slide layout. All slides use the same layout.
SLD_LAYOUT_CUSTOM_NAME = "docx2pptx"

# Desired output directory/folder to save the pptx in
OUTPUT_PPTX_FOLDER = SCRIPT_DIR / "output"
# e.g., r"c:\my_presentations"
# If you leave it blank it'll save in the root of where you run the script from the command line

# Desired output filename; Note that this will clobber an existing file of the same name!
OUTPUT_PPTX_FILENAME = r"sample_slides_output.pptx"


# Input file to process. First, copy your docx file into the docx2slides-py/resources folder,
# then update the name at the end of the next line from "sample_doc.docx" to the real name.
INPUT_DOCX_FILE = SCRIPT_DIR / "resources" / "sample_doc.docx"


# Which chunking method to use to divide the docx into slides. This enum lists the available choices:
class ChunkType(Enum):
    """Chunk Type Choices"""

    HEADING_NESTED = "heading_nested"
    HEADING_FLAT = "heading_flat"
    PARAGRAPH = "paragraph"
    PAGE = "page"


# And this is where to set what will be used in this run
CHUNK_TYPE: ChunkType = ChunkType.HEADING_FLAT


# Toggle on/off whether to print debug_prints() to the console
DEBUG_MODE = True  # TODO, v1 POLISH: set to false before publishing; update: TODO, UX: please god replace all these consts with a config class or something for v1


DISPLAY_COMMENTS: bool = True
DISPLAY_FOOTNOTES: bool = True
DISPLAY_ENDNOTES: bool = True

DISPLAY_DOCX_ANNOTATIONS_IN_SLIDE_SPEAKER_NOTES: bool = (
    DISPLAY_COMMENTS or DISPLAY_FOOTNOTES or DISPLAY_ENDNOTES
)

# We ought to support some way to leave speaker notes completely empty if the user really wants that, it's a valid use case.
# Documentation and tooltips should make it clear that this means metadata loss for round-trip pipeline data.
PRESERVE_DOCX_METADATA_IN_SPEAKER_NOTES: bool = True

COMMENTS_SORT_BY_DATE: bool = True
COMMENTS_KEEP_AUTHOR_AND_DATE: bool = True

EXPERIMENTAL_FORMATTING_ON: bool = True

# ========== pptx2docxtext pipeline consts

INPUT_PPTX_FILE = (
    SCRIPT_DIR / "resources" / "sample_slides_output.pptx"
)  # "sample_slides.pptx"

TEMPLATE_DOCX = SCRIPT_DIR / "resources" / "docx_template.docx"

OUTPUT_DOCX_FOLDER = SCRIPT_DIR / "output"
# e.g., r"c:\my_manuscripts"

OUTPUT_DOCX_FILENAME = r"sample_pptx2docxtext_output.docx"
# endregion


# region models.py
@dataclass
class Comment_docx_custom:
    """A custom wrapper for the python-docx Comment class, allowing us to capture reference text."""

    comment_obj: Comment_docx
    reference_text: str | None = None  # The text this comment is attached to

    @property
    def note_id(self) -> int:
        """Alias for comment_id to provide a common interface with other note types."""
        return self.comment_obj.comment_id


@dataclass
class Footnote_docx:
    """
    Represents a footnote extracted from a docx.

    Contains the footnote ID, text content, and any hyperlinks found within.
    Used for preserving footnote information when python-docx doesn't provide
    direct access to footnote content.
    """

    footnote_id: str
    text_body: str
    hyperlinks: list[str] = field(default_factory=list)
    reference_text: str | None = None

    @property
    def note_id(self) -> str:
        """Alias for footnote_id to provide a common interface with other note types."""
        return self.footnote_id


@dataclass
class Endnote_docx:
    """
    Represents a endnote extracted from a docx.

    Contains the endnote ID, text content, and any hyperlinks found within.
    Used for preserving endnote information when python-docx doesn't provide
    direct access to endnote content.
    """

    endnote_id: str
    text_body: str
    hyperlinks: list[str] = field(default_factory=list)
    reference_text: str | None = None

    @property
    def note_id(self) -> str:
        """Alias for endnote_id to provide a common interface with other note types."""
        return self.endnote_id


@dataclass
class Chunk_docx:
    """Class for Chunk objects made from docx paragraphs and their associated annotations."""

    # Use "default_factory" to ensure every chunk gets its own list.
    # (Lists are mutable; it is a common error/bug to accidentally assign one list
    # shared amongst every instance of a class, rather than one per instance.)
    paragraphs: list[Paragraph_docx] = field(default_factory=list[Paragraph_docx])

    comments: list[Comment_docx_custom] = field(
        default_factory=list[Comment_docx_custom]
    )
    footnotes: list[Footnote_docx] = field(default_factory=list[Footnote_docx])
    endnotes: list[Endnote_docx] = field(default_factory=list[Endnote_docx])

    @classmethod
    def create_with_paragraph(cls, paragraph: Paragraph_docx) -> "Chunk_docx":
        """Create a new instance of a Chunk_docx object but also populate the paragraphs list with an initial element."""
        return cls(paragraphs=[paragraph])

    def add_paragraph(self, new_paragraph: Paragraph_docx) -> None:
        """Add a paragraph to this Chunk object's paragraphs list."""
        self.paragraphs.append(new_paragraph)

    def add_paragraphs(self, new_paragraphs: list[Paragraph_docx]) -> None:
        """Add a list of paragraphs to this Chunk object's paragraphs list."""
        self.paragraphs.extend(new_paragraphs)  # Add multiple at once

    def add_comment(self, comment: Comment_docx_custom) -> None:
        """Add a comment to this Chunk object's comment list."""
        self.comments.append(comment)

    def add_footnote(self, footnote: Footnote_docx) -> None:
        """Add a footnote to this Chunk object's footnote list."""
        self.footnotes.append(footnote)

    def add_endnote(self, endnote: Endnote_docx) -> None:
        """Add a endnote to this Chunk object's endnote list."""
        self.endnotes.append(endnote)


@dataclass
class SlideNotes:
    """User notes and metadata extracted from a slide's speaker notes."""

    metadata: dict = field(default_factory=dict)
    user_notes: str = ""
    comments: list = field(default_factory=list)
    headings: list = field(default_factory=list)
    experimental_formatting: list = field(default_factory=list)

    @property
    def has_metadata(self) -> bool:
        """Returns a bool to indicate whether we did or did not find/store JSON metadata from these SlideNotes."""
        return bool(self.metadata)  # True if dict is non-empty

    @property
    def has_user_notes(self) -> bool:
        """
        Returns a bool to indicate whether we did or did not find/store unique user notes (not JSON metadata, and not
        copied annotations from earlier docx2pptx pipeline runs) from these SlideNotes.
        """
        return bool(self.user_notes.strip())


# endregion


# region __main__.py & run_pipeline.py
# eventual destination: ./src/docx2pptx/__main__.py
def main() -> None:
    """Entry point for program flow."""
    setup_console_encoding()
    debug_print("Hello, manuscript parser!")

    # run_docx2pptx_pipeline(INPUT_DOCX_FILE)

    run_pptx2docx_pipeline(INPUT_PPTX_FILE)


# endregion


# region pptx2docxtext
def run_pptx2docx_pipeline(pptx_path: Path) -> None:
    """Orchestrates the pptx2docxtext pipeline."""

    # Validate the user's pptx filepath
    try:
        validated_pptx_path = validate_pptx_path(pptx_path)
    except FileNotFoundError:
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)
    except ValueError as e:
        print(f"Error: {e}")
        sys.exit(1)
    except PermissionError:
        print(f"I don't have permission to read that file ({pptx_path})!")
        sys.exit(1)

    # Load the pptx at that validated filepath
    try:
        user_prs: presentation.Presentation = load_and_validate_pptx(
            validated_pptx_path
        )
    except Exception as e:
        print(
            f"Content of powerpoint file invalid for pptx2docxtext pipeline run. Error: {e}."
        )
        sys.exit(1)

    # Create an empty docx
    new_doc = docx.Document(str(TEMPLATE_DOCX))

    copy_slides_to_docx_body(user_prs, new_doc)

    debug_print("Attempting to save new docx file.")

    save_output(new_doc)


METADATA_MARKER_HEADER: str = "START OF JSON METADATA FROM SOURCE DOCUMENT"
METADATA_MARKER_FOOTER: str = "END OF JSON METADATA FROM SOURCE DOCUMENT"
NOTES_MARKER_HEADER: str = "START OF COPIED NOTES FROM SOURCE DOCX"
NOTES_MARKER_FOOTER: str = "END OF COPIED NOTES FROM SOURCE DOCX"

# endregion


# region pptx2docx-text helpers
def split_speaker_notes(speaker_notes_text: str) -> SlideNotes:
    """
    Extract user notes, JSON metadata, and discard copied annotations.
    Returns a SlideNotes object with parsed sections (or an empty object if parsing failed).
    """

    # Find all marker positions
    json_start = speaker_notes_text.find(METADATA_MARKER_HEADER)
    json_end = speaker_notes_text.find(METADATA_MARKER_FOOTER)

    notes_start = speaker_notes_text.find(NOTES_MARKER_HEADER)
    notes_end = speaker_notes_text.find(NOTES_MARKER_FOOTER)

    # Extract JSON if present
    json_content = None
    if json_start != -1 and json_end != -1:
        json_text = speaker_notes_text[
            json_start + len(METADATA_MARKER_HEADER) : json_end
        ]
        # Clean up both ends - remove separator lines
        json_text = json_text.strip().lstrip("=").rstrip("=").strip()
        try:
            json_content = json.loads(json_text)
            # TODO: add deeper JSON parsing. See `2025-09-24b JSON Validation TODO.md`
        except json.JSONDecodeError:
            json_content = None
            debug_print(
                "We found what looked like a JSON section but then failed to load it."
            )

    # Build list of ranges to remove (inclusive of markers)
    ranges_to_remove = []

    # If we found a range of JSON start/end, mark that for removal
    if json_start != -1 and json_end != -1:
        ranges_to_remove.append((json_start, json_end + len(METADATA_MARKER_FOOTER)))

    # If we found a range of notes start/end, mark that for removal
    if notes_start != -1 and notes_end != -1:
        ranges_to_remove.append((notes_start, notes_end + len(NOTES_MARKER_FOOTER)))

    # Sort ranges by start position (for consistent removal)
    ranges_to_remove.sort()

    # Extract user notes by removing the marked sections
    user_notes = remove_ranges_from_text(speaker_notes_text, ranges_to_remove)

    # Create and populate the chunk
    slide_notes = SlideNotes()
    slide_notes.user_notes = user_notes.strip()

    if json_content:
        slide_notes.metadata = json_content
        slide_notes.comments = json_content.get("comments", [])
        slide_notes.headings = json_content.get("headings", [])
        slide_notes.experimental_formatting = json_content.get(
            "experimental_formatting", []
        )

    return slide_notes


def remove_ranges_from_text(text: str, ranges: list) -> str:
    """Remove multiple ranges from text, working backwards to preserve positions."""

    # First merge any overlapping ranges
    ranges_merged = merge_overlapping_ranges(ranges)

    # Sort ranges by start position, then work backwards
    ranges_sorted = sorted(ranges_merged, reverse=True)  # Start from end

    result = text
    for start, end in ranges_sorted:
        result = result[:start] + result[end:]

    return result


def merge_overlapping_ranges(ranges: list) -> list:
    """If any (int, int) index ranges overlap, merge them."""
    if not ranges:
        return []

    # Sort by start position
    sorted_ranges = sorted(ranges)
    merged = [sorted_ranges[0]]

    for current in sorted_ranges[1:]:
        last = merged[-1]

        # Check if current overlaps with last merged range
        if current[0] <= last[1]:  # start of current <= end of last
            # Merge: extend the end position if needed
            merged[-1] = (last[0], max(last[1], current[1]))
        else:
            # No overlap, add as new range
            merged.append(current)

    return merged


def add_hyperlink_to_docx_paragraph(paragraph: Paragraph_docx, url: str) -> Run_docx:
    """
    Custom function to add Hyperlink objects to docx paragraphs using XML manipulation.

    - Create a regular run using paragraph.add_run()
    - Create the hyperlink XML element structure
    - Move the run's XML element from being a direct child of the paragraph to being a nested child of the Hyperlink element
    - Add the Hyperlink element (which now contains the run) to the paragraph

    Adapted from https://stackoverflow.com/questions/47666642/adding-an-hyperlink-in-msword-by-using-python-docx
    and https://github.com/python-openxml/python-docx/issues/384#issuecomment-294853130
    """
    # Create a new run on this paragraph
    run = paragraph.add_run()

    # Create the hyperlink structure
    part = paragraph.part
    r_id = part.relate_to(url, constants.RELATIONSHIP_TYPE.HYPERLINK, is_external=True)
    hyperlink = OxmlElement_docx("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)

    # Move the run from within the paragraph to within the Hyperlink
    run_element = run._element
    run_element.getparent().remove(run_element)  # Remove from paragraph
    hyperlink.append(run_element)  # Add to hyperlink
    paragraph._p.append(hyperlink)  # Add hyperlink to paragraph

    return run


def process_slide_paragraphs(
    slide: Slide, slide_notes: SlideNotes, new_doc: document.Document
) -> None:
    """
    Process a slide's body content, using any metadata stored in the speaker notes to restore formatting and annotation
    anchors from an earlier docx2pptx pipeline run. If we find metadata but aren't able to attach it to text content from
    the body paragraphs/runs, attach that as a new comment to the very last copied paragraph/run from that slide.

    Additionally, store non-metadata speaker notes content as a new comment, too.
    """

    slide_paragraphs: list[Paragraph_pptx] = get_slide_paragraphs(slide)

    unmatched_annotations = []
    matched_comment_ids = set()

    last_run = None

    # For every pptx paragraph.....
    for pptx_para in slide_paragraphs:

        # Make a new docx para
        new_para = new_doc.add_paragraph()

        # If the text of this paragraph exactly matches a previous heading's text, apply that heading style
        if slide_notes and slide_notes.has_metadata and slide_notes.headings:
            for heading in slide_notes.headings:
                if heading["text"].strip() == pptx_para.text.strip():
                    new_para.style = heading["name"]
                    break  # we should only ever apply one style to a paragraph

        for run in pptx_para.runs:

            # Handle adding hyperlinks versus regular runs, and the runs' basic formatting.
            if run.hyperlink.address:
                debug_print("Hyperlink address found.")
                run_from_hyperlink = add_hyperlink_to_docx_paragraph(
                    new_para, run.hyperlink.address
                )
                last_run = run_from_hyperlink
                copy_run_formatting_pptx2docx(run, run_from_hyperlink)
            else:
                new_docx_run = new_para.add_run()
                last_run = new_docx_run
                copy_run_formatting_pptx2docx(run, new_docx_run)

            # Check if this run contains matching text for comments from the speaker notes' stored JSON
            # metadata, from previous docx2pptx pipeline processing
            if slide_notes.comments:
                # Check to see if the run text matches any comments' ref text
                for comment in slide_notes.comments:
                    # if there is a match, create a new comment based on the original, attached to this run.
                    if (
                        comment["reference_text"] in run.text
                        and comment["id"] not in matched_comment_ids
                    ):
                        original = comment["original"]
                        text = original["text"]
                        author = original["author"]
                        initials = original["initials"]
                        new_doc.add_comment(new_docx_run, text, author, initials)
                        matched_comment_ids.add(comment["id"])
                    # don't break because there can be multiple comments added to a single run

            if EXPERIMENTAL_FORMATTING_ON and slide_notes.experimental_formatting:
                for exp_fmt in slide_notes.experimental_formatting:
                    if exp_fmt["ref_text"] in run.text:
                        _apply_experimental_formatting_from_metadata(
                            new_docx_run, exp_fmt
                        )

    # Find all the unmatched annotations by getting the complement set to each of the matched_comments set
    unmatched_comments = [
        c for c in slide_notes.comments if c["id"] not in matched_comment_ids
    ]

    unmatched_annotations.extend(unmatched_comments)

    # Put the slide's user notes into a new comment attached to the last run
    if slide_notes and slide_notes.has_user_notes is True and last_run is not None:
        # append as a comment to the last run
        raw_comment_text = slide_notes.user_notes
        comment_header = "Copied from the PPTX Speaker Notes: \n\n"
        comment_text = comment_header + sanitize_xml_text(raw_comment_text)

        if comment_text.strip():
            new_doc.add_comment(last_run, comment_text)

    # If we have any unmatched annotations from the slide_notes.metadata, attach that as a new comment to the last run
    if unmatched_annotations and last_run is not None:
        unmatched_parts = []
        # append as a comment to the last run
        for annotation in unmatched_annotations:
            if "original" in annotation:  # comments
                unmatched_parts.append(f"Comment: {annotation['original']['text']}")
            elif "text_body" in annotation:  # footnotes/endnotes
                unmatched_parts.append(f"Note: {annotation['text_body']}")
        combined = "\n\n".join(unmatched_parts)

        raw_comment_text = combined
        comment_header = "We found metadata for these annotations (comments, footnotes, or endnotes), but weren't able to match them to specific text in this paragraph: \n\n"
        comment_text = comment_header + sanitize_xml_text(raw_comment_text)

        if comment_text.strip():
            new_doc.add_comment(last_run, comment_text)

def _exp_fmt_issue(formatting_type: str, run_text: str, e: Exception) -> str:
    """Construct error message string per experimental formatting type."""
    message = f"We found a {formatting_type} in the experimental formatting JSON from a previous docx2pptx run, but we couldn't apply it. \n Run text: {run_text[:50]}... \n Error: {e}"
    return message

# TODO: test try/except and catch exceptions
def _apply_experimental_formatting_from_metadata(
    target_run: Run_docx, format_info: dict
) -> None:
    """Using JSON metadata from an earlier docx2pptx-text run, try to restore experimental formatting metadata to a run during the reverse pipeline."""

    tfont = target_run.font
    formatting_type = format_info.get("formatting_type")

    if formatting_type == "highlight":
        highlight_enum = format_info.get("highlight_color_enum")
        if highlight_enum:
            try:
                color_index = getattr(WD_COLOR_INDEX, highlight_enum, None)
                tfont.highlight_color = color_index
            except Exception as e:
                debug_print(
                    _exp_fmt_issue(formatting_type, target_run.text, e)
            )

    elif formatting_type == "strike":
        try:
            tfont.strike = True
        except Exception as e:
            debug_print(
                _exp_fmt_issue(formatting_type, target_run.text, e)
        )

    elif formatting_type == "double_strike":
        try:
            tfont.double_strike = True
        except Exception as e:
            debug_print(
                _exp_fmt_issue(formatting_type, target_run.text, e)
        )
        

    elif formatting_type == "subscript":
        try:
            tfont.subscript = True
        except Exception as e:
            debug_print(
                _exp_fmt_issue(formatting_type, target_run.text, e)
        )

    elif formatting_type == "superscript":
        try:
            tfont.superscript = True
        except Exception as e:
            debug_print(
                _exp_fmt_issue(formatting_type, target_run.text, e)
        )

    elif formatting_type == "all_caps":
        try:
            tfont.all_caps = True
        except Exception as e:
            debug_print(
                _exp_fmt_issue(formatting_type, target_run.text, e)
        )

    elif formatting_type == "small_caps":
        try:
            tfont.small_caps = True
        except Exception as e:
            debug_print(
                _exp_fmt_issue(formatting_type, target_run.text, e)
        )

def copy_slides_to_docx_body(
    prs: presentation.Presentation, new_doc: document.Document
) -> None:
    """
    Sequentially process each slide in the deck by copying the paragraphs from the slide body into the docx's body. Analyze
    the speaker notes of the slide, seeking stored JSON metadata there from previous docx2pptx runs, and use that to apply
    formatting and/or annotations.
    """

    # Make a list of all slides
    slide_list = list(prs.slides)

    # For each slide...
    for slide in slide_list:

        # If there's slide notes, process them into a SlideNotes object; otherwise, make an empty one.
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            slide_notes = split_speaker_notes(slide.notes_slide.notes_text_frame.text)
        else:
            slide_notes = SlideNotes()

        process_slide_paragraphs(slide, slide_notes, new_doc)


def copy_run_formatting_pptx2docx(source_run: Run_pptx, target_run: Run_docx) -> None:
    """Mutates a docx Run object to apply text and formatting from a pptx _Run object."""
    sfont = source_run.font
    tfont = target_run.font

    target_run.text = source_run.text

    _copy_basic_run_formatting(sfont, tfont)

    _copy_run_color_formatting(sfont, tfont)

    if source_run.text and source_run.text.strip() and EXPERIMENTAL_FORMATTING_ON:
        _copy_experimental_formatting_pptx2docx(source_run, target_run)


def _copy_experimental_formatting_pptx2docx(
    source_run: Run_pptx, target_run: Run_docx
) -> None:
    """
    Extract experimental formatting from the pptx _Run and attempt to apply it to the docx Run.
    (Unlike in the docx2pptx pipeline, we don't additionally store this as metadata anywhere.)
    """
    sfont = source_run.font
    tfont = target_run.font

    try:
        sfont_xml = sfont._element.xml

        # Quick string checks before parsing
        if (
            "strike=" not in sfont_xml
            and "baseline=" not in sfont_xml
            and "cap=" not in sfont_xml
            and "a:highlight" not in sfont_xml
        ):
            return  # No experimental formatting to apply

        root = ET.fromstring(sfont_xml)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        # Check for highlight nested element
        highlight = root.find(".//a:highlight/a:srgbClr", ns)
        if highlight is not None:
            debug_print("found a highlight")
            # Extract the color HEX out of the XML
            hex_color = highlight.get("val")
            if hex_color:
                # Convert the hex using the map map
                color_index = COLOR_MAP_FROM_HEX.get(hex_color)
                if color_index:
                    target_run.font.highlight_color = color_index

        # Check for strike/double-strike attribute
        strike = root.get("strike")
        if strike:
            if strike == "sngStrike":
                tfont.strike = True
            elif strike == "dblStrike":
                tfont.double_strike = True

        # Check for super/subscript
        baseline = root.get("baseline")
        if baseline:
            baseline_val = int(baseline)
            if baseline_val < 0:
                tfont.subscript = True
            elif baseline_val > 0:
                tfont.superscript = True

        # Check for all/small caps
        cap = root.get("cap")
        if cap:
            if cap == "all":
                tfont.all_caps = True
            elif cap == "small":
                tfont.small_caps = True

    except Exception as e:
        debug_print(f"Failed to parse pptx formatting from XML: {e}")


# endregion

# region ===========
# endregion

# region docxtext2pptx pipeline
"""
Below are all the functions written for the docx2pptx original pipeline flow; we'll start the reverse flow above.
"""
# endregion


# region run_docxtext2pptx_pipeline()
def run_docx2pptx_pipeline(docx_path: Path) -> None:
    """Orchestrates the docx2pptx pipeline."""
    user_path = docx_path

    # Validate it's a real path of the correct type. If it's not, return the error.
    try:
        user_path_validated = validate_docx_path(user_path)
    except FileNotFoundError:
        print(f"Error: File not found: {user_path}")
        sys.exit(1)
    except ValueError as e:
        print(f"Error: {e}")
        sys.exit(1)
    except PermissionError:
        print(f"I don't have permission to read that file ({user_path})!")
        sys.exit(1)

    # Load the docx file at that path.
    user_docx = load_and_validate_docx(user_path_validated)

    # Chunk the docx by ___
    chunks = create_docx_chunks(user_docx, CHUNK_TYPE)

    if PRESERVE_DOCX_METADATA_IN_SPEAKER_NOTES:
        chunks = process_chunk_annotations(chunks, user_docx)

    # Create the presentation object from template
    try:
        output_prs = create_empty_slide_deck()
    except Exception as e:
        print(f"Could not load template file (may be corrupted): {e}")
        sys.exit(1)

    # Mutate the presentation object by adding slides
    slides_from_chunks(user_docx, output_prs, chunks)

    # Save the presentation to an actual pptx on disk
    save_output(output_prs)


# endregion

# region create_slides.py
# eventual destination: ./src/docx2pptx/create_slides.py


def _copy_basic_run_formatting(
    source_font: Union[Font_docx, Font_pptx], target_font: Union[Font_docx, Font_pptx]
) -> None:
    """Extract common formatting logic for Runs."""

    # Bold/Italics: Only overwrite when explicitly set on the source (avoid clobbering inheritance)
    if source_font.bold is not None:
        target_font.bold = source_font.bold
    if source_font.italic is not None:
        target_font.italic = source_font.italic

    # Underline: collapse any explicit value (True/False/WD_UNDERLINE.*) to bool
    if source_font.underline is not None:
        target_font.underline = bool(source_font.underline)

    # TODO: THIS IS UNTESTED; TEST IT.
    if source_font.size is not None:
        target_font.size = Pt(source_font.size.pt)
        """
        <a:r>
            <a:rPr lang="en-US" sz="8800" i="1" dirty="0"/>
            <a:t>MAKE this text BIG!</a:t>
        </a:r>
        """


def _copy_run_color_formatting(
    source_font: Union[Font_docx, Font_pptx], target_font: Union[Font_docx, Font_pptx]
) -> None:
    # Color: copy only if source has an explicit RGB
    src_rgb = getattr(getattr(source_font, "color", None), "rgb", None)
    if src_rgb is not None:
        if isinstance(target_font, Font_pptx):
            target_font.color.rgb = RGBColor_pptx(*src_rgb)
        elif isinstance(target_font, Font_docx):
            target_font.color.rgb = RGBColor_docx(*src_rgb)


def copy_run_formatting_docx2pptx(
    source_run: Run_docx,
    target_run: Run_pptx,
    experimental_formatting_metadata: list,
) -> None:
    """Mutates a pptx _Run object to apply text and formatting from a docx Run object."""
    sfont = source_run.font
    tfont = target_run.font

    target_run.text = source_run.text

    _copy_basic_run_formatting(sfont, tfont)

    _copy_run_color_formatting(sfont, tfont)

    if EXPERIMENTAL_FORMATTING_ON:
        if source_run.text and source_run.text.strip():
            _copy_experimental_formatting_docx2pptx(
                source_run, target_run, experimental_formatting_metadata
            )


def _copy_experimental_formatting_docx2pptx(
    source_run: Run_docx,
    target_run: Run_pptx,
    experimental_formatting_metadata: list,
) -> None:
    """
    Extract experimental formatting from the docx Run and attempt to apply it to the pptx run. Additionally,
    store the formatting information in a metadata list (for the purpose of saving to JSON and enabling restoration
    during the reverse pipeline).
    """

    sfont = source_run.font
    tfont = target_run.font

    # The following code, which extends formatting support beyond python-pptx's capabilities,
    # is adapted from the md2pptx project, particularly from ./paragraph.py
    # Original source: https://github.com/MartinPacker/md2pptx
    # Author: Martin Packer
    # License: MIT
    if sfont.highlight_color is not None:
        experimental_formatting_metadata.append(
            {
                "ref_text": source_run.text,
                "highlight_color_enum": sfont.highlight_color.name,
                "formatting_type": "highlight",
            }
        )
        try:
            # Convert the docx run highlight color to a hex string
            tfont_hex_str = COLOR_MAP_HEX.get(sfont.highlight_color)

            # Create an object to represent this run in memory
            rPr = target_run._r.get_or_add_rPr()  # type: ignore[reportPrivateUsage]

            # Create a highlight Oxml object in memory
            hl = OxmlElement_pptx("a:highlight")

            # Create a srgbClr Oxml object in memory
            srgbClr = OxmlElement_pptx("a:srgbClr")

            # Set the attribute val of the srgbClr Oxml object in memory to the desired color
            setattr(srgbClr, "val", tfont_hex_str)

            # Add srgbClr object inside the hl Oxml object
            hl.append(srgbClr)  # type: ignore[reportPrivateUsage]

            # Add the hl object to the run representation object, which will add all our Oxml elements inside it
            rPr.append(hl)  # type: ignore[reportPrivateUsage]

        except Exception as e:
            
            debug_print(
                f"We found a highlight in a docx run but couldn't apply it. \n Run text: {source_run.text[:50]}... \n Error: {e}"
            )
        """
        Reference pptx XML for highlighting:
        <a:r>
            <a:rPr>
                <a:highlight>
                    <a:srgbClr val="FFFF00"/>
                </a:highlight>
            </a:rPr>
            <a:t>Highlight this text.</a:t>
        </a:r>
        """

    if sfont.strike is not None:
        experimental_formatting_metadata.append(
            {"ref_text": source_run.text, "formatting_type": "strike"}
        )
        try:
            tfont._element.set("strike", "sngStrike")  # type: ignore[reportPrivateUsage]
        except Exception as e:
            debug_print(
                f"Failed to apply single-strikethrough. \nRun text: {source_run.text[:50]}... \n Error: {e}"
            )

        """
        Reference pptx XML for single strikethrough:
        <a:p>
            <a:r>
                <a:rPr lang="en-US" strike="sngStrike" dirty="0"/>
                <a:t>Strike this text.</a:t>
            </a:r>
        </a:p>        
        """

    if sfont.double_strike is not None:
        experimental_formatting_metadata.append(
            {"ref_text": source_run.text, "formatting_type": "double_strike"}
        )
        try:
            tfont._element.set("strike", "dblStrike")  # type: ignore[reportPrivateUsage]
        except Exception as e:
            debug_print(
                f"""
                        Failed to apply double-strikthrough.
                        \nRun text: {source_run.text[:50]}... \n Error: {e}
                        \nWe'll attempt single strikethrough."""
            )
            tfont._element.set("strike", "sngStrike")  # type: ignore[reportPrivateUsage]
        """
        Reference pptx XML for double strikethrough:
        <a:p>
            <a:r>
                <a:rPr lang="en-US" strike="dblStrike" dirty="0" err="1"/>
                <a:t>Double strike this text.</a:t>
            </a:r>        
        </a:p>
        """

    if sfont.subscript is not None:
        experimental_formatting_metadata.append(
            {"ref_text": source_run.text, "formatting_type": "subscript"}
        )
        try:
            if tfont.size is None or tfont.size < Pt(24):
                tfont._element.set("baseline", BASELINE_SUBSCRIPT_SMALL_FONT)  # type: ignore[reportPrivateUsage]
            else:
                tfont._element.set("baseline", BASELINE_SUBSCRIPT_LARGE_FONT)  # type: ignore[reportPrivateUsage]

        except Exception as e:
            debug_print(
                f"""
                        Failed to apply subscript. 
                        \nRun text: {source_run.text[:50]}... 
                        \n Error: {e}"""
            )
        """
        Reference pptx XML for subscript:
        <a:r>
            <a:rPr lang="en-US" baseline="-25000" dirty="0" err="1"/>
            <a:t>Subscripted text</a:t>
        </a:r>
        """

    if sfont.superscript is not None:
        experimental_formatting_metadata.append(
            {"ref_text": source_run.text, "formatting_type": "superscript"}
        )
        try:
            if tfont.size is None or tfont.size < Pt(24):
                tfont._element.set("baseline", BASELINE_SUPERSCRIPT_SMALL_FONT)  # type: ignore[reportPrivateUsage]
            else:
                tfont._element.set("baseline", BASELINE_SUPERSCRIPT_LARGE_FONT)  # type: ignore[reportPrivateUsage]

        except Exception as e:
            debug_print(
                f"""
                        Failed to apply superscript. 
                        \nRun text: {source_run.text[:50]}... 
                        \n Error: {e}"""
            )
        """
        Reference pptx XML for superscript
        <a:r>
            <a:rPr lang="en-US" baseline="30000" dirty="0" err="1"/>
            <a:t>Superscript this text.</a:t>
        </a:r>
        """

    # The below caps-handling code is not directly from md2pptx,
    # but is heavily influenced by it.
    if sfont.all_caps is not None:
        experimental_formatting_metadata.append(
            {"ref_text": source_run.text, "formatting_type": "all_caps"}
        )
        try:
            tfont._element.set("cap", "all")  # type: ignore[reportPrivateUsage]
        except Exception as e:
            debug_print(
                f"""
                        Failed to apply all caps. 
                        \nRun text: {source_run.text[:50]}... 
                        \n Error: {e}"""
            )
        """
        Reference XML for all caps:
        <a:r>
            <a:rPr lang="en-US" cap="all" dirty="0" err="1"/>
            <a:t>Put this text in all caps.</a:t>
        </a:r>
        """

    if sfont.small_caps is not None:
        experimental_formatting_metadata.append(
            {"ref_text": source_run.text, "formatting_type": "small_caps"}
        )
        try:
            tfont._element.set("cap", "small")  # type: ignore[reportPrivateUsage]
        except Exception as e:
            debug_print(
                f"""
                        Failed to apply small caps on run with text body: 
                        \nRun text: {source_run.text[:50]}... 
                        \n Error: {e}"""
            )
        """
        Reference pptx XML for small caps:
        <a:r>
            <a:rPr lang="en-US" cap="small" dirty="0" err="1"/>
            <a:t>Put this text in small caps.</a:t>
        </a:r>
        """


def create_blank_slide_for_chunk(
    prs: presentation.Presentation, slide_layout: SlideLayout
) -> tuple[Slide, TextFrame]:
    """Initialize an empty slide so that we can populate it with a chunk."""
    new_slide = prs.slides.add_slide(slide_layout)
    content = new_slide.placeholders[1]

    if not isinstance(content, SlidePlaceholder):
        raise TypeError(f"Expected SlidePlaceholder, got {type(content)}")

    text_frame: TextFrame = content.text_frame
    text_frame.clear()

    # Access the slide's notes_slide attribute in order to initialize it.
    notes_slide_ptr = new_slide.notes_slide  # type: ignore # noqa

    return new_slide, text_frame


def slides_from_chunks(
    doc: document.Document,
    prs: presentation.Presentation,
    chunks: list[Chunk_docx],
) -> None:
    """Generate slide objects, one for each chunk created by earlier pipeline steps."""

    # Specify which slide layout to use
    slide_layout = prs.slide_layouts.get_by_name(SLD_LAYOUT_CUSTOM_NAME)

    if slide_layout is None:
        raise KeyError(
            f"No slide layout found to match provided custom name, {SLD_LAYOUT_CUSTOM_NAME}"
        )

    for chunk in chunks:
        # Create a new slide for this chunk.
        new_slide, text_frame = create_blank_slide_for_chunk(prs, slide_layout)

        # Store custom metadata for this chunk that we'll want to tuck into the speaker notes as JSON
        # (for the purposes of restoring during reverse pipeline runs).
        slide_metadata = {}
        headings = []
        experimental_formatting: list[dict] = []

        # For each paragraph in this chunk, handle adding it
        for i, paragraph in enumerate(chunk.paragraphs):
            
            # Creating a new slide and a text frame leaves an empty paragraph in place, even when clearing it.
            # So if we're at the start of our list, use that existing empty paragraph.
            if (
                i == 0
                and len(text_frame.paragraphs) > 0
                and not text_frame.paragraphs[0].text
            ):
                # Use the existing first/0th paragraph
                pptx_paragraph = text_frame.paragraphs[0]
            else:
                pptx_paragraph = text_frame.add_paragraph()

            # Process the docx's paragraph contents, including both runs & hyperlinks
            para_experimental_formatting = process_chunk_paragraph_inner_contents(
                paragraph, pptx_paragraph
            )

            if para_experimental_formatting:
                experimental_formatting.extend(para_experimental_formatting)

            if (
                paragraph.style
                and paragraph.style.name
                and is_standard_heading(paragraph.style.name)
            ):
                headings.append(
                    {
                        "text": paragraph.text.strip(),
                        "style_id": paragraph.style.style_id,
                        "name": paragraph.style.name,
                    }
                )

        if headings:
            slide_metadata["headings"] = headings
        if experimental_formatting:
            slide_metadata["experimental_formatting"] = experimental_formatting

        notes_text_frame = new_slide.notes_slide.notes_text_frame

        if notes_text_frame is None:
            raise ValueError("This slide doesn't seem to have a notes text frame. This should never happen, but it's possible for the notes_slide or notes_text_frame properties to return None if the notes placeholder has been removed from the notes master or the notes slide itself.")

        if DISPLAY_DOCX_ANNOTATIONS_IN_SLIDE_SPEAKER_NOTES:
            annotate_slide(chunk, notes_text_frame)

        if PRESERVE_DOCX_METADATA_IN_SPEAKER_NOTES:
            add_metadata_to_slide_notes(notes_text_frame, chunk, slide_metadata)


def add_metadata_to_slide_notes(
    notes_text_frame: TextFrame, chunk: Chunk_docx, slide_body_metadata: dict
) -> None:
    """
    Populate the slide notes text frame with docx metadata so that we may restore it during a round-trip pptx2docx pipeline.
    """
    comments = []

    if chunk.comments:
        comments = [
            {
                "original": {
                    "text": c.comment_obj.text,
                    "author": c.comment_obj.author,
                    "comment_id": c.comment_obj.comment_id,
                    "initials": c.comment_obj.initials,
                    "timestamp": (
                        str(c.comment_obj.timestamp)
                        if c.comment_obj.timestamp
                        else None
                    ),
                },
                "reference_text": c.reference_text,
                "id": c.note_id,
            }
            for c in chunk.comments
        ]

    if slide_body_metadata or comments:
        header_para = notes_text_frame.add_paragraph()
        header_run = header_para.add_run()
        header_run.text = f"\n\n\n\n\n\n\n{METADATA_MARKER_HEADER}\n" + "=" * 40

        notes_text_frame.add_paragraph()  # blank paragraph to ensure separation for JSON block

        json_para = notes_text_frame.add_paragraph()
        json_run = json_para.add_run()

        combined_metadata = {**slide_body_metadata}

        if comments:  # only add if the list has content
            combined_metadata["comments"] = comments

        json_run.text = json.dumps(combined_metadata, indent=2)

        footer_para = notes_text_frame.add_paragraph()
        footer_run = footer_para.add_run()
        footer_run.text = "=" * 40 + f"\n{METADATA_MARKER_FOOTER}"


def process_chunk_paragraph_inner_contents(
    paragraph: Paragraph_docx, pptx_paragraph: Paragraph_pptx
) -> list[dict]:
    """
    Iterate through a paragraph's runs and hyperlinks, in document order, and:
    - copy it into the slide, with formatting
    - capture any additional metadata (like experimental formatting) that we cannot apply directly to the copied-over
        pptx objects
    """
    items_processed = False
    experimental_formatting_metadata = []

    for item in paragraph.iter_inner_content():
        items_processed = True
        if isinstance(item, Run_docx):

            # If this Run has a field code for instrText and it begins with HYPERLINK, this is an old-style
            # word hyperlink, which we cannot handle the same way as normal docx hyperlinks. But we try to detect
            # when it happens and report it to the user.
            field_code_URL = detect_field_code_hyperlinks(item)
            if field_code_URL:
                item.text = f" [Link: {field_code_URL}] "

            process_run(item, pptx_paragraph, experimental_formatting_metadata)

        # Run and Hyperlink objects are peers in docx, but Hyperlinks can contain lists of Runs.
        # We check the item.url field because that seems the most reliable way to see if this is a
        # basic run versus a Hyperlink containing its own nested runs.
        # https://python-docx.readthedocs.io/en/latest/api/text.html#docx.text.hyperlink.Hyperlink.url
        elif hasattr(item, "url"):
            # Process all runs within the Hyperlink
            for run in item.runs:
                process_run(
                    run,
                    pptx_paragraph,
                    experimental_formatting_metadata,
                    item.url,
                )
        # elif hasattr(item, "fragment"):
        #   ...
        #   TODO, leafy: We need to handle document anchors differently from other hyperlinks.
        #   We need to 1) process the nested runs as if it were a Hyperlink object, and
        #   2) preserve the anchor somewhere, maybe the experimental formatting metadata.
        #   https://python-docx.readthedocs.io/en/latest/api/text.html#docx.text.hyperlink.Hyperlink.fragment

        else:
            debug_print(f"Unknown content type in paragraph: {type(item)}")

    # Fallback: if no content was processed but paragraph has text
    if not items_processed and paragraph.text:
        debug_print(
            f"Fallback: paragraph has text but no runs/hyperlinks: {paragraph.text[:50]}"
        )
        pptx_run = pptx_paragraph.add_run()
        pptx_run.text = paragraph.text

    return experimental_formatting_metadata


def detect_field_code_hyperlinks(run: Run_docx) -> None | str:
    """
    Detect if this Run has a field code for instrText and it begins with HYPERLINK.
    If so, report it to the user, because we do not handle adding these to the pptx output.
    """
    try:
        run_xml: str = run.element.xml
        if "instrText" not in run_xml or "HYPERLINK" not in run_xml:
            return None
        root = ET.fromstring(run_xml)

        # Find instrText elements
        instr_texts = root.findall(
            ".//w:instrText",
            {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"},
        )
        for instr in instr_texts:
            if instr.text and instr.text.startswith("HYPERLINK"):
                match = re.search(r'HYPERLINK\s+"([^"]+)"', instr.text)
                if match and match.group(1):
                    return match.group(1)

    except (AttributeError, ET.ParseError) as e:
        debug_print(
            f"WARNING: Could not parse run XML for field codes: {e} while seeking instrText"
        )

    return None


def process_run(
    run: Run_docx,
    pptx_paragraph: Paragraph_pptx,
    experimental_formatting_metadata: list,
    hyperlink: str | None = None,
) -> Run_pptx:
    """Copy a run from the docx parent to the pptx paragraph, including copying its formatting."""
    # Handle formatting

    pptx_run = pptx_paragraph.add_run()
    copy_run_formatting_docx2pptx(run, pptx_run, experimental_formatting_metadata)

    if hyperlink:
        pptx_run_url = pptx_run.hyperlink
        pptx_run_url.address = hyperlink

    return pptx_run


# endregion

# region GET annotation helpers
# eventual destination: ./src/docx2pptx/annotations/annotate_chunks.py
# endregion


# region annotate_chunks.py
def process_chunk_annotations(
    chunks: list[Chunk_docx], doc: document.Document
) -> list[Chunk_docx]:
    """For a list of Chunk_docx objects, populate the annotation dicts for each one."""

    # Gather all the doc annotations
    all_raw_comments = get_all_docx_comments(doc)
    all_footnotes = get_all_docx_footnotes(doc)
    all_endnotes = get_all_docx_endnotes(doc)

    for chunk in chunks:
        for paragraph in chunk.paragraphs:
            for item in paragraph.iter_inner_content():

                if isinstance(item, Run_docx):
                    process_run_annotations(
                        chunk,
                        paragraph,
                        item,
                        all_raw_comments=all_raw_comments,
                        all_footnotes=all_footnotes,
                        all_endnotes=all_endnotes,
                    )

                # If the item has the attribute "url" we assume it is of type Hyperlink instead of Run;
                # that means it contains its own child runs, so we need to go one step inward to process them.
                elif hasattr(item, "url"):
                    # Process all runs within the hyperlink
                    for run in item.runs:
                        process_run_annotations(
                            chunk,
                            paragraph,
                            run,
                            all_raw_comments=all_raw_comments,
                            all_footnotes=all_footnotes,
                            all_endnotes=all_endnotes,
                        )
                else:
                    debug_print(f"Unknown content type in paragraph: {type(item)}")

    return chunks


def process_run_annotations(
    chunk: Chunk_docx,
    paragraph: Paragraph_docx,
    run: Run_docx,
    all_raw_comments: dict[str, Comment_docx],
    all_footnotes: dict[str, Footnote_docx],
    all_endnotes: dict[str, Endnote_docx],
) -> None:
    """Get the annotations from a run object and adding them into its (grand)parent chunk object."""
    try:
        # Get XML from the run using public API
        run_xml = run.element.xml

        # Parse it safely with ElementTree
        root = ET.fromstring(run_xml)

        # Define namespace
        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

        # get the reference text to be used by comments, footnotes, or endnotes
        ref_text = get_ref_text(run, paragraph)

        # Find comment references
        comment_refs = root.findall(".//w:commentReference", ns)
        for ref in comment_refs:
            comment_id = ref.get(f'{{{ns["w"]}}}id')
            if comment_id and comment_id in all_raw_comments:
                comment_object = all_raw_comments[comment_id]

                custom_comment_obj = Comment_docx_custom(
                    comment_obj=comment_object, reference_text=ref_text
                )
                chunk.add_comment(custom_comment_obj)

        # Find footnote references
        footnote_refs = root.findall(".//w:footnoteReference", ns)
        for ref in footnote_refs:
            footnote_id = ref.get(f'{{{ns["w"]}}}id')
            if footnote_id and footnote_id in all_footnotes:
                footnote_obj = all_footnotes[footnote_id]
                footnote_obj.reference_text = ref_text
                chunk.add_footnote(footnote_obj)

        # Find endnote references - same pattern
        endnote_refs = root.findall(".//w:endnoteReference", ns)
        for ref in endnote_refs:
            endnote_id = ref.get(f'{{{ns["w"]}}}id')
            if endnote_id and endnote_id in all_endnotes:
                endnote_obj = all_endnotes[endnote_id]
                endnote_obj.reference_text = ref_text
                chunk.add_endnote(endnote_obj)

    except (AttributeError, ET.ParseError) as e:
        debug_print(f"WARNING: Could not parse run XML for references: {e}")


# endregion


def get_ref_text(run: Run_docx, paragraph: Paragraph_docx) -> str | None:
    """
    Get the Run or Paragraph text with which a piece of metadata is associated in the docx so that we can store that in
    metadata and reference it on reverse-pipeline runs.
    """
    if run.text and run.text.strip():
        ref_text = run.text
    elif paragraph.text and paragraph.text.strip():
        # Grab the first (up to 10) words of this paragraph if the run text is empty
        ref_text = " ".join(paragraph.text.split()[:10])
    else:
        ref_text = None

    return ref_text


# eventual destination: ./src/docx2pptx/annotations/get_docx_annotations.py
# region Get all notes in this docx
def get_all_docx_comments(doc: document.Document) -> dict[str, Comment_docx]:
    """
    Get all the comments in this document as a dictionary.

    Elements of the dictionary are formatted as:
    {
        "comment_id_#" : this_comment_docx_object,
        "3": "<docx.comments.Comment object at 0x00000###>
    }
    """
    all_comments_dict: dict[str, Comment_docx] = {}

    if hasattr(doc, "comments") and doc.comments:
        for comment in doc.comments:
            all_comments_dict[str(comment.comment_id)] = comment
    return all_comments_dict


def get_all_docx_footnotes(doc: document.Document) -> dict[str, Footnote_docx]:
    """
    Extract all footnotes from a docx document.
    Returns {id: {footnote_id: str, text_body: str, hyperlinks: list of str} }.
    """

    if not DISPLAY_FOOTNOTES:
        return {}

    try:
        footnotes_parts = find_xml_parts(doc, "footnotes.xml")

        if not footnotes_parts:
            return {}

        # We think this will always be a list of one item, so assign that item to a variable.
        root = parse_xml_blob(footnotes_parts[0].blob)
        return extract_notes_from_xml(root, Footnote_docx)

    except Exception as e:
        debug_print(f"Warning: Could not extract footnotes: {e}")
        return {}


def get_all_docx_endnotes(doc: document.Document) -> dict[str, Endnote_docx]:
    """
    Extract all endnotes from a docx document.
    Returns {id: {footnote_id: str, text_body: str, hyperlinks: list of str} }.
    """
    if not DISPLAY_ENDNOTES:
        return {}

    try:
        endnotes_parts = find_xml_parts(doc, "endnotes.xml")

        if not endnotes_parts:
            return {}

        root = parse_xml_blob(endnotes_parts[0].blob)
        return extract_notes_from_xml(root, Endnote_docx)

    except Exception as e:
        debug_print(f"Warning: Could not extract endnotes: {e}")
        return {}


# endregion


# region Utils - XML parsing
def find_xml_parts(doc: document.Document, part_name: str) -> list[Part]:
    """Find XML parts matching the given name (e.g., 'footnotes.xml')"""
    # The zip package inspection logic
    # Inspect the docx package as a zip
    zip_package = doc.part.package

    if zip_package is None:
        debug_print("WARNING: Could not access docx package.")
        return []

    part_name_parts: list[Part] = []
    for part in zip_package.parts:
        if part_name in str(part.partname):
            debug_print(f"We found a {part_name} part!")
            part_name_parts.append(part)

    return part_name_parts


def parse_xml_blob(xml_blob: bytes | str) -> ET.Element:
    """Parse an XML blob into a string, from bytes."""
    if isinstance(xml_blob, str):
        xml_string = xml_blob
    else:
        # If footnote_blob is in bytes, or is bytes-like,
        # convert it to a string
        xml_string = bytes(xml_blob).decode("utf-8")

    # Create an ElementTree object by deserializing the footnotes.xml contents into a Python object
    root: ET.Element = ET.fromstring(xml_string)

    return root


# TODO, multi-file split: move to the top of whatever file this function ends up living in
# This allows for a generic type parameter - when you pass Footnote_docx into the extract_notes_from_xml(...) function, you will get dict[str, Footnote_docx] back
NOTE_TYPE = TypeVar("NOTE_TYPE", Footnote_docx, Endnote_docx)


def extract_notes_from_xml(
    root: ET.Element, note_class: type[NOTE_TYPE]
) -> dict[str, NOTE_TYPE]:
    """Extract footnotes or endnotes from XML, depending on note_class provided."""

    # Construct the strings we need to use in the XML search.
    # First, define the prefix and the namespace to which it will refer.
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

    # Second, construct the uri as a lookup in that dict to match how the XML works
    namespace_uri = ns["w"]

    # Third, construct the actual lookup strings. These are the full attribute name we're looking for in the data structure.
    # We must use double-curly braces to indicate we want a real curly brace in the node string.
    # And we also need an outer curly brace pair for the f-string syntax. That's why there's 3 total.
    id_attribute = f"{{{namespace_uri}}}id"  # "{http://...}id"
    type_attribute = f"{{{namespace_uri}}}type"

    notes_dict: dict[str, NOTE_TYPE] = {}

    for note in root:
        note_id = note.get(id_attribute)
        note_type = note.get(type_attribute)

        if note_id is None or note_type in ["separator", "continuationSeparator"]:
            continue

        note_full_text = "".join(note.itertext())
        note_hyperlinks = extract_hyperlinks_from_note(note)

        note_obj = note_class(
            note_id, text_body=note_full_text, hyperlinks=note_hyperlinks
        )

        notes_dict[note_id] = note_obj

    return notes_dict


def extract_hyperlinks_from_note(element: ET.Element) -> list[str]:
    """Extract all hyperlinks from a footnote element."""
    hyperlinks: list[str] = []
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

    for hyperlink in element.findall(".//w:hyperlink", ns):
        # Get the link text
        link_text = "".join([t.text or "" for t in hyperlink.findall(".//w:t", ns)])
        if link_text.strip():
            hyperlinks.append(link_text.strip())

    return hyperlinks


# endregion


# region annotate_slides.py
# eventual destination: ./src/docx2pptx/annotations/annotate_slides.py


def annotate_slide(chunk: Chunk_docx, notes_text_frame: TextFrame) -> None:
    """
    Pull a chunk's preserved annotations and copy them into the slide's speaker notes text frame.

    NOTE: We DO NOT PRESERVE any anchoring to the slide's body for annotations. That means we don't preserve
    comments' selected text ranges, nor do we preserve footnote or endnote numbering.
    """
    notes_text_frame.add_paragraph()  # add a blank first line for actual annotations by the user

    header_para = notes_text_frame.add_paragraph()
    header_run = header_para.add_run()
    header_run.text = f"\n\n\n\n\n\n\n{NOTES_MARKER_HEADER}\n" + "=" * 40 + "\n"

    if DISPLAY_COMMENTS and chunk.comments:
        add_comments_to_notes(chunk.comments, notes_text_frame)

    if DISPLAY_FOOTNOTES and chunk.footnotes:
        add_notes_to_speaker_notes(chunk.footnotes, notes_text_frame, Footnote_docx)

    if DISPLAY_ENDNOTES and chunk.endnotes:
        add_notes_to_speaker_notes(chunk.endnotes, notes_text_frame, Endnote_docx)

    footer_para = notes_text_frame.add_paragraph()
    footer_run = footer_para.add_run()
    footer_run.text = "=" * 40 + f"\n{NOTES_MARKER_FOOTER}"


# region Add annotations to pptx speaker notes text frame
def add_comments_to_notes(
    comments_list: list[Comment_docx_custom], notes_text_frame: TextFrame
) -> None:
    """Copy logic for appending the comments portion of the speaker notes."""

    if comments_list:
        if COMMENTS_SORT_BY_DATE:
            # Sort comments by date (newest first, or change reverse=False for oldest first)
            sorted_comments = sorted(
                comments_list,
                key=lambda c: getattr(c.comment_obj, "timestamp", None) or datetime.min,
                reverse=False,
            )
        else:
            sorted_comments = comments_list

        comment_para = notes_text_frame.add_paragraph()
        comment_run = comment_para.add_run()
        comment_run.text = "COMMENTS FROM SOURCE DOCUMENT:\n" + "=" * 40

        for i, comment in enumerate(sorted_comments, 1):
            # Check if comment has paragraphs attribute
            if hasattr(comment.comment_obj, "paragraphs"):
                # Get paragraphs safely, default to empty list if not present
                this_comment_paragraphs = getattr(comment.comment_obj, "paragraphs", [])
                for para in this_comment_paragraphs:
                    if hasattr(para, "text") and para.text.rstrip():
                        notes_para = notes_text_frame.add_paragraph()
                        comment_header = notes_para.add_run()

                        if COMMENTS_KEEP_AUTHOR_AND_DATE:
                            author = getattr(
                                comment.comment_obj, "author", "Unknown Author"
                            )
                            timestamp = getattr(comment.comment_obj, "timestamp", None)

                            if timestamp and hasattr(timestamp, "strftime"):
                                timestamp_str = timestamp.strftime(
                                    "%A, %B %d, %Y at %I:%M %p"
                                )
                            else:
                                timestamp_str = "Unknown Date"

                            comment_header.text = (
                                f"\n {i}. {author} ({timestamp_str}):\n"
                            )

                        else:
                            comment_header.text = "\n"
                        process_chunk_paragraph_inner_contents(para, notes_para)


def add_notes_to_speaker_notes(
    notes_list: list[NOTE_TYPE],
    notes_text_frame: TextFrame,
    note_class: type[NOTE_TYPE],
) -> None:
    """Generic function for adding footnotes or endnotes to speaker notes."""

    if note_class is Footnote_docx:
        header_text = "FOOTNOTES FROM SOURCE DOCUMENT"
    elif note_class is Endnote_docx:
        header_text = "ENDNOTES FROM SOURCE DOCUMENT"
    else:
        header_text = "Unknown Note Type from Source Document:"

    if notes_list:
        note_para = notes_text_frame.add_paragraph()
        note_run = note_para.add_run()
        note_run.text = f"\n{header_text}:\n" + "=" * 40

        for note_obj in notes_list:
            notes_para = notes_text_frame.add_paragraph()
            note_run = notes_para.add_run()

            # Start with the main note text
            note_text = f"\n{note_obj.note_id}. {note_obj.text_body}\n"

            # Add hyperlinks if they exist
            if note_obj.hyperlinks:
                note_text += "\nHyperlinks:"
                for hyperlink in note_obj.hyperlinks:
                    note_text += f"\n{hyperlink}"

            note_run.text = note_text


# endregion


# region create_chunks.py
# eventual destination: ./src/docx2pptx/create_chunks.py
def create_docx_chunks(
    doc: document.Document, chunk_type: ChunkType = ChunkType.PARAGRAPH
) -> list[Chunk_docx]:
    """
    Orchestrator function to create chunks (that will become slides) from the document
    contents, either from paragraph, heading (heading_nested or heading_flat),
    or page. Defaults to paragraph.
    """
    if chunk_type == ChunkType.HEADING_FLAT:
        chunks = chunk_by_heading_flat(doc)
    elif chunk_type == ChunkType.HEADING_NESTED:
        chunks = chunk_by_heading_nested(doc)
    elif chunk_type == ChunkType.PAGE:
        chunks = chunk_by_page(doc)
    else:
        chunks = chunk_by_paragraph(doc)
    return chunks


# endregion


# region by Paragraph
def chunk_by_paragraph(doc: document.Document) -> list[Chunk_docx]:
    """
    Creates chunks (which will become slides) based on paragraph, which are blocks of content
    separated by whitespace.
    """
    paragraph_chunks: list[Chunk_docx] = []

    for para in doc.paragraphs:

        # Skip empty paragraphs (but keep those that are new-lines to respect intentional whitespace newlines)
        if para.text == "":
            continue

        new_chunk = Chunk_docx.create_with_paragraph(para)
        paragraph_chunks.append(new_chunk)

    return paragraph_chunks


# endregion


# region by Page
def chunk_by_page(doc: document.Document) -> list[Chunk_docx]:
    """Creates chunks based on page breaks"""

    # Start building the chunks
    all_chunks: list[Chunk_docx] = []

    # Start with a current chunk ready-to-go
    current_page_chunk: Chunk_docx = Chunk_docx()

    for para in doc.paragraphs:
        # Skip empty paragraphs (keep intentional whitespace newlines)
        if para.text == "":
            continue

        # If the current_page_chunk is empty, append the current para regardless of style & continue to next para.
        if not current_page_chunk.paragraphs:
            current_page_chunk.add_paragraph(para)
            continue

        # Handle page breaks - create new chunk and start fresh
        if para.contains_page_break:
            # Add the current_chunk to chunks list (if it has content)
            if current_page_chunk:
                all_chunks.append(current_page_chunk)

            # Start new chunk with this paragraph
            current_page_chunk = Chunk_docx.create_with_paragraph(para)

            continue

        # If there was no page break, just append this paragraph to the current_chunk
        current_page_chunk.add_paragraph(para)

    # Ensure final chunk from loop is added to chunks list
    if current_page_chunk:
        all_chunks.append(current_page_chunk)

    print(f"This document has {len(all_chunks)} page chunks.")
    return all_chunks


# endregion

# region by Heading (nested)


## === chunk by heading helpers ===


def is_standard_heading(style_name: str) -> bool:
    """Check if paragraph.style.name is a standard Word Heading (Heading 1, Heading 2, ..., Heading 6)"""
    return style_name.startswith("Heading") and style_name[8:].isdigit()


def get_heading_level(style_name: str) -> int | float:
    """
    Extract the numeric level from a heading style name (e.g., 'Heading 2' -> 2),
    or return infinity if the style name doesn't have a number.
    """
    try:
        return int(style_name[8:])
    except (ValueError, IndexError):
        return float("inf")  # Treat non-headings as "deepest possible"


## ===


def chunk_by_heading_nested(doc: document.Document) -> list[Chunk_docx]:
    """
    Creates chunks based on headings, using nesting logic to group "deeper" headings

    New chunks are begun when:
    - a page break happens in the middle of a paragraph
    - we reach a heading-depth that is equal to or "higher" than (number is less than) the current depth
    Otherwise, appends paragraph to the current chunk.

    E.g.:

    CHUNK:
    H1
    Normal Paragraph
    H2
    Normal Paragraph
    Normal Paragraph

    NEXT_CHUNK:
    H2
    Normal Paragraph
    Normal Paragraph
    Normal Paragraph

    NEXT_CHUNK:
    H1
    H2
    Pararaph
    Normal Paragraph
    H3
    Normal Paragraph

    NEXT_CHUNK:
    H2
    Normal Paragraph
    Normal Paragraph

    """
    # Start building the chunks
    all_chunks: list[Chunk_docx] = []
    current_chunk: Chunk_docx = Chunk_docx()

    # Initialize current_heading_style_name
    current_heading_style_name = "Normal"  # Default for documents without headings

    for i, para in enumerate(doc.paragraphs):

        # Skip empty paragraphs
        if para.text == "":
            continue

        # Set a style_name to make Pylance happy (it gets mad if we direct-check para.style.style_name later)
        style_name = para.style.name if para.style and para.style.name else "Normal"

        debug_print(f"Paragraph begins: {para.text[:30]}... and is index: {i}")

        # If the current_chunk is empty, append the current para regardless of style & continue to next para.
        if not current_chunk.paragraphs:
            current_chunk.add_paragraph(para)
            if is_standard_heading(style_name):
                current_heading_style_name = style_name
            continue

        # Handle page breaks - create new chunk and start fresh
        if para.contains_page_break:
            # Add the current chunk to chunks list (if it has content)
            if current_chunk:
                all_chunks.append(current_chunk)

            # Start new chunk with this paragraph
            current_chunk = Chunk_docx.create_with_paragraph(para)

            # Update heading depth if this paragraph is a heading
            if is_standard_heading(style_name):
                current_heading_style_name = style_name
            continue

        # Handle headings
        if is_standard_heading(style_name):
            # Check if this heading is at same level or higher (less deep) than current. Smaller numbers are higher up in the hierarchy.
            if get_heading_level(style_name) <= get_heading_level(
                current_heading_style_name
            ):
                # If yes, start a new chunk
                if current_chunk:
                    all_chunks.append(current_chunk)
                current_chunk = Chunk_docx.create_with_paragraph(para)
                current_heading_style_name = style_name
            else:
                # This heading is deeper, add to current chunk
                current_chunk.add_paragraph(para)
                current_heading_style_name = style_name
        else:
            # Normal paragraph - add to current chunk
            current_chunk.add_paragraph(para)

    if current_chunk:
        all_chunks.append(current_chunk)

    print(f"This document has {len(all_chunks)} nested heading chunks.")
    return all_chunks


# endregion


# region by Heading (flat)
def chunk_by_heading_flat(doc: document.Document) -> list[Chunk_docx]:
    """
    Creates chunks based on headings; no nesting logic used. New chunks are created when:
    - a page break happens in the middle of a paragraph
    - we reach any paragraph that is any heading

    CHUNK:
    H1
    Normal Paragraph

    NEXT_CHUNK:
    H2
    Normal Paragraph
    Normal Paragraph

    NEXT_CHUNK:
    H2
    Normal Paragraph
    Normal Paragraph
    Normal Paragraph

    NEXT_CHUNK:
    H1

    NEXT_CHUNK:
    H2
    Pararaph
    Normal Paragraph

    NEXT_CHUNK:
    H3
    Normal Paragraph

    NEXT_CHUNK:
    H2
    Normal Paragraph
    Normal Paragraph
    """

    # Start building the chunks
    all_chunks: list[Chunk_docx] = []
    current_chunk: Chunk_docx = Chunk_docx()

    for para in doc.paragraphs:
        # Skip empty paragraphs
        if para.text == "":
            continue

        # Set a style_name to make Pylance happy (it gets mad if we direct-check para.style.name later)
        style_name = para.style.name if para.style and para.style.name else "Normal"

        debug_print(f"Paragraph begins: {para.text[:30]}...")

        # If the current_chunk is empty, append the current para regardless of style & continue to next para.
        if not current_chunk.paragraphs:
            current_chunk.add_paragraph(para)
            continue

        # Handle page breaks - always start a new chunk
        if para.contains_page_break:
            # Add the current chunk to chunks list (if it has content)
            if current_chunk:
                all_chunks.append(current_chunk)

            # Start new chunk with this paragraph
            current_chunk = Chunk_docx.create_with_paragraph(para)
            continue

        # If this paragraph is a heading, start a new chunk
        if is_standard_heading(style_name):
            # If we already have content in current_chunk, save it and start fresh
            if current_chunk:
                all_chunks.append(current_chunk)

            # Start new chunk with this paragraph
            current_chunk = Chunk_docx.create_with_paragraph(para)

        else:
            # This is a normal paragraph - add it to current chunk
            current_chunk.add_paragraph(para)

    if current_chunk:
        all_chunks.append(current_chunk)

    print(f"This document has {len(all_chunks)} flat heading chunks.")
    return all_chunks


# endregion

import inspect


# region Utils - Basic
def debug_print(msg: str | list[str]) -> None:
    """Basic debug printing function"""
    if DEBUG_MODE:
        caller = inspect.stack()[1].function
        print(f"DEBUG [{caller}]: {msg}")


def setup_console_encoding() -> None:
    """Configure UTF-8 encoding for Windows console to prevent UnicodeEncodeError when printing non-ASCII characters (like emojis)."""
    if platform.system() == "Windows":
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")


# endregion


# region Utils - I/O (./io.py)
def validate_path(user_path: str | Path) -> Path:
    """Ensure filepath exists and is a file."""
    path = Path(user_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {user_path}")
    if not path.is_file():
        raise ValueError("That's not a file.")
    return path


# endregion


# region docx load and validate
def validate_docx_path(user_path: str | Path) -> Path:
    """Validates the user-provided filepath exists and is actually a docx file."""
    path = validate_path(user_path)

    # Verify it's the right extension:
    if path.suffix.lower() == ".doc":
        raise ValueError(
            "This tool only supports .docx files right now. Please convert your .doc file to .docx format first."
        )
    if path.suffix.lower() != ".docx":
        raise ValueError(f"Expected a .docx file, but got: {path.suffix}")
    return path


# eventual destination: ./src/docx2pptx/io.py (?)
def load_and_validate_docx(input_filepath: Path) -> document.Document:
    """Use python-docx to read in the docx file contents and store to a runtime variable."""

    # Try to load the docx
    try:
        doc = docx.Document(input_filepath)  # type: ignore
    except Exception as e:
        raise ValueError(f"Document appears to be corrupted: {e}")

    # Validate it contains content
    if not doc.paragraphs:
        raise ValueError("Document contains no paragraphs.")

    first_para_w_text = find_first_docx_paragraph_with_text(doc.paragraphs)
    if first_para_w_text is None:
        raise ValueError("Document contains no text content.")

    # Report content information to the user
    paragraph_count = len(doc.paragraphs)
    debug_print(f"This document has {paragraph_count} paragraphs in it.")

    text = first_para_w_text.text
    preview = text[:20] + ("..." if len(text) > 20 else "")
    debug_print(f"The first paragraph containing text begins with: {preview}")

    return doc


def find_first_docx_paragraph_with_text(
    paragraphs: list[Paragraph_docx],
) -> Paragraph_docx | None:
    """Find the first paragarph that contains any text content in a docx."""
    for paragraph in paragraphs:
        if paragraph.text and paragraph.text.strip():
            return paragraph
    return None


# endregion


# region pptx load and validate
def validate_pptx_path(user_path: str | Path) -> Path:
    """Validates the pptx template filepath exists and is actually a pptx file."""
    path = validate_path(user_path)
    # Verify it's the right extension:
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"Expected a .pptx file, but got: {path.suffix}")
    return path


# Observations about slides and slide_ids:
# 1)    Sequential iteration order matters: slide_ids will iterate here in the order that they appear in the visual slide deck.
# 2)    slide_id does NOT imply its place in the slide order: slide_ids appear to be n+1 generated,
#       but a person can easily move slides around in the sequence of slides as desired. The only meaning that you can gain
#       by sorting by slide_id is (maybe) the order in which the slide items were added to the deck, not the order that the
#       user wants them to be viewed/read.


def load_and_validate_pptx(pptx_path: Path | str) -> presentation.Presentation:
    """
    Read in pptx file contents, validate minimum content is present, and store to a runtime object. (pptx2docx-text pipeline)
    """

    # Try to load the pptx
    try:
        prs = pptx.Presentation(str(pptx_path))
    except Exception as e:
        raise ValueError(f"Presentation appears to be corrupted: {e}")

    # Validate the pptx contains slides, and at least one contains content.
    if not prs.slides:
        raise ValueError("Presentation contains no slides.")

    first_slide = find_first_slide_with_text(list(prs.slides))
    if first_slide is None:
        raise ValueError(
            f"No slides in {pptx_path} contain text content, so there's nothing for the pipeline to do."
        )

    # Report content information to the user
    slide_count = len(prs.slides)
    debug_print(f"The pptx file {pptx_path} has {slide_count} slide(s) in it.")

    first_slide_paragraphs = get_slide_paragraphs(first_slide)
    debug_print(
        f"The first slide detected with text content is slide_id: {first_slide.slide_id}. The text content is: \n"
    )

    for p in first_slide_paragraphs:
        if p.text.strip():
            text = p.text.strip()
            preview = text[:20] + ("..." if len(text) > 20 else "")
            debug_print(f"The first paragraph containing text begins with: {preview}")
            break
    # An else on a for-loop runs if we never hit break. This is here because I'm maybe-overly defensive in programming style.
    else:
        debug_print("(Could not extract preview text)")

    # Return the runtime object
    return prs


def find_first_slide_with_text(slides: list[Slide]) -> Slide | None:
    """Find the first slide that contains any paragraphs with text content."""
    for slide in slides:
        if get_slide_paragraphs(slide):
            return slide
    return None


def get_slide_paragraphs(slide: Union[Slide, NotesSlide]) -> list[Paragraph_pptx]:
    """Extract all paragraphs from all text placeholders in a slide."""
    paragraphs: list[Paragraph_pptx] = []

    for placeholder in slide.placeholders:
        if (
            isinstance(placeholder, SlidePlaceholder)
            and hasattr(placeholder, "text_frame")
            and placeholder.text_frame
        ):
            textf: TextFrame = placeholder.text_frame
            for para in textf.paragraphs:
                if para.runs or para.text:
                    paragraphs.append(para)

    return paragraphs


def create_empty_slide_deck() -> presentation.Presentation:
    """Load the PowerPoint template, create a new presentation object, and validate it contains the custom layout. (docx2pptx-text pipeline)"""

    # Try to load the pptx
    try:
        template_path = validate_pptx_path(Path(TEMPLATE_PPTX))
        prs = pptx.Presentation(str(template_path))
    except Exception as e:
        raise ValueError(f"Could not load template file (may be corrupted): {e}")

    # Validate it has the required slide layout for the pipeline
    layout_names = [layout.name for layout in prs.slide_layouts]
    if SLD_LAYOUT_CUSTOM_NAME not in layout_names:
        raise ValueError(
            f"Template is missing the required layout: '{SLD_LAYOUT_CUSTOM_NAME}'. "
            f"Available layouts: {', '.join(layout_names)}"
        )

    return prs


# TODO, multi-file split: Another TypeVar to move to the top of whatever file these funcs live in later
OUTPUT_TYPE = TypeVar("OUTPUT_TYPE", document.Document, presentation.Presentation)


# TODO, leafy: I'd really like this to validate we're not about to save 100MB+ files. But that's not easy to
# estimate from the runtime object. 
# For now we'll check for absolutely insane slide or paragraph counts, and just report it to the
# debug/logger.
# TODO, polish: Around here is where we ought to add an option to split the output into multiple files, 
# by X-number of slides or pages. There probably needs to be a default for each output type and a way for the
# user to specify an override for the default.
def _validate_content_size(save_object: OUTPUT_TYPE) -> None:
    """Report if the output content we're about to save is excessively large."""
    if isinstance(save_object, document.Document):
        max_p_count = 10000
        if len(save_object.paragraphs) > max_p_count:
            debug_print(
                f"This is about to save a docx file with over {max_p_count} paragraphs ... that seems a bit long!"
            )
    elif isinstance(save_object, presentation.Presentation):
        max_s_count = 1000
        if len(list(save_object.slides)) > max_s_count:
            debug_print(
                f"This is about to save a pptx file with over {max_s_count} slides ... that seems a bit long!"
            )


def _determine_output_path(save_object: OUTPUT_TYPE) -> tuple[Path, str]:
    """Construct output folder and filename in memory based on output type."""
    if isinstance(save_object, document.Document):
        save_folder = OUTPUT_DOCX_FOLDER
        save_filename = OUTPUT_DOCX_FILENAME
        return save_folder, save_filename
    elif isinstance(save_object, presentation.Presentation):
        save_folder = OUTPUT_PPTX_FOLDER
        save_filename = OUTPUT_PPTX_FILENAME
        return save_folder, save_filename
    else:
        raise RuntimeError(f"Unexpected output object type: {save_object}")


def save_output(save_object: OUTPUT_TYPE) -> None:
    """Save the generated output object to disk as a file. Genericized to output either docx or pptx depending on which pipeline is running."""

    # Build the output path components based on filetype.
    save_folder, save_filename = _determine_output_path(save_object)

    # Report if the content we're about to save is excessively huge
    _validate_content_size(save_object)

    # Create the output folder if we need to
    save_folder.mkdir(parents=True, exist_ok=True)

    # Add a timestamp to the filename
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    name, ext = save_filename.rsplit(
        ".", 1
    )  # The 1 is telling rsplit() to split from the right side and do a maximum of 1 split.
    timestamped_filename = f"{name}_{timestamp}.{ext}"
    output_filepath = save_folder / timestamped_filename

    # Attempt to save
    try:
        save_object.save(str(output_filepath))
        print(f"Successfully saved to {output_filepath}")
    except PermissionError:
        raise PermissionError("Save failed: File may be open in another program")
    except OSError as e:
        raise OSError(f"Save failed (disk space or IO issue): {e}")
    except Exception as e:
        raise RuntimeError(f"Save failed with unexpected error: {e}")


# endregion


# region sanitize xml
def sanitize_xml_text(text: str) -> str:
    """Remove characters that aren't valid in XML."""
    if not text:
        return ""

    # Remove NULL bytes and control characters (except tab, newline, carriage return)
    sanitized = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]", "", text)

    # Ensure it's a proper string
    return str(sanitized)


# endregion

# region call main
if __name__ == "__main__":
    main()
# endregion
