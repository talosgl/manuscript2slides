# templates.py
"""Load docx and pptx templates from disk, validate shape, and create in-memory python objects from them."""
# pyright: reportArgumentType=false, reportIndexIssue=false,  reportAttributeAccessIssue=false
# mypy: disable-error-code="import-untyped"

import logging
from pathlib import Path

import docx
import pptx
from docx import document
from docx.text.paragraph import Paragraph as Paragraph_docx
from pptx import presentation

from manuscript2slides import io
from manuscript2slides.internals import constants
from manuscript2slides.internals.define_config import UserConfig

log = logging.getLogger("manuscript2slides")


def create_empty_slide_deck(cfg: UserConfig) -> presentation.Presentation:
    """Load the PowerPoint template, create a new presentation object, and validate it contains the custom layout. (manuscript2slides pipeline)"""

    # Try to load the pptx
    try:
        template_path = cfg.get_template_pptx_path()
        validated_template = io.validate_pptx_path(Path(template_path))
        prs: presentation.Presentation = pptx.Presentation(str(validated_template))  # pyright: ignore[reportPrivateImportUsage]
    except Exception as e:
        log.error(f"Could not load template file at path {e}")
        raise ValueError(f"Could not load template file (may be corrupted): {e}")

    # Validate it has the required slide layout for the pipeline
    layout_names = [layout.name for layout in prs.slide_layouts]
    if constants.SLD_LAYOUT_CUSTOM_NAME not in layout_names:
        log.error(
            f"Could not find required slide layout {constants.SLD_LAYOUT_CUSTOM_NAME} in template."
        )
        raise ValueError(
            f"Template is missing the required layout: '{constants.SLD_LAYOUT_CUSTOM_NAME}'. "
            f"Available layouts: {', '.join(layout_names)}"
            f"If error persists, try renaming the Documents/manuscript2slides/templates/ folder to templates_old/ or deleting it."
        )

    num_slides = len(prs.slides)

    # If the user passed in a deck that contained slides, delete them from our in-memory copy.
    if num_slides > 0:
        delete_all_prs_slides(prs)

    return prs


def delete_all_prs_slides(prs: presentation.Presentation) -> None:
    """Safely remove all slides from a Presentation object."""
    num_slides = len(prs.slides)

    # Iterate through all slides in reverse order and delete them
    # Iterating backward prevents index issues after deletion
    # This backward looping pattern is crucial when we delete items from a collection while iterating over it.
    for i in range(
        num_slides - 1,  # start: index of the last slide
        -1,  # stop: stop just before index -1 (which includes index 0)
        -1,  # step: count down by 1 (aka count by -1)
    ):
        delete_pptx_slide(prs, i)


def delete_pptx_slide(prs: presentation.Presentation, i: int) -> None:
    """Safely delete a slide from a Presentation object by index."""
    # Access internal document structure to remove the slide<->prs relationship
    rId = prs.slides._sldIdLst[i].rId
    prs.part.drop_rel(rId)

    # Finally, delete the slide from the slides list
    del prs.slides._sldIdLst[i]


def create_empty_document(cfg: UserConfig) -> document.Document:
    """
    Load Word template and create document object.

    Validates the template is a valid docx file.

    Raises:
        ValueError: If template is corrupted or invalid.
    """
    from manuscript2slides.io import validate_docx_path  # Avoid circular import

    try:
        template_path = cfg.get_template_docx_path()
        validated_template = validate_docx_path(Path(template_path))
        doc = docx.Document(str(validated_template))
    except Exception as e:
        raise ValueError(f"Could not load docx template (may be corrupted): {e}")

    required_styles = {"Normal"}

    # Get a list of all style names in the document
    template_styles = [style.name for style in doc.styles]
    log.debug(f"Available styles in document: {template_styles}")

    missing_styles = []
    for style_name in required_styles:
        if style_name not in template_styles:
            missing_styles.append(style_name)

    if missing_styles:
        log.error(
            "Could not find required default Word style.name 'Normal' in provided template."
        )
        raise ValueError(
            f"Template is missing required styles: '{required_styles}'. "
            f"Available layouts: {template_styles}"
            f"If error persists, try renaming the Documents/manuscript2slides/templates/ folder to templates_old/ or deleting it, then re-running the program."
        )

    # Remove any existing text
    num_para = len(doc.paragraphs)

    if num_para > 0:
        # Iterate backward through the paragraphs list for safe deletion
        delete_all_docx_paragraphs(doc)

    return doc


def delete_all_docx_paragraphs(doc: document.Document) -> None:
    """Safely delete all paragraphs from a Document object."""
    for i in range(len(doc.paragraphs) - 1, -1, -1):
        paragraph = doc.paragraphs[i]
        delete_docx_paragraph(paragraph)


def delete_docx_paragraph(paragraph: Paragraph_docx) -> None:
    """Safely deletes a paragraph from a Document object."""
    p = paragraph._element
    p.getparent().remove(p)

    # Optional: Clear internal references to prevent errors if the object is used later
    paragraph._p = paragraph._element = None  # type: ignore[assignment]
