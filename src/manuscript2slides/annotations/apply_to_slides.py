"""Add annotations we pulled from the docx to PowerPoint slide notes."""

# mypy: disable-error-code="import-untyped"
# region imports
import json
import logging
from datetime import datetime

from pptx.text.text import TextFrame

from manuscript2slides.annotations.extract import NOTE_TYPE
from manuscript2slides.internals.constants import (
    METADATA_MARKER_FOOTER,
    METADATA_MARKER_HEADER,
    NOTES_MARKER_FOOTER,
    NOTES_MARKER_HEADER,
)
from manuscript2slides.internals.define_config import UserConfig
from manuscript2slides.models import (
    Chunk_docx,
    Comment_docx_custom,
    Endnote_docx,
    Footnote_docx,
)
from manuscript2slides.processing.run_processing import (
    process_docx_paragraph_inner_contents,
)

# endregion
log = logging.getLogger("manuscript2slides")


# region annotate_slide
def annotate_slide(
    chunk: Chunk_docx, notes_text_frame: TextFrame, cfg: UserConfig
) -> None:
    """
    Pull a chunk's preserved annotations and copy them into the slide's speaker notes text frame.

    NOTE: We DO NOT PRESERVE any anchoring to the slide's body for annotations. That means we don't replicate
    comments' selected text ranges, nor do we preserve footnote or endnote numbering.
    """
    # Check if there's anything to annotate first
    has_content = (
        (cfg.display_comments and chunk.comments)
        or (cfg.display_footnotes and chunk.footnotes)
        or (cfg.display_endnotes and chunk.endnotes)
    )

    if not has_content:
        return  # Nothing to add, exit early

    notes_text_frame.add_paragraph()  # add a blank first line for actual annotations by the user

    header_para = notes_text_frame.add_paragraph()
    header_run = header_para.add_run()
    header_run.text = f"\n\n\n\n\n\n\n{NOTES_MARKER_HEADER}\n" + "=" * 40 + "\n"

    if cfg.display_comments and chunk.comments:
        add_comments_to_speaker_notes(chunk.comments, notes_text_frame, cfg)

    if cfg.display_footnotes and chunk.footnotes:
        add_notes_to_speaker_notes(chunk.footnotes, notes_text_frame, Footnote_docx)

    if cfg.display_endnotes and chunk.endnotes:
        add_notes_to_speaker_notes(chunk.endnotes, notes_text_frame, Endnote_docx)

    footer_para = notes_text_frame.add_paragraph()
    footer_run = footer_para.add_run()
    footer_run.text = "=" * 40 + f"\n{NOTES_MARKER_FOOTER}"


# endregion


# region add_comments_to_speaker_notes
def add_comments_to_speaker_notes(
    comments_list: list[Comment_docx_custom],
    notes_text_frame: TextFrame,
    cfg: UserConfig,
) -> None:
    """Copy logic for appending the comments portion of the speaker notes."""

    if comments_list:
        if cfg.comments_sort_by_date:
            # Sort comments by date (newest first, or change reverse=False for oldest first)
            sorted_comments = sorted(
                comments_list,
                key=lambda c: getattr(c.comment_obj, "timestamp", None) or datetime.min,
                reverse=False,
            )
        else:
            sorted_comments = comments_list

        comment_para = notes_text_frame.add_paragraph()
        comment_run = comment_para.add_run()
        comment_run.text = "COMMENTS FROM SOURCE DOCUMENT:\n" + "=" * 40

        for i, comment in enumerate(sorted_comments, 1):
            # Check if comment has paragraphs attribute
            if hasattr(comment.comment_obj, "paragraphs"):
                # Get paragraphs safely, default to empty list if not present
                this_comment_paragraphs = getattr(comment.comment_obj, "paragraphs", [])
                for para in this_comment_paragraphs:
                    if hasattr(para, "text") and para.text.rstrip():
                        notes_para = notes_text_frame.add_paragraph()
                        comment_header = notes_para.add_run()

                        if cfg.comments_keep_author_and_date:
                            author = getattr(
                                comment.comment_obj, "author", "Unknown Author"
                            )
                            timestamp = getattr(comment.comment_obj, "timestamp", None)

                            if timestamp and hasattr(timestamp, "strftime"):
                                timestamp_str = timestamp.strftime(
                                    "%A, %B %d, %Y at %I:%M %p"
                                )
                            else:
                                timestamp_str = "Unknown Date"

                            comment_header.text = (
                                f"\n {i}. {author} ({timestamp_str}):\n"
                            )

                        else:
                            comment_header.text = "\n"
                        process_docx_paragraph_inner_contents(para, notes_para, cfg)


# endregion


# region add_notes_to_speaker_notes
def add_notes_to_speaker_notes(
    notes_list: list[NOTE_TYPE],
    notes_text_frame: TextFrame,
    note_class: type[NOTE_TYPE],
) -> None:
    """Generic function for adding footnotes or endnotes to speaker notes."""

    if note_class is Footnote_docx:
        header_text = "FOOTNOTES FROM SOURCE DOCUMENT"
    elif note_class is Endnote_docx:
        header_text = "ENDNOTES FROM SOURCE DOCUMENT"
    else:
        header_text = "Unknown Note Type from Source Document:"

    if notes_list:
        note_para = notes_text_frame.add_paragraph()
        note_run = note_para.add_run()
        note_run.text = f"\n{header_text}:\n" + "=" * 40

        for note_obj in notes_list:
            notes_para = notes_text_frame.add_paragraph()
            note_run = notes_para.add_run()

            # Start with the main note text
            note_text = f"\n{note_obj.note_id}. {note_obj.text_body}\n"

            # Add hyperlinks if they exist
            if note_obj.hyperlinks:
                note_text += "\nHyperlinks:"
                for hyperlink in note_obj.hyperlinks:
                    note_text += f"\n{hyperlink}"

            note_run.text = note_text


# endregion


# region add_metadata_to_slide_notes
def add_metadata_to_slide_notes(
    notes_text_frame: TextFrame, chunk: Chunk_docx, slide_body_metadata: dict
) -> None:
    """
    Populate the slide notes text frame with docx metadata so that we may restore it during a round-trip pptx2docx pipeline.
    """
    comments = []

    if chunk.comments:
        comments = [
            {
                "original": {
                    "text": c.comment_obj.text,
                    "author": c.comment_obj.author,
                    "comment_id": c.comment_obj.comment_id,
                    "initials": c.comment_obj.initials,
                    "timestamp": (
                        str(c.comment_obj.timestamp)
                        if c.comment_obj.timestamp
                        else None
                    ),
                },
                "reference_text": c.reference_text,
                "id": c.note_id,
            }
            for c in chunk.comments
        ]

    footnotes = []
    if chunk.footnotes:
        footnotes = [
            {
                "id": f.note_id,
                "text_body": f.text_body,
                "hyperlinks": f.hyperlinks,
                "reference_text": f.reference_text,
                "note_type": "footnote",
            }
            for f in chunk.footnotes
        ]

    endnotes = []
    if chunk.endnotes:
        endnotes = [
            {
                "id": e.note_id,
                "text_body": e.text_body,
                "hyperlinks": e.hyperlinks,
                "reference_text": e.reference_text,
                "note_type": "endnote",
            }
            for e in chunk.endnotes
        ]

    if slide_body_metadata or comments or footnotes or endnotes:
        metadata_para = notes_text_frame.add_paragraph()
        metadata_run = metadata_para.add_run()
        metadata_run.text = f"\n\n\n\n\n\n\n{METADATA_MARKER_HEADER}\n" + "=" * 40

        notes_text_frame.add_paragraph()  # blank paragraph to ensure separation for JSON block

        json_para = notes_text_frame.add_paragraph()
        json_run = json_para.add_run()

        combined_metadata = {**slide_body_metadata}

        if comments:  # only add if the list has content
            combined_metadata["comments"] = comments

        if footnotes:
            combined_metadata["footnotes"] = footnotes

        if endnotes:
            combined_metadata["endnotes"] = endnotes

        json_run.text = json.dumps(combined_metadata, indent=2)

        footer_para = notes_text_frame.add_paragraph()
        footer_run = footer_para.add_run()
        footer_run.text = "=" * 40 + f"\n{METADATA_MARKER_FOOTER}"


# endregion
