# io.py
"""File I/O operations for docx and pptx files."""

# mypy: disable-error-code="import-untyped"
import logging
import sys
from datetime import datetime
from pathlib import Path
from typing import TypeVar

import docx
import pptx
from docx import document
from docx.text.paragraph import Paragraph as Paragraph_docx
from pptx import presentation
from pptx.slide import Slide

from manuscript2slides.internals import constants
from manuscript2slides.internals.define_config import UserConfig
from manuscript2slides.internals.run_context import get_pipeline_run_id
from manuscript2slides.processing.populate_docx import get_slide_paragraphs

log = logging.getLogger("manuscript2slides")

OUTPUT_TYPE = TypeVar("OUTPUT_TYPE", document.Document, presentation.Presentation)

# region validate path/extension helpers (Read I/O)


# region _validate_path
def _validate_path(user_path: str | Path) -> Path:
    """Ensure filepath exists and is a file."""
    path = Path(user_path)
    pipeline_id = get_pipeline_run_id()
    if not path.exists():
        # Check if this looks like a Windows path on a non-Windows system
        path_str = str(user_path)
        if sys.platform != "win32" and "\\" in path_str and ":" in path_str:
            log.error(
                f"File not found: {user_path} [pipeline:{pipeline_id}]\n"
                f"This appears to be a Windows-style path (contains backslashes and drive letter), "
                f"but you're running on {sys.platform}. "
                f"Please use forward slashes (/) in your config file."
            )
        else:
            log.error(f"File not found: {user_path} [pipeline:{pipeline_id}]")
        raise FileNotFoundError(f"File not found: {user_path}")
    if not path.is_file():
        log.error(
            f"Path is not a file (might be a directory): {user_path} [pipeline:{pipeline_id}]"
        )
        raise ValueError(f"Path is not a file: {user_path}")
    return path


# endregion


# region validate_docx_path
def validate_docx_path(user_path: str | Path) -> Path:
    """Validates the user-provided filepath exists and is actually a docx file."""
    path = _validate_path(user_path)

    # Get pipeline ID for logs.
    pipeline_id = get_pipeline_run_id()

    # Verify it's the right extension:
    if path.suffix.lower() == ".doc":
        log.error(f"Unsupported .doc file: {path} [pipeline:{pipeline_id}]")
        raise ValueError(
            "This tool only supports .docx files. Please convert your .doc file to .docx format first."
        )
    if path.suffix.lower() != ".docx":
        log.error(
            f"Wrong file extension: expected .docx, got {path.suffix} [pipeline:{pipeline_id}]"
        )
        raise ValueError(f"Expected a .docx file, but got: {path.suffix}")
    return path


# endregion


# region validate_pptx_path
def validate_pptx_path(user_path: str | Path) -> Path:
    """Validates the pptx template filepath exists and is actually a pptx file."""
    path = _validate_path(user_path)

    # Get pipeline ID for logs.
    pipeline_id = get_pipeline_run_id()

    # Verify it's the right extension:
    if path.suffix.lower() == ".ppt":
        log.error(f"Unsupported .ppt file: {path} [pipeline:{pipeline_id}]")
        raise ValueError(
            "This tool only supports .pptx files. Please convert your .ppt file to .pptx format first."
        )
    if path.suffix.lower() != ".pptx":
        log.error(
            f"Wrong file extension: expected .pptx, got {path.suffix} [pipeline:{pipeline_id}]"
        )
        raise ValueError(f"Expected a .pptx file, but got: {path.suffix}")
    return path


# endregion


# endregion


# region validate file content helpers (read I/O)


# region load_and_validate_docx
def load_and_validate_docx(input_filepath: Path) -> document.Document:
    """Use python-docx to read in the (input) docx file contents and store to a runtime variable."""
    pipeline_id = get_pipeline_run_id()

    # Try to load the docx
    try:
        doc = docx.Document(input_filepath)  # type: ignore
    except FileNotFoundError:
        # Should never happen (validated earlier), but be defensive
        log.error(f"File not found: {input_filepath} [pipeline:{pipeline_id}]")
        raise FileNotFoundError(f"File not found: {input_filepath}")
    except PermissionError as e:
        log.error(
            f"Permission denied reading {input_filepath} [pipeline:{pipeline_id}]: {e}"
        )
        raise PermissionError(
            f"Permission denied reading file. Is it open in another program?\n{input_filepath}"
        )
    except Exception as e:
        # Corrupted file, wrong format, etc.
        log.error(
            f"Could not load document {input_filepath} [pipeline:{pipeline_id}]: {e}"
        )
        raise ValueError(f"Document may be corrupted or in wrong format:\n{e}")

    # Validate it contains content
    if not doc.paragraphs:
        log.error(
            f"Document {str(input_filepath)} contains no paragraphs [pipeline:{pipeline_id}]"
        )
        raise ValueError("Document contains no paragraphs.")

    first_para_w_text = _find_first_docx_paragraph_with_text(doc.paragraphs)
    if first_para_w_text is None:
        log.error(
            f"Document {str(input_filepath)} contains no text content [pipeline:{pipeline_id}]"
        )
        raise ValueError(
            "Document contains no text content, so there's nothing for the pipeline to do."
        )

    # Report content information to the user
    paragraph_count = len(doc.paragraphs)
    log.info(
        f"This document has {paragraph_count} paragraphs in it. [pipeline:{pipeline_id}]"
    )

    text = first_para_w_text.text
    preview = text[:20] + ("..." if len(text) > 20 else "")
    log.info(
        f"The first paragraph containing text begins with: {preview}. [pipeline:{pipeline_id}]"
    )

    return doc


# endregion


# region load_and_validate_pptx
def load_and_validate_pptx(pptx_path: Path | str) -> presentation.Presentation:
    """
    Read in input pptx file contents, validate minimum content is present, and store to a runtime object. (pptx2docx-text pipeline)
    """
    pipeline_id = get_pipeline_run_id()

    # Try to load the pptx
    try:
        prs = pptx.Presentation(pptx_path)
    except FileNotFoundError:
        log.error(f"File not found: {pptx_path} [pipeline:{pipeline_id}]")
        raise FileNotFoundError(f"File not found: {pptx_path}")
    except PermissionError as e:
        log.error(
            f"Permission denied reading {pptx_path} [pipeline:{pipeline_id}]: {e}"
        )
        raise PermissionError(
            f"Permission denied reading file. Is it open in another program?\n{pptx_path}"
        )
    except Exception as e:
        log.error(
            f"Could not load PowerPoint file {pptx_path} [pipeline:{pipeline_id}]: {e}"
        )
        raise ValueError(f"Presentation may be corrupted or in wrong format:\n{e}")

    # Validate the pptx contains slides, and at least one contains content.
    if not prs.slides:
        log.error(
            f"Document {str(pptx_path)} contains no slides. [pipeline:{pipeline_id}]"
        )
        raise ValueError("Presentation contains no slides.")

    first_slide = _find_first_slide_with_text(list(prs.slides))
    if first_slide is None:
        log.error(
            f"No slides in {str(pptx_path)} contain text content. [pipeline:{pipeline_id}]"
        )
        raise ValueError(
            f"No slides in {str(pptx_path)} contain text content, so there's nothing for the pipeline to do."
        )

    # Report content information to the user
    slide_count = len(prs.slides)
    log.info(
        f"The pptx file {pptx_path} has {slide_count} slide(s) in it. [pipeline:{pipeline_id}]"
    )

    first_slide_paragraphs = get_slide_paragraphs(first_slide)

    log.info(
        f"The first slide detected with text content is slide_id: {first_slide.slide_id} (inside presentation.xml). [pipeline:{pipeline_id}]"
    )

    for p in first_slide_paragraphs:
        if p.text.strip():
            text = p.text.strip()
            preview = text[:20] + ("..." if len(text) > 20 else "")
            log.info(f"The first paragraph containing text begins with: {preview}")
            break
    # An else on a for-loop runs if we never hit break. This is here because I'm maybe-overly defensive in programming style.
    else:
        log.warning(f"(Could not extract preview text) [pipeline:{pipeline_id}]")

    # Return the runtime object
    return prs


# endregion


# region _find_first_docx_paragraph_with_text
def _find_first_docx_paragraph_with_text(
    paragraphs: list[Paragraph_docx],
) -> Paragraph_docx | None:
    """Find the first paragraph that contains any text content in a docx."""
    for paragraph in paragraphs:
        if paragraph.text and paragraph.text.strip():
            return paragraph
    return None


# endregion


# region _find_first_slide_with_text
def _find_first_slide_with_text(slides: list[Slide]) -> Slide | None:
    """Find the first slide that contains any paragraphs with text content."""
    for slide in slides:
        if get_slide_paragraphs(slide):
            return slide
    return None


# endregion
# endregion


# region Save file helpers (write I/O)


# region save_output
def save_output(save_object: OUTPUT_TYPE, cfg: UserConfig) -> Path:
    """
    Save the generated output object to disk as a file. Genericized to output either docx or pptx depending on which pipeline is running.
    Returns the path to which we save.
    """
    pipeline_id = get_pipeline_run_id()

    # Get the output folder from the config object
    save_folder = cfg.get_output_folder()

    # Apply a timestamp to the base filename
    timestamped_filename = _build_timestamped_output_filename(save_object)

    # Report if the content we're about to save is excessively huge
    _validate_content_size(save_object)

    # Create the output folder if we need to
    save_folder.mkdir(parents=True, exist_ok=True)

    output_filepath = save_folder / timestamped_filename

    # Attempt to save
    try:
        save_object.save(str(output_filepath))
        log.info(f"Successfully saved to {output_filepath}. [pipeline:{pipeline_id}]")
    except PermissionError as e:
        log.error(f"Save failed - permission denied [pipeline:{pipeline_id}]: {e}")
        raise PermissionError(
            f"Could not save file. It may be open in another program or you lack write permissions.\n"
            f"Location: {output_filepath}\n"
            f"Error: {e}"
        )
    except OSError as e:
        log.error(f"Save failed - OS error [pipeline:{pipeline_id}]: {e}")
        raise OSError(
            f"Could not save file. Check disk space and that the path is accessible.\n"
            f"Location: {output_filepath}\n"
            f"Error: {e}"
        )
    except Exception as e:
        log.error(f"Save failed - unexpected error [pipeline:{pipeline_id}]: {e}")
        raise RuntimeError(
            f"Save failed with unexpected error.\n"
            f"Location: {output_filepath}\n"
            f"Error: {e}"
        )

    return output_filepath


# endregion


# region _build_timestamped_output_filename
def _build_timestamped_output_filename(save_object: OUTPUT_TYPE) -> str:
    """Apply a per-run timestamp to the output's base filename."""
    # Get pipeline ID for logs.
    pipeline_id = get_pipeline_run_id()

    # Get the base filename string
    if isinstance(save_object, document.Document):
        save_filename = constants.OUTPUT_DOCX_FILENAME
    elif isinstance(save_object, presentation.Presentation):
        save_filename = constants.OUTPUT_PPTX_FILENAME
    else:
        log.error(
            f"Unexpected output object type passed to _build_timestamped_output_filename(): {save_object}. Only document.Document or presentation.Presentation objects are supported. [pipeline:{pipeline_id}]"
        )
        raise RuntimeError(f"Unexpected output object type: {save_object}")

    # Add a timestamp to the filename
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    name, ext = save_filename.rsplit(
        ".", 1
    )  # The 1 is telling rsplit() to split from the right side and do a maximum of 1 split.
    timestamped_filename = f"{name}_{timestamp}.{ext}"

    return timestamped_filename


# endregion


# region _validate_content_size
def _validate_content_size(save_object: OUTPUT_TYPE) -> None:
    """Report if the output content we're about to save is excessively large."""
    if isinstance(save_object, document.Document):
        max_p_count = 10000
        if len(save_object.paragraphs) > max_p_count:
            log.warning(
                f"This is about to save a docx file with over {max_p_count} paragraphs ... that seems a bit long!"
            )
    elif isinstance(save_object, presentation.Presentation):
        max_s_count = 1000
        if len(list(save_object.slides)) > max_s_count:  # type: ignore[reportAttributeAccessIssue]
            log.warning(
                f"This is about to save a pptx file with over {max_s_count} slides ... that seems a bit long!"
            )


# endregion

# endregion
