# populate_docx.py
"""Process slides from a presentation and copy their content into a Word document.

This module handles the reverse pipeline (pptx -> docx), extracting text content
from slide paragraphs, restoring formatting and annotations from speaker notes metadata,
and creating comments for user notes and unmatched annotations.
"""


import logging
from typing import Union

from docx import document
from docx.comments import Comment as Comment_docx
from docx.text.run import Run as Run_docx
from pptx import presentation
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.slide import NotesSlide, Slide
from pptx.text.text import TextFrame
from pptx.text.text import _Paragraph as Paragraph_pptx

from manuscript2slides import utils
from manuscript2slides.annotations.restore_from_slides import split_speaker_notes
from manuscript2slides.internals.config.define_config import UserConfig
from manuscript2slides.models import SlideNotes
from manuscript2slides.run_processing import process_pptx_run

log = logging.getLogger("manuscript2slides")


# region copy slides to docx body orchestrator
def copy_slides_to_docx_body(
    prs: presentation.Presentation, new_doc: document.Document, cfg: UserConfig
) -> None:
    """
    Sequentially process each slide in the deck by copying the paragraphs from the slide body into the docx's body. Analyze
    the speaker notes of the slide, seeking stored JSON metadata there from previous docx2pptx runs, and use that to apply
    formatting and/or annotations.
    """

    # Make a list of all slides
    slide_list = list(prs.slides)

    # For each slide...
    for slide in slide_list:

        # If there's slide notes, process them into a SlideNotes object; otherwise, make an empty one.
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            slide_notes = split_speaker_notes(slide.notes_slide.notes_text_frame.text)
        else:
            slide_notes = SlideNotes()

        process_slide_paragraphs(slide, slide_notes, new_doc, cfg)


# endregion


# region process_slide_paragraphs
def process_slide_paragraphs(
    slide: Slide, slide_notes: SlideNotes, new_doc: document.Document, cfg: UserConfig
) -> None:
    """
    Process a slide's body content, using any metadata stored in the speaker notes to restore formatting and annotation
    anchors from an earlier docx2pptx pipeline run. If we find metadata but aren't able to attach it to text content from
    the body paragraphs/runs, attach that as a new comment to the very last copied paragraph/run from that slide.

    Additionally, store non-metadata speaker notes content as a new comment, too.
    """

    slide_paragraphs: list[Paragraph_pptx] = get_slide_paragraphs(slide)

    unmatched_annotations = []
    matched_comment_ids = set()

    last_run = None

    # For every pptx paragraph.....
    for pptx_para in slide_paragraphs:

        # Make a new docx para
        new_para = new_doc.add_paragraph()

        # If the text of this paragraph exactly matches a previous heading's text, apply that heading style
        if slide_notes and slide_notes.has_metadata and slide_notes.headings:
            for heading in slide_notes.headings:
                if heading["text"].strip() == pptx_para.text.strip():
                    new_para.style = heading["name"]
                    break  # we should only ever apply one style to a paragraph

        for run in pptx_para.runs:
            last_run = process_pptx_run(
                run, new_para, new_doc, slide_notes, matched_comment_ids, cfg
            )

    # Put the slide's user notes into a new comment attached to the last run
    if slide_notes and slide_notes.has_user_notes is True and last_run is not None:
        user_notes_comment = copy_user_notes_to_new_comment(
            slide_notes, last_run, new_doc
        )
        if user_notes_comment:
            log.debug(
                f"Added a new comment with this slide's user notes: {user_notes_comment}"
            )

    # Find all the unmatched annotations for this slide by getting the complement set(s)
    # (Only comments are supported, for now, but if we ever add footnote/endnote support,
    # we'll need 3 sets2lists here.)
    unmatched_comments = [
        c for c in slide_notes.comments if c["id"] not in matched_comment_ids
    ]

    unmatched_annotations.extend(unmatched_comments)

    # If we have any unmatched annotations from the slide_notes.metadata, attach them as a new comment to the last run
    if unmatched_annotations and last_run is not None:
        unmatched_comment = copy_unmatched_comments_to_new_comment(
            last_run, unmatched_annotations, new_doc
        )
        if unmatched_comment:
            log.debug(
                f"Added comment for {len(unmatched_annotations)} unmatched annotations"
            )


# endregion


# region copy speaker notes items into new docx comments
def copy_user_notes_to_new_comment(
    slide_notes: SlideNotes, last_run: Run_docx, new_doc: document.Document
) -> Comment_docx | None:
    """Append this slide's user notes a comment to the last-copied-in run to the docx."""

    # Verify there's actually text present to copy in; return None if not
    raw_comment_text = slide_notes.user_notes
    if not raw_comment_text.strip():
        return None

    # Prepare the header + body text for the comment
    comment_header = "Copied from the PPTX Speaker Notes: \n\n"
    comment_text = comment_header + _sanitize_xml_text(raw_comment_text)

    # Add the comment to the doc & run
    new_comment = new_doc.add_comment(last_run, comment_text)

    # Return the new comment for testing/logging purposes
    return new_comment


def copy_unmatched_comments_to_new_comment(
    last_run: Run_docx, unmatched_annotations: list, new_doc: document.Document
) -> Comment_docx | None:
    """Compile any unmatched annotations into a combined string and attach it as a comment to the final run of its parent slide."""
    unmatched_parts = []

    for annotation in unmatched_annotations:
        if "original" in annotation:  # Detect comments
            # Verify there's actually text present to copy in
            raw_original_text = annotation["original"]["text"]
            if not raw_original_text.strip():
                continue
            # Add the comment to unmatched list
            unmatched_parts.append(f"Comment: {annotation['original']['text']}")

        # elif "text_body" in annotation:  # footnotes/endnotes, cut feature for now
        #     unmatched_parts.append(f"Note: {annotation['text_body']}")

    # combine all unmatched annotations into one string
    combined = "\n\n".join(unmatched_parts)

    # Verify there's actually text present to copy into a comment; return None if not
    if not combined.strip():
        return None

    # Prepare the header + body text for the comment
    raw_comment_text = combined
    comment_header = "We found metadata for these annotations (comments, footnotes, or endnotes), but weren't able to match them to specific text in this paragraph: \n\n"
    comment_text = comment_header + _sanitize_xml_text(raw_comment_text)

    # Add the comment to the doc & run
    new_comment = new_doc.add_comment(last_run, comment_text)

    return new_comment


import re


def _sanitize_xml_text(text: str) -> str:
    """Remove characters that aren't valid in XML."""
    if not text:
        return ""

    # Remove NULL bytes and control characters (except tab, newline, carriage return)
    sanitized = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]", "", text)

    # Ensure it's a proper string
    return str(sanitized)


# region get_slide_paragraphs
def get_slide_paragraphs(slide: Union[Slide, NotesSlide]) -> list[Paragraph_pptx]:
    """Extract all paragraphs from all text placeholders in a slide."""
    paragraphs: list[Paragraph_pptx] = []

    for placeholder in slide.placeholders:
        if (
            isinstance(placeholder, SlidePlaceholder)
            and hasattr(placeholder, "text_frame")
            and placeholder.text_frame
        ):
            textf: TextFrame = placeholder.text_frame
            for para in textf.paragraphs:
                if para.runs or para.text:
                    paragraphs.append(para)

    return paragraphs


#
