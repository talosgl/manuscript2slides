# create_slides.py
"""Take chunks we built from the input docx file and turn them into slide body content."""
# mypy: disable-error-code="import-untyped"

import logging

from pptx import presentation
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.slide import Slide, SlideLayout
from pptx.text.text import TextFrame

from manuscript2slides.annotations.apply_to_slides import (
    add_metadata_to_slide_notes,
    annotate_slide,
)
from manuscript2slides.internals import constants
from manuscript2slides.internals.define_config import UserConfig
from manuscript2slides.internals.run_context import get_pipeline_run_id
from manuscript2slides.models import Chunk_docx
from manuscript2slides.processing.chunking import is_standard_heading
from manuscript2slides.processing.run_processing import (
    process_docx_paragraph_inner_contents,
)

log = logging.getLogger("manuscript2slides")


# region slides_from_chunks
def slides_from_chunks(
    prs: presentation.Presentation, chunks: list[Chunk_docx], cfg: UserConfig
) -> None:
    """Generate slide objects, one for each chunk created by earlier pipeline steps."""
    pipeline_id = get_pipeline_run_id()
    log.info(f"Creating new slides from chunks. [pipeline:{pipeline_id}]")

    # Specify which slide layout to use
    slide_layout = prs.slide_layouts.get_by_name(constants.SLD_LAYOUT_CUSTOM_NAME)

    if slide_layout is None:
        log.error(
            f"No slide layout found to match provided custom name, {constants.SLD_LAYOUT_CUSTOM_NAME}. [pipeline:{pipeline_id}]"
        )
        raise KeyError(
            f"No slide layout found to match provided custom name, {constants.SLD_LAYOUT_CUSTOM_NAME}"
        )

    for chunk in chunks:
        # Skip chunks whose page range is outside the user-specified start/end range
        if (cfg.range_start and chunk.original_sequence_number < cfg.range_start) or (
            cfg.range_end and chunk.original_sequence_number > cfg.range_end
        ):
            log.info(
                f"Skipping chunk from page {chunk.original_sequence_number} as user-specified."
            )
            continue

        # Create a new slide for this chunk.
        new_slide, text_frame = create_blank_slide_for_chunk(prs, slide_layout)

        # Store custom metadata for this chunk that we'll want to tuck into the speaker notes as JSON
        # (for the purposes of restoring during reverse pipeline runs).
        slide_metadata = {}
        headings = []
        experimental_formatting = []

        # For each paragraph in this chunk, handle adding it
        for i, paragraph in enumerate(chunk.paragraphs):
            # Creating a new slide and a text frame leaves an empty paragraph in place, even when clearing it.
            # So if we're at the start of our list, use that existing empty paragraph.
            if (
                i == 0
                and len(text_frame.paragraphs) > 0
                and not text_frame.paragraphs[0].text
            ):
                # Use the existing first/0th paragraph
                pptx_paragraph = text_frame.paragraphs[0]
            else:
                pptx_paragraph = text_frame.add_paragraph()

            # Process the docx's paragraph contents, including both runs & hyperlinks
            para_experimental_formatting = process_docx_paragraph_inner_contents(
                paragraph, pptx_paragraph, cfg
            )

            if para_experimental_formatting:
                experimental_formatting.extend(para_experimental_formatting)

            if (
                paragraph.style
                and paragraph.style.name
                and is_standard_heading(paragraph.style.name)
            ):
                headings.append(
                    {
                        "text": paragraph.text.strip(),
                        "style_id": paragraph.style.style_id,
                        "name": paragraph.style.name,
                    }
                )

        if headings:
            slide_metadata["headings"] = headings
        if experimental_formatting:
            slide_metadata["experimental_formatting"] = experimental_formatting

        notes_text_frame = new_slide.notes_slide.notes_text_frame

        if notes_text_frame is None:
            log.error(
                f"No notes text frame found in slide {new_slide.slide_id}. [pipeline:{get_pipeline_run_id()}]"
            )
            raise ValueError(
                "This slide doesn't seem to have a notes text frame. This should never happen, but it's possible for the notes_slide or notes_text_frame properties to return None if the notes placeholder has been removed from the notes master or the notes slide itself."
            )

        if cfg.display_comments or cfg.display_footnotes or cfg.display_endnotes:
            annotate_slide(chunk, notes_text_frame, cfg)

        if cfg.preserve_docx_metadata_in_speaker_notes:
            add_metadata_to_slide_notes(notes_text_frame, chunk, slide_metadata)


# endregion


# region create_blank_slide_for_chunk
def create_blank_slide_for_chunk(
    prs: presentation.Presentation, slide_layout: SlideLayout
) -> tuple[Slide, TextFrame]:
    """Initialize an empty slide so that we can populate it with a chunk."""
    new_slide = prs.slides.add_slide(  # pyright: ignore[reportAttributeAccessIssue]
        slide_layout
    )
    content = new_slide.placeholders[1]

    if not isinstance(content, SlidePlaceholder):
        log.error(
            f"Expected SlidePlaceholder, got {type(content)} [pipeline:{get_pipeline_run_id()}]"
        )
        raise TypeError(f"Expected SlidePlaceholder, got {type(content)}")

    text_frame: TextFrame = content.text_frame
    text_frame.clear()

    return new_slide, text_frame


# endregion
