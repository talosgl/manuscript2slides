"""TODO: docstring"""
import pptx
import docx
from pathlib import Path
from pptx import presentation
from typing import TypeVar
from pptx.slide import NotesSlide
from pptx.slide import Slide
from docx import document
from typing import Union
from pptx.text.text import TextFrame, _Paragraph as Paragraph_pptx, _Run as Run_pptx  # type: ignore
from pptx.shapes.placeholder import SlidePlaceholder
from src.docx2pptx import config
from docx.text.paragraph import Paragraph as Paragraph_docx
from src.docx2pptx.utils import debug_print
from datetime import datetime

# TODO, multi-file split: Another TypeVar to move to the top of whatever file these funcs live in later
OUTPUT_TYPE = TypeVar("OUTPUT_TYPE", document.Document, presentation.Presentation)


# region not sure these belong in io
# TODO: figure out where this live after we move the rest of the monolith
def get_slide_paragraphs(slide: Union[Slide, NotesSlide]) -> list[Paragraph_pptx]:
    """Extract all paragraphs from all text placeholders in a slide."""
    paragraphs: list[Paragraph_pptx] = []

    for placeholder in slide.placeholders:
        if (
            isinstance(placeholder, SlidePlaceholder)
            and hasattr(placeholder, "text_frame")
            and placeholder.text_frame
        ):
            textf: TextFrame = placeholder.text_frame
            for para in textf.paragraphs:
                if para.runs or para.text:
                    paragraphs.append(para)

    return paragraphs

# TODO: I think the problem with this one is the name, not the location
def create_empty_slide_deck() -> presentation.Presentation:
    """Load the PowerPoint template, create a new presentation object, and validate it contains the custom layout. (docx2pptx-text pipeline)"""

    # Try to load the pptx
    try:
        template_path = validate_pptx_path(Path(config.TEMPLATE_PPTX))
        prs = pptx.Presentation(str(template_path))
    except Exception as e:
        raise ValueError(f"Could not load template file (may be corrupted): {e}")

    # Validate it has the required slide layout for the pipeline
    layout_names = [layout.name for layout in prs.slide_layouts]
    if config.SLD_LAYOUT_CUSTOM_NAME not in layout_names:
        raise ValueError(
            f"Template is missing the required layout: '{config.SLD_LAYOUT_CUSTOM_NAME}'. "
            f"Available layouts: {', '.join(layout_names)}"
        )

    return prs
# endregion

# region Path Helpers
def validate_path(user_path: str | Path) -> Path:
    """Ensure filepath exists and is a file."""
    path = Path(user_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {user_path}")
    if not path.is_file():
        raise ValueError("That's not a file.")
    return path

def validate_pptx_path(user_path: str | Path) -> Path:
    """Validates the pptx template filepath exists and is actually a pptx file."""
    path = validate_path(user_path)
    # Verify it's the right extension:
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"Expected a .pptx file, but got: {path.suffix}")
    return path

def validate_docx_path(user_path: str | Path) -> Path:
    """Validates the user-provided filepath exists and is actually a docx file."""
    path = validate_path(user_path)

    # Verify it's the right extension:
    if path.suffix.lower() == ".doc":
        raise ValueError(
            "This tool only supports .docx files right now. Please convert your .doc file to .docx format first."
        )
    if path.suffix.lower() != ".docx":
        raise ValueError(f"Expected a .docx file, but got: {path.suffix}")
    return path

def _determine_output_path(save_object: OUTPUT_TYPE) -> tuple[Path, str]:
    """Construct output folder and filename in memory based on output type."""
    if isinstance(save_object, document.Document):
        save_folder = config.OUTPUT_DOCX_FOLDER
        save_filename = config.OUTPUT_DOCX_FILENAME
        return save_folder, save_filename
    elif isinstance(save_object, presentation.Presentation):
        save_folder = config.OUTPUT_PPTX_FOLDER
        save_filename = config.OUTPUT_PPTX_FILENAME
        return save_folder, save_filename
    else:
        raise RuntimeError(f"Unexpected output object type: {save_object}")

# end region

# region Disk I/O - Write
def save_output(save_object: OUTPUT_TYPE) -> None:
    """Save the generated output object to disk as a file. Genericized to output either docx or pptx depending on which pipeline is running."""

    # Build the output path components based on filetype.
    save_folder, save_filename = _determine_output_path(save_object)

    # Report if the content we're about to save is excessively huge
    _validate_content_size(save_object)

    # Create the output folder if we need to
    save_folder.mkdir(parents=True, exist_ok=True)

    # Add a timestamp to the filename
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    name, ext = save_filename.rsplit(
        ".", 1
    )  # The 1 is telling rsplit() to split from the right side and do a maximum of 1 split.
    timestamped_filename = f"{name}_{timestamp}.{ext}"
    output_filepath = save_folder / timestamped_filename

    # Attempt to save
    try:
        save_object.save(str(output_filepath))
        print(f"Successfully saved to {output_filepath}")
    except PermissionError:
        raise PermissionError("Save failed: File may be open in another program")
    except OSError as e:
        raise OSError(f"Save failed (disk space or IO issue): {e}")
    except Exception as e:
        raise RuntimeError(f"Save failed with unexpected error: {e}")



# TODO, leafy: I'd really like `_validate_content_size()` to validate we're not about to save 100MB+ files. But that's not easy to
# estimate from the runtime object.
# For now we'll check for absolutely insane slide or paragraph counts, and just report it to the
# debug/logger.

# TODO, polish: Around here is where we ought to add an option to split the output into multiple files,
# by X-number of slides or pages. There probably needs to be a default for each output type and a way for the
# user to specify an override for the default.
def _validate_content_size(save_object: OUTPUT_TYPE) -> None:
    """Report if the output content we're about to save is excessively large."""
    if isinstance(save_object, document.Document):
        max_p_count = 10000
        if len(save_object.paragraphs) > max_p_count:
            debug_print(
                f"This is about to save a docx file with over {max_p_count} paragraphs ... that seems a bit long!"
            )
    elif isinstance(save_object, presentation.Presentation):
        max_s_count = 1000
        if len(list(save_object.slides)) > max_s_count:
            debug_print(
                f"This is about to save a pptx file with over {max_s_count} slides ... that seems a bit long!"
            )

# endregion


# region Disk I/O - Read & Validate
def load_and_validate_docx(input_filepath: Path) -> document.Document:
    """Use python-docx to read in the docx file contents and store to a runtime variable."""

    # Try to load the docx
    try:
        doc = docx.Document(input_filepath)  # type: ignore
    except Exception as e:
        raise ValueError(f"Document appears to be corrupted: {e}")

    # Validate it contains content
    if not doc.paragraphs:
        raise ValueError("Document contains no paragraphs.")

    first_para_w_text = _find_first_docx_paragraph_with_text(doc.paragraphs)
    if first_para_w_text is None:
        raise ValueError("Document contains no text content.")

    # Report content information to the user
    paragraph_count = len(doc.paragraphs)
    debug_print(f"This document has {paragraph_count} paragraphs in it.")

    text = first_para_w_text.text
    preview = text[:20] + ("..." if len(text) > 20 else "")
    debug_print(f"The first paragraph containing text begins with: {preview}")

    return doc

def load_and_validate_pptx(pptx_path: Path | str) -> presentation.Presentation:
    """
    Read in pptx file contents, validate minimum content is present, and store to a runtime object. (pptx2docx-text pipeline)
    """

    # Try to load the pptx
    try:
        prs = pptx.Presentation(str(pptx_path))
    except Exception as e:
        raise ValueError(f"Presentation appears to be corrupted: {e}")

    # Validate the pptx contains slides, and at least one contains content.
    if not prs.slides:
        raise ValueError("Presentation contains no slides.")

    first_slide = _find_first_slide_with_text(list(prs.slides))
    if first_slide is None:
        raise ValueError(
            f"No slides in {pptx_path} contain text content, so there's nothing for the pipeline to do."
        )

    # Report content information to the user
    slide_count = len(prs.slides)
    debug_print(f"The pptx file {pptx_path} has {slide_count} slide(s) in it.")

    first_slide_paragraphs = get_slide_paragraphs(first_slide)
    debug_print(
        f"The first slide detected with text content is slide_id: {first_slide.slide_id}. The text content is: \n"
    )

    for p in first_slide_paragraphs:
        if p.text.strip():
            text = p.text.strip()
            preview = text[:20] + ("..." if len(text) > 20 else "")
            debug_print(f"The first paragraph containing text begins with: {preview}")
            break
    # An else on a for-loop runs if we never hit break. This is here because I'm maybe-overly defensive in programming style.
    else:
        debug_print("(Could not extract preview text)")

    # Return the runtime object
    return prs

# endregion



# region load & validate helpers
def _find_first_docx_paragraph_with_text(
    paragraphs: list[Paragraph_docx],
) -> Paragraph_docx | None:
    """Find the first paragraph that contains any text content in a docx."""
    for paragraph in paragraphs:
        if paragraph.text and paragraph.text.strip():
            return paragraph
    return None

def _find_first_slide_with_text(slides: list[Slide]) -> Slide | None:
    """Find the first slide that contains any paragraphs with text content."""
    for slide in slides:
        if get_slide_paragraphs(slide):
            return slide
    return None
# endregion