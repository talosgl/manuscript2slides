"""TODO"""
from pptx.slide import Slide
from docx2pptx.models import SlideNotes
from src.docx2pptx import utils
from src.docx2pptx import io
from src.docx2pptx import config
from src.docx2pptx.utils import debug_print
from src.docx2pptx.formatting import copy_run_formatting_pptx2docx, _apply_experimental_formatting_from_metadata
from src.docx2pptx.annotations.restore_from_slides import split_speaker_notes, safely_extract_comment_data, safely_extract_experimental_formatting_data
from docx import document
from docx.opc import constants
from docx.oxml.ns import qn
from docx.oxml.parser import OxmlElement as OxmlElement_docx
from docx.text.paragraph import Paragraph as Paragraph_docx
from docx.text.run import Run as Run_docx
from pptx.text.text import _Paragraph as Paragraph_pptx # type: ignore
from pptx.slide import Slide
from pptx import presentation

# region process_slide_paragraphs
def process_slide_paragraphs(
    slide: Slide, slide_notes: SlideNotes, new_doc: document.Document
) -> None:
    """
    Process a slide's body content, using any metadata stored in the speaker notes to restore formatting and annotation
    anchors from an earlier docx2pptx pipeline run. If we find metadata but aren't able to attach it to text content from
    the body paragraphs/runs, attach that as a new comment to the very last copied paragraph/run from that slide.

    Additionally, store non-metadata speaker notes content as a new comment, too.
    """

    slide_paragraphs: list[Paragraph_pptx] = io.get_slide_paragraphs(slide)

    unmatched_annotations = []
    matched_comment_ids = set()

    last_run = None

    # For every pptx paragraph.....
    for pptx_para in slide_paragraphs:

        # Make a new docx para
        new_para = new_doc.add_paragraph()

        # If the text of this paragraph exactly matches a previous heading's text, apply that heading style
        if slide_notes and slide_notes.has_metadata and slide_notes.headings:
            for heading in slide_notes.headings:
                if heading["text"].strip() == pptx_para.text.strip():
                    new_para.style = heading["name"]
                    break  # we should only ever apply one style to a paragraph


        for run in pptx_para.runs:

            # Handle adding hyperlinks versus regular runs, and the runs' basic formatting.
            if run.hyperlink.address:
                debug_print("Hyperlink address found.")
                run_from_hyperlink = add_hyperlink_to_docx_paragraph(
                    new_para, run.hyperlink.address
                )
                last_run = run_from_hyperlink
                copy_run_formatting_pptx2docx(run, run_from_hyperlink)
            else:
                new_docx_run = new_para.add_run()
                last_run = new_docx_run
                copy_run_formatting_pptx2docx(run, new_docx_run)

            # Check if this run contains matching text for comments from the speaker notes' stored JSON
            # metadata, from previous docx2pptx pipeline processing
            if slide_notes.comments:
                # Check to see if the run text matches any comments' ref text
                for comment in slide_notes.comments:
                    comment_data = safely_extract_comment_data(comment)

                    if comment_data is None:
                        debug_print(f"Skipping invalid comment: {comment}")
                        continue

                    if (
                        comment_data["reference_text"] in run.text
                        and comment_data["id"] not in matched_comment_ids
                    ):
                        new_doc.add_comment(
                            new_docx_run,
                            comment_data["text"],
                            comment_data["author"],
                            comment_data["initials"],
                        )
                        matched_comment_ids.add(comment_data["id"])
                    # don't break because there can be multiple comments added to a single run

            if config.EXPERIMENTAL_FORMATTING_ON and slide_notes.experimental_formatting:
                for exp_fmt in slide_notes.experimental_formatting:
                    fmt_info = safely_extract_experimental_formatting_data(exp_fmt)
                    if fmt_info is None:
                        continue
                    if exp_fmt["ref_text"] in run.text:
                        _apply_experimental_formatting_from_metadata(
                            new_docx_run, exp_fmt
                        )

    # Find all the unmatched annotations by getting the complement set to each of the matched_comments set
    unmatched_comments = [
        c for c in slide_notes.comments if c["id"] not in matched_comment_ids
    ]

    unmatched_annotations.extend(unmatched_comments)

    # Put the slide's user notes into a new comment attached to the last run
    if slide_notes and slide_notes.has_user_notes is True and last_run is not None:
        # append as a comment to the last run
        raw_comment_text = slide_notes.user_notes
        comment_header = "Copied from the PPTX Speaker Notes: \n\n"
        comment_text = comment_header + utils.sanitize_xml_text(raw_comment_text)

        if comment_text.strip():
            new_doc.add_comment(last_run, comment_text)

    # If we have any unmatched annotations from the slide_notes.metadata, attach that as a new comment to the last run
    if unmatched_annotations and last_run is not None:
        unmatched_parts = []
        # append as a comment to the last run
        for annotation in unmatched_annotations:
            if "original" in annotation:  # comments
                unmatched_parts.append(f"Comment: {annotation['original']['text']}")
            elif "text_body" in annotation:  # footnotes/endnotes
                unmatched_parts.append(f"Note: {annotation['text_body']}")
        combined = "\n\n".join(unmatched_parts)

        raw_comment_text = combined
        comment_header = "We found metadata for these annotations (comments, footnotes, or endnotes), but weren't able to match them to specific text in this paragraph: \n\n"
        comment_text = comment_header + utils.sanitize_xml_text(raw_comment_text)

        if comment_text.strip():
            new_doc.add_comment(last_run, comment_text)
# endregion

# region pptx2docx para/run processing helpers
def add_hyperlink_to_docx_paragraph(paragraph: Paragraph_docx, url: str) -> Run_docx:
    """
    Custom function to add Hyperlink objects to docx paragraphs using XML manipulation.

    - Create a regular run using paragraph.add_run()
    - Create the hyperlink XML element structure
    - Move the run's XML element from being a direct child of the paragraph to being a nested child of the Hyperlink element
    - Add the Hyperlink element (which now contains the run) to the paragraph

    Adapted from https://stackoverflow.com/questions/47666642/adding-an-hyperlink-in-msword-by-using-python-docx
    and https://github.com/python-openxml/python-docx/issues/384#issuecomment-294853130
    """
    # Create a new run on this paragraph
    run = paragraph.add_run()

    # Create the hyperlink structure
    part = paragraph.part
    r_id = part.relate_to(url, constants.RELATIONSHIP_TYPE.HYPERLINK, is_external=True)
    hyperlink = OxmlElement_docx("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)

    # Move the run from within the paragraph to within the Hyperlink
    run_element = run._element
    run_element.getparent().remove(run_element)  # Remove from paragraph
    hyperlink.append(run_element)  # Add to hyperlink
    paragraph._p.append(hyperlink)  # Add hyperlink to paragraph

    return run
# endregion



# region copy slides to docx body orchestrator
def copy_slides_to_docx_body(
    prs: presentation.Presentation, new_doc: document.Document
) -> None:
    """
    Sequentially process each slide in the deck by copying the paragraphs from the slide body into the docx's body. Analyze
    the speaker notes of the slide, seeking stored JSON metadata there from previous docx2pptx runs, and use that to apply
    formatting and/or annotations.
    """

    # Make a list of all slides
    slide_list = list(prs.slides)

    # For each slide...
    for slide in slide_list:

        # If there's slide notes, process them into a SlideNotes object; otherwise, make an empty one.
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            slide_notes = split_speaker_notes(slide.notes_slide.notes_text_frame.text)
        else:
            slide_notes = SlideNotes()

        process_slide_paragraphs(slide, slide_notes, new_doc)
